#!/usr/bin/env python3
"""Generate the 6-slide section 2.3 architecture-route deck.

Output: reports/6g-data-fabric-section-23-deck.pptx
Does not overwrite the older 16-page 2.3/3/4 deck.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
ASSET = ROOT / "reports" / "data-fabric-assets"
OUT = ROOT / "reports" / "6g-data-fabric-section-23-deck.pptx"

W_IN, H_IN = 13.333, 7.5
W, H = Inches(W_IN), Inches(H_IN)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
CORAL = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xB8, 0x86, 0x0B)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF4, 0xF7, 0xF8)
PALE = RGBColor(0xE8, 0xF0, 0xF2)
LINE = RGBColor(0xC5, 0xD0, 0xD4)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
SOFT_TEAL = RGBColor(0xD7, 0xEC, 0xEA)
SOFT_AMBER = RGBColor(0xF6, 0xED, 0xD0)
SOFT_NAVY = RGBColor(0xE6, 0xEC, 0xF2)
SOFT_CORAL = RGBColor(0xF8, 0xE6, 0xE4)
FONT = "WenQuanYi Micro Hei"
TOTAL = 6
EMU = 914400


def _set_run_east_asia(run) -> None:
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT)


def rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def lined(shape, fill: RGBColor, border: RGBColor, width: float = 0.75) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(width)


def write_box(tf, lines, *, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT, spacing=1.05, anchor=None):
    tf.word_wrap = True
    tf.auto_size = None
    if anchor is not None:
        tf.paragraphs[0].alignment = align
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(1)
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
        _set_run_east_asia(run)


def add_text(slide, l, t, w, h, lines, **kw):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    write_box(box.text_frame, lines, **kw)
    return box


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rgb(bg, WHITE)
    return s


def header(slide, number: str, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.07))
    rgb(bar, TEAL)
    add_text(slide, 0.38, 0.12, 12.55, 0.22, [number], size=10, color=TEAL, bold=True)
    add_text(slide, 0.38, 0.32, 12.55, 0.36, [title], size=20, color=NAVY, bold=True)


def footer(slide, page: int, chapter: str, refs: str = "") -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(7.22), Inches(12.55), Inches(0.01))
    rgb(line, LINE)
    left = f"6G 数据架构 × 数据编织 · 2.3 架构路线 · {chapter}"
    if refs:
        left = f"{left}    {refs}"
    add_text(slide, 0.38, 7.26, 10.4, 0.20, [left], size=8, color=MUTED)
    add_text(slide, 11.35, 7.26, 1.58, 0.20, [f"{page} / {TOTAL}"], size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def thesis(slide, text: str, top: float = 0.70, height: float = 0.54) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(top), Inches(12.55), Inches(height)
    )
    lined(box, PALE, TEAL)
    tf = box.text_frame
    tf.word_wrap = True
    write_box(tf, [text], size=12, color=NAVY, bold=True, spacing=1.04)
    tf.paragraphs[0].space_before = Pt(2)


def card(slide, l, t, w, h, title, body, accent=TEAL, title_size=11, body_size=9.5, fill=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    lined(shp, fill, LINE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(0.07), Inches(h))
    rgb(bar, accent)
    add_text(slide, l + 0.16, t + 0.05, w - 0.26, 0.24, [title], size=title_size, color=NAVY, bold=True)
    add_text(slide, l + 0.16, t + 0.30, w - 0.26, h - 0.36, body, size=body_size, color=SLATE, spacing=1.04)


def chip(slide, l, t, w, h, text, fill, color=WHITE, size=8):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    rgb(shp, fill)
    add_text(slide, l, t + 0.02, w, h - 0.02, [text], size=size, color=color, bold=True, align=PP_ALIGN.CENTER)


def number_badge(slide, l, t, n, fill=TEAL):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(0.26), Inches(0.26))
    rgb(circ, fill)
    add_text(slide, l, t + 0.01, 0.26, 0.24, [str(n)], size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def evidence_bar(slide, items, y: float = 6.28, h: float = 0.86) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(y), Inches(12.55), Inches(h))
    lined(box, OFF, LINE)
    n = len(items)
    gap = 0.10
    inner = 12.55 - 0.24
    cw = (inner - gap * (n - 1)) / n
    for i, (label, text, accent) in enumerate(items):
        x = 0.50 + i * (cw + gap)
        add_text(slide, x, y + 0.04, cw, 0.18, [label], size=8, color=accent, bold=True)
        add_text(slide, x, y + 0.22, cw, h - 0.28, [text], size=8, color=SLATE, spacing=1.02)


def add_table(slide, l, t, w, h, rows, col_w=None, font=8.5):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(l), Inches(t), Inches(w), Inches(h))
    table = table_shape.table
    if col_w:
        for i, cw in enumerate(col_w):
            table.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = NAVY if r == 0 else (OFF if r % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font)
            run.font.bold = r == 0 or c == 0
            run.font.color.rgb = WHITE if r == 0 else INK
            run.font.name = FONT
            _set_run_east_asia(run)
            cell.text_frame.word_wrap = True
            cell.text_frame.margin_left = Inches(0.05)
            cell.text_frame.margin_right = Inches(0.04)
            cell.text_frame.margin_top = Inches(0.03)
            cell.text_frame.margin_bottom = Inches(0.03)
    return table


def ensure_png(svg_name: str, png_name: str) -> Path:
    png = ASSET / png_name
    svg = ASSET / svg_name
    if png.exists():
        return png
    if svg.exists():
        import cairosvg

        cairosvg.svg2png(url=str(svg), write_to=str(png), output_width=2200)
        return png
    raise FileNotFoundError(png_name)


def fit_picture(slide, path: Path, box_l, box_t, box_w, box_h):
    with Image.open(path) as im:
        ar = im.size[0] / im.size[1]
    box_ar = box_w / box_h
    if ar >= box_ar:
        pw, ph = box_w, box_w / ar
    else:
        ph, pw = box_h, box_h * ar
    pl = box_l + (box_w - pw) / 2
    pt = box_t + (box_h - ph) / 2
    slide.shapes.add_picture(str(path), Inches(pl), Inches(pt), Inches(pw), Inches(ph))
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(pl), Inches(pt), Inches(pw), Inches(ph))
    frame.fill.background()
    frame.line.color.rgb = LINE
    frame.line.width = Pt(0.75)
    return pl, pt, pw, ph


def layer_box(slide, l, t, w, h, kicker, title, body, fill, title_color=WHITE, body_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    rgb(shp, fill)
    add_text(slide, l + 0.16, t + 0.06, w - 0.32, 0.18, [kicker], size=9, color=GOLD, bold=True)
    add_text(slide, l + 0.16, t + 0.24, w - 0.32, 0.26, [title], size=13, color=title_color, bold=True)
    add_text(slide, l + 0.16, t + 0.52, w - 0.32, h - 0.60, body if isinstance(body, list) else [body], size=10, color=body_color, spacing=1.04)


def slide_p1(prs) -> None:
    s = new_slide(prs)
    header(s, "2.3.1  ·  ARCHITECTURE CHOICE", "6G数据能力不会收敛为单一Fabric；三条路线围绕同一组控制能力竞争")
    thesis(
        s,
        "三条路线来自同一组架构矛盾，不是对Fabric理解不同，更不是成熟度阶梯。Ericsson优先低风险演进，Huawei优先近源执行，Nokia优先产品化速度；真正要比较的是控制权放在哪里、以什么代价换什么优势。",
    )

    constraints = [
        ("约束 1  ·  时标", "事实在近源", "无线状态、移动性和断链生存，要求毫秒执行留在网元、边缘或域控制器。把快环事实外移到远端平台，会同时失去时延和故障域。"),
        ("约束 2  ·  复用", "语义要跨域", "模型与产品要跨RAN、Core、OSS/BSS复用，必须有稳定对象、上下文、QoD和版本。名称相似不等于跨域语义已经对齐。"),
        ("约束 3  ·  成本", "少搬运、多计算", "感知、训练和遥测规模不允许默认把全量数据搬到中心。过滤、特征化与推理应先在源侧完成，再按用途交付。"),
        ("约束 4  ·  风险", "动作要回证", "Agent输出会改变真实网络。数据、模型、策略、动作、结果和回滚必须连成可解释、可追责的证据链，否则自治权限扩不出去。"),
    ]
    for i, (kicker, title, body) in enumerate(constraints):
        y = 1.32 + i * 1.08
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(y), Inches(4.15), Inches(1.00))
        lined(box, WHITE, LINE)
        number_badge(s, 0.50, y + 0.08, i + 1)
        add_text(s, 0.84, y + 0.05, 3.55, 0.16, [kicker], size=8.5, color=TEAL, bold=True)
        add_text(s, 0.84, y + 0.21, 3.55, 0.20, [title], size=12, color=NAVY, bold=True)
        add_text(s, 0.50, y + 0.42, 3.87, 0.54, [body], size=9.5, color=SLATE, spacing=1.02)

    collector = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.58), Inches(1.72), Inches(0.045), Inches(3.70))
    rgb(collector, TEAL)
    arr_in = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.58), Inches(3.58), Inches(0.24), Inches(0.22))
    rgb(arr_in, TEAL)
    hub = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.86), Inches(3.00), Inches(1.46), Inches(1.48))
    rgb(hub, NAVY)
    add_text(s, 4.92, 3.12, 1.34, 1.24, ["控制权", "如何分配", "决定路线"], size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, spacing=1.12)
    arr_out = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.36), Inches(3.58), Inches(0.20), Inches(0.22))
    rgb(arr_out, TEAL)

    routes = [
        (TEAL, "Ericsson  ·  存量联邦", "连接既有数据岛", "以复杂协调换迁移连续性。权威源与Mediation、NWDAF、湖仓不可跳过；统一的是发现、策略和消费体验，不是把全部运行时收进一个SKU。"),
        (AMBER, "Huawei  ·  网络原生拓扑", "重构数据拓扑", "以新平面复杂度换近源控制。把多源、多处理、多汇做成可编排执行，而不是先建设企业Catalog；缺少通用目录不能直接判为路线不完整。"),
        (NAVY, "Nokia  ·  分层商品化", "组合数据与自治软件", "以分层SKU换产品化速度。先让目录、AI和应用可分别销售并进入现网，再验证各层是否共享同一契约和控制体系。"),
    ]
    for i, (accent, kicker, title, body) in enumerate(routes):
        y = 1.32 + i * 1.44
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.58), Inches(y), Inches(6.35), Inches(1.36))
        lined(box, WHITE, LINE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.58), Inches(y), Inches(0.08), Inches(1.36))
        rgb(bar, accent)
        add_text(s, 6.82, y + 0.06, 5.95, 0.18, [kicker], size=9, color=accent, bold=True)
        add_text(s, 6.82, y + 0.24, 5.95, 0.22, [title], size=13, color=NAVY, bold=True)
        add_text(s, 6.82, y + 0.48, 5.95, 0.82, [body], size=10, color=SLATE, spacing=1.04)

    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(5.70), Inches(12.55), Inches(1.42))
    lined(band, SOFT_TEAL, TEAL)
    add_text(s, 0.54, 5.76, 12.23, 0.20, ["共同绕不开的不是平台名称，而是四类控制能力及其外部接口"], size=11, color=NAVY, bold=True)
    pills = [
        ("跨域语义", "稳定ID、权威源、配置版本和对象映射"),
        ("QoD / 产品契约", "用途适用性、Owner、SLO、计量与退出"),
        ("近源运行时", "算子放置、策略下发、断链生存与回退"),
        ("动作证据", "数据到回滚可关联、可导出、可测试"),
    ]
    for i, (label, body) in enumerate(pills):
        x = 0.54 + i * 3.10
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(6.02), Inches(2.95), Inches(0.96))
        rgb(pill, TEAL)
        add_text(s, x + 0.10, 6.08, 2.75, 0.28, [label], size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.10, 6.38, 2.75, 0.52, [body], size=9, color=WHITE, align=PP_ALIGN.CENTER, spacing=1.02)
    footer(s, 1, "2.3.1")


def slide_p2(prs) -> None:
    s = new_slide(prs)
    header(s, "2.3.2  ·  STANDARD BOUNDARY", "标准先冻结外部行为，不会先冻结统一平台")
    thesis(
        s,
        "ITU-R、SA2/SA5、O-RAN和外部运营接口正在形成一条可组合链路，但并不共同要求一个名为Data Fabric的网元。标准冻结跨厂商必须一致的对象与接口，不替厂商选择内部拓扑，因此三条路线可以长期并存。",
    )

    layers = [
        (
            1.32,
            RGBColor(0x1F, 0x4E, 0x5F),
            "LAYER 3  ·  外部产品行为",
            "API  ·  SLA  ·  用途  ·  计量  ·  结算",
            ["消费者看见的是可订购、可计量、可追责的服务，而不是内部网元名称。TMF/GSMA/CAMARA会先冻结这些外部行为，倒逼上游给出稳定语义和责任映射。"],
        ),
        (
            3.00,
            TEAL,
            "LAYER 2  ·  跨域控制语义",
            "稳定ID  ·  最小元数据  ·  QoD  ·  授权  ·  动作证据",
            ["跨域复用真正依赖对象能否对齐、质量能否声明、授权能否传递、动作能否回证。这是标准最可能先下手的控制层，也是三条路线最终都要交出的接口。"],
        ),
        (
            4.68,
            NAVY,
            "LAYER 1  ·  网络权威事实",
            "RAN  ·  Core  ·  OAM  ·  O-RAN接口",
            ["无线、会话和运维事实仍由网络域产生。标准先承认这些权威源，而不是另造一个总库来替代RAN快环、会话状态或OAM对象。"],
        ),
    ]
    for y, fill, kicker, title, body in layers:
        layer_box(s, 0.38, y, 7.20, 1.50, kicker, title, body, fill)
    for y in (4.40, 2.72):
        arr = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(3.68), Inches(y), Inches(0.42), Inches(0.26))
        rgb(arr, GOLD)

    card(
        s,
        7.78,
        1.32,
        5.15,
        2.36,
        "较可能先冻结",
        [
            "角色、稳定标识、最小元数据、交付控制、QoD/授权/证据挂钩点，以及一致性测试。",
            "ZSM 029对注册、发现、资产和工作流的研究也指向同一组外部行为，而不是指向某个平台形态。",
        ],
        TEAL,
        title_size=13,
        body_size=10.5,
        fill=SOFT_TEAL,
    )
    card(
        s,
        7.78,
        3.80,
        5.15,
        2.38,
        "仍然保持开放",
        [
            "既有NF增强、新NF、边缘Agent、管理侧服务、湖仓和消息运行时，都可能承载同一外部行为。",
            "是否存在独立Fabric仍是实现选择。标准不替厂商选定内部拓扑，也不把某家组件写成行业唯一答案。",
        ],
        AMBER,
        title_size=13,
        body_size=10.5,
        fill=SOFT_AMBER,
    )
    evidence_bar(
        s,
        [
            ("制度含义", "竞争从“组件是否存在”转向“接口能否互操作、后端能否替换、责任能否跨域传递”。名称相似不等于控制体系已经统一。", TEAL),
            ("判断含义", "同一外部行为可以由网络原生功能、管理域服务或外部数据层分别实现；不同数据族也可以选择不同后端。", NAVY),
            ("停止线", "Kafka、湖仓、完整知识图谱或独立Fabric都不是标准预设终点。不能用接口清单反推统一平台已经形成。", CORAL),
        ],
        y=6.28,
        h=0.86,
    )
    footer(s, 2, "2.3.2", "[S1][S5][S6][S26][S29][S7][S43][S15][S16]")


def slide_p3(prs, fig18: Path, north: Path) -> None:
    s = new_slide(prs)
    header(s, "2.3.4  ·  ERICSSON", "Ericsson以联邦保留存量：迁移连续性最强，跨产品一致性是关键风险")
    thesis(
        s,
        "这是一条“先联邦、后治理”的路线：数据在域内加工、按需联邦消费，统一的是治理视图而不是全部运行时。公开材料支持架构思想连续，不支持已经形成一个统一商用Data Fabric SKU。官方产品血缘是 Mediation→Telco DataOps，不是 EDCA→DataOps。",
    )

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.30), Inches(7.72), Inches(4.88))
    lined(panel, OFF, LINE)
    fit_picture(s, fig18, 0.50, 1.36, 7.48, 2.72)
    chip(s, 0.56, 1.42, 1.55, 0.22, "官方架构  Figure 18", NAVY, WHITE, 7.5)
    add_text(s, 0.50, 4.10, 7.48, 0.18, ["主图：2026 Data Ingestion Architecture。EDCA/DRG仍在，客户侧DMF与Ericsson侧EFDL/Data Fabric构成联邦边界。"], size=8, color=MUTED)

    fit_picture(s, north, 0.50, 4.30, 3.40, 1.78)
    chip(s, 0.56, 4.36, 1.35, 0.20, "参考架构  North Star", TEAL, WHITE, 7.5)
    add_text(
        s,
        4.00,
        4.30,
        3.96,
        1.78,
        [
            "① 近源摄取仍保留：EDCA、DRG在2026 Figure 18继续存在，证明联邦摄取机制没有被废弃。",
            "② 联邦边界被扩展：客户侧DMF连接Ericsson侧EFDL/Data Fabric，权威数据不必先搬迁。",
            "③ 治理范围上移：North Star增加目录、质量、血缘、语义/KG和数据产品，但是future-state参考架构，不是已交付统一SKU。",
        ],
        size=9,
        color=SLATE,
        spacing=1.06,
    )

    card(
        s,
        8.24,
        1.30,
        4.71,
        1.54,
        "数据与控制如何流动",
        [
            "数据在域内采集加工，经联邦访问和智能移动供跨域消费，重点不是把数据集中到一个湖。",
            "目录、策略、质量与血缘提供跨岛视图；执行仍分散在DataOps、NWDAF、EIAP/rApps和域系统。",
        ],
        TEAL,
    )
    card(
        s,
        8.24,
        2.94,
        4.71,
        1.54,
        "为什么选择联邦",
        [
            "RAN、Core、OSS/BSS和边缘拥有不同对象、时标和故障域，权威源天然分散。",
            "Mediation、NWDAF、SMO/RIC和客户湖仓不可跳过。先连接数据岛、对齐外部行为，再决定内部控制是否值得收敛。",
        ],
        NAVY,
    )
    card(
        s,
        8.24,
        4.58,
        4.71,
        1.60,
        "竞争含义",
        [
            "优势是渐进接入、多后端和源侧权威，迁移风险低，也更容易进入异构数据栈。",
            "代价是跨产品Schema、QoD、策略和版本可能长期碎片化；协调成本若高于替换成本，联邦优势会下降。",
        ],
        AMBER,
        fill=SOFT_AMBER,
    )

    evidence_bar(
        s,
        [
            ("已证明", "Figure 18延续EDCA/DRG；Mediation、EIAP、NWDAF和API渠道分别有产品或部署锚点。North Star把范围扩到语义/KG、质量和数据产品，但仍是参考架构。", TEAL),
            ("关键边界", "官方产品血缘是 Mediation→Telco DataOps，不是 EDCA→DataOps。DDC/GDC名称消失只能说明功能重组。“no data copy”应读成减少不必要复制。", CORAL),
            ("决定性验证", "同一版本化数据产品能否跨DataOps、EIAP、NWDAF及第三方后端复用，无需项目级重写，并在断链、回退与动作回证上满足通信预算。", AMBER),
        ],
    )
    footer(s, 3, "2.3.4", "[S9][S10][S21][S30][S35][S36][S57][S58]")


def slide_p4(prs, huawei: Path) -> None:
    s = new_slide(prs)
    header(s, "2.3.5  ·  HUAWEI", "Huawei重构数据拓扑：近源执行控制最强，电信级可靠性是验证门")
    thesis(
        s,
        "这是一条“先重构执行、再决定治理边界”的路线：DO编排拓扑，DA近源处理，DCP异步分发。它首先是6G数据面机制，能吸收部分发现、编排和交付能力，但不会自动解决跨域语义、QoD、血缘与产品责任。不能把“补成Fabric”当成唯一终点。",
    )

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.30), Inches(7.72), Inches(4.88))
    lined(panel, OFF, LINE)
    fit_picture(s, huawei, 0.50, 1.36, 7.48, 3.28)
    chip(s, 0.56, 1.42, 1.55, 0.22, "研究架构  非3GPP方案", AMBER, NAVY, 7.5)
    notes = [
        (TEAL, "① DO", "能力注册、操作DAG分解、选点与策略下发。登记的是“谁能执行什么”，不是通用资产目录。"),
        (NAVY, "② DA / DPF / DSF", "在UE、RAN、边缘和Core近源采集、过滤、推理与存储，把计算放到数据所在故障域。"),
        (CORAL, "③ DCP / Data Spine", "以Pub/Sub解耦生产者和多消费者，服务AI/ISAC扇出，而不是沿用GTP会话模型。"),
    ]
    for i, (accent, title, body) in enumerate(notes):
        x = 0.50 + i * 2.50
        add_text(s, x, 4.72, 2.40, 0.18, [title], size=9.5, color=accent, bold=True)
        add_text(s, x, 4.90, 2.40, 1.16, [body], size=9, color=SLATE, spacing=1.03)

    card(
        s,
        8.24,
        1.30,
        4.71,
        1.54,
        "控制流",
        [
            "服务请求向DO声明采集、过滤、聚合、转换、分析或推理目标。",
            "DO从能力注册中选择执行节点，拆解操作DAG，并下发拓扑、访问、隐私与资源策略。",
        ],
        TEAL,
    )
    card(
        s,
        8.24,
        2.94,
        4.71,
        1.54,
        "数据流",
        [
            "源DA发布topic，DCP按订阅或带内拓扑把数据异步送往多个处理/消费节点。",
            "DA/DPF沿途过滤、压缩、特征化和本地推理，上游模型输出还可成为下游输入。",
        ],
        NAVY,
    )
    card(
        s,
        8.24,
        4.58,
        4.71,
        1.60,
        "为什么需要新平面",
        [
            "SBI偏控制消息，PDU会话围绕连接；AI/ISAC需要多源、多处理、多汇和派生数据链。",
            "跨域语义与产品契约可能被网络原生吸收，也可能由管理域补充。AUTINOps不是DCP商用证明。",
        ],
        AMBER,
        fill=SOFT_AMBER,
    )
    evidence_bar(
        s,
        [
            ("原型信号", "DCP公开实验约2ms；10 broker、128KB消息约2,394.6MB/s，说明精简路径与横向吞吐具有潜力，仍属研究原型。", TEAL),
            ("比较局限", "DCP移除持久化，而RocketMQ对照开启持久化；结果不能外推为同功能全面领先，也没有证明长期可靠性。", CORAL),
            ("决定性验证", "同等持久化条件下的跨站故障、移动性、断链回退、多租户SLA和第三方算子测试。增强SBA/UP若已够用，新增平面必要性下降。", AMBER),
        ],
    )
    footer(s, 4, "2.3.5", "[S12][S14][S37][S38]")


def slide_p5(prs, nokia: Path) -> None:
    s = new_slide(prs)
    header(s, "2.3.6  ·  NOKIA", "Nokia分层商品化：产品进入现网最快，统一控制与开放性仍待证明")
    thesis(
        s,
        "这是一条“先产品化能力、再验证统一控制”的路线。Data Suite是数据管理层，AN Fabric是智能消费层；动作仍回OSS/RIC/SON，数据默认留域。两套Fabric不是同一层。真正的竞争问题是它们能否共享契约并贯通到动作证据。",
    )

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.30), Inches(7.72), Inches(4.88))
    lined(panel, OFF, LINE)
    add_text(s, 0.52, 1.34, 4.55, 0.18, ["主视觉  ·  自绘五段链"], size=8.5, color=TEAL, bold=True)

    stages = [
        (NAVY, "1  域数据底座", "Domain lakehouse  ·  各域本地采集、加工与存储"),
        (TEAL, "2  Common Fabric / Lakehouse", "适配、抽象和联邦访问，提供逻辑视图而非强制集中"),
        (TEAL, "3  Data Suite", "目录、语义、QoD、血缘和数据产品，与治理最直接重合"),
        (AMBER, "4  AN Fabric / AN Apps", "Sense–Think–Act 智能消费与编排，不是数据管理Fabric"),
        (NAVY, "5  域控制器 + API / Marketplace", "高风险动作回OSS/RIC/SON，再向计量和渠道外延"),
    ]
    for i, (fill, title, body) in enumerate(stages):
        y = 1.54 + i * 0.86
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.52), Inches(y), Inches(4.55), Inches(0.72))
        rgb(box, fill)
        add_text(s, 0.66, y + 0.05, 4.27, 0.26, [title], size=11, color=WHITE, bold=True)
        add_text(s, 0.66, y + 0.32, 4.27, 0.34, [body], size=9, color=WHITE)
        if i < len(stages) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.62), Inches(y + 0.70), Inches(0.22), Inches(0.14))
            rgb(arr, GOLD)

    fit_picture(s, nokia, 5.20, 1.54, 2.72, 2.42)
    chip(s, 5.26, 1.60, 2.20, 0.20, "官方图  仅证智能消费层", AMBER, NAVY, 7.5)
    add_text(
        s,
        5.20,
        4.04,
        2.72,
        2.04,
        [
            "官方AN Fabric图只证明Sense–Think–Act智能层和Marketplace外延。",
            "Data Suite、域数据底座和公共湖仓不在这张图里，必须靠左侧自绘链路补齐，否则会把两套Fabric读成同一层。",
            "高风险动作仍经AN Apps回到域控制器；公共层偏Non-RT，不能替代RAN快环。",
        ],
        size=9,
        color=SLATE,
        spacing=1.06,
    )

    card(
        s,
        8.24,
        1.30,
        4.71,
        1.54,
        "必须分层",
        [
            "Data Suite管目录、语义、质量、血缘和产品；AN Fabric管智能消费与编排。",
            "数据默认留域，公共层只做适配、抽象和联邦访问。支持多厂商只有经过第三方互操作和双向迁移，才构成中立性证据。",
        ],
        TEAL,
    )
    card(
        s,
        8.24,
        2.94,
        4.71,
        1.54,
        "数据与控制如何流动",
        [
            "域湖仓完成本地加工，经公共层进入Data Suite，再进入Sense，建立观察与chain-of-custody。",
            "Think使用电信模型和Glass-Box推理；Act通过AN Apps调用既有OSS/RIC/SON，再向计量和Marketplace延伸。",
        ],
        NAVY,
    )
    card(
        s,
        8.24,
        4.58,
        4.71,
        1.60,
        "竞争含义",
        [
            "优势是SKU可组合、进入现网快，不必等待6G新网元标准。",
            "代价是各层可能只在品牌层组合，并形成语义、应用与云后端的软锁定。客户若只买域App，分层复利就落空。",
        ],
        AMBER,
        fill=SOFT_AMBER,
    )
    evidence_bar(
        s,
        [
            ("域应用已证", "KDDI覆盖20万4G/5G RAN；stc披露超过1万次自动动作和性能改善；Indosat证明节能应用可扩展。它们证明域级闭环，不证明统一跨域Fabric。", TEAL),
            ("规模信号", "匿名Tier-1：25万小区、27万报告/秒，仅支持摄取与关联潜力。因客户未具名，不能外推广泛跨域商用。", AMBER),
            ("决定性验证", "Data Suite与AN Fabric是否共享同一契约，第三方后端/Agent能否双向替换，动作与回滚证据能否贯通域控制器。", CORAL),
        ],
    )
    footer(s, 5, "2.3.6", "[S23][S24][S31][S32][S33][S34][S39]")


def slide_p6(prs) -> None:
    s = new_slide(prs)
    header(s, "2.3.3 + 2.3.7  ·  SYNTHESIS", "三条路线不构成成熟度序列；最终胜负由四类接口控制权决定")
    thesis(
        s,
        "Ericsson争联邦治理，Huawei争近源运行时，Nokia争语义产品与应用入口。谁能定义这些能力的外部接口，同时允许执行层和后端替换，谁才更接近长期平台位置；谁只提供封闭组件或单点应用，谁就仍停留在项目层。",
    )

    add_table(
        s,
        0.38,
        1.28,
        12.55,
        2.68,
        [
            ["比较项", "Ericsson  ·  存量联邦", "Huawei  ·  网络原生拓扑", "Nokia  ·  分层商品化"],
            ["架构起点", "分散数据岛和既有产品如何少迁移、可复用", "AI/ISAC多源多处理多汇如何低时延流动", "数据管理、自治应用和商业外延如何形成可售组合"],
            ["控制中心", "目录、治理和策略逻辑联邦；产品控制仍分散", "DO编排执行能力与数据拓扑", "Data Suite管数据产品，AN Fabric管智能消费"],
            ["执行位置", "域内产品、EDCA/DMF与数据湖按需访问", "DA/DPF/DSF在UE/RAN/边缘/Core近源执行", "数据留域；动作经AN Apps回OSS/RIC/SON"],
            ["换取优势", "保护存量、降低迁移风险、允许多后端", "降低全量搬运，获得可编排近源处理", "能力可分别产品化并快速进入现网"],
            ["结构性代价", "跨产品语义、QoD、策略和版本协调持续存在", "新增顺序、持久化、容灾、移动性与可移植性", "多层SKU可能只在品牌层组合并形成软锁定"],
            ["决定性验证", "同一数据产品跨DataOps、EIAP、NWDAF和第三方复用", "同功能对比下通过持久化、故障、移动性和租户SLA", "三层共享契约，并支持第三方双向替换"],
        ],
        col_w=[1.55, 3.67, 3.67, 3.66],
        font=8.5,
    )

    controls = [
        ("1  跨域语义与身份", "定义稳定ID、权威源、配置版本和对象映射，使同一小区、会话、模型与业务对象能跨接口关联。谁掌握这组映射，谁就能让模型和产品复用。"),
        ("2  QoD与产品契约", "定义数据对特定任务是否适用、谁负责、满足何种SLO、如何计量、何时降级或退出。没有契约，跨域交付只能停留在项目集成。"),
        ("3  近源运行时", "决定过滤、特征、推理与缓存放在哪里，策略如何下发，断链如何生存，版本如何回退。快环不能外包给远端目录或云平台。"),
        ("4  动作证据与接口", "把数据、模型、策略、动作、结果和回滚关联起来，并能被第三方审计、导出与一致性测试。没有证据链，自治权限扩不出去。"),
    ]
    for i, (title, body) in enumerate(controls):
        x = 0.38 + i * 3.20
        card(s, x, 4.06, 3.08, 1.36, title, [body], TEAL, title_size=11, body_size=9)

    land = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(5.50), Inches(12.55), Inches(1.62))
    lined(land, SOFT_NAVY, NAVY)
    add_text(s, 0.54, 5.54, 12.23, 0.18, ["路线落点与竞合边界"], size=10, color=TEAL, bold=True)
    add_text(
        s,
        0.54,
        5.74,
        12.23,
        0.36,
        ["Ericsson争夺联邦目录与渐进治理控制权。Huawei争夺数据拓扑与近源运行时控制权。Nokia争夺语义产品与自治应用入口控制权。三者可替代也可互补，但都不是一条通向统一Fabric的成熟度阶梯。"],
        size=10.5,
        color=NAVY,
        bold=True,
        spacing=1.03,
    )
    add_text(
        s,
        0.54,
        6.12,
        12.23,
        0.40,
        ["真正的平台地位不是把所有数据收到一个平台，而是让跨域消费者遵循自己的语义、QoD和证据契约，同时证明运行时、后端和渠道仍可互操作、可替换。标准更可能冻结这些外部行为，而不是替市场选出唯一Data Fabric。"],
        size=10.5,
        color=SLATE,
        spacing=1.03,
    )
    add_text(
        s,
        0.54,
        6.54,
        12.23,
        0.48,
        ["竞合边界：设备商更接近网络语义、近源执行与动作权限；云和湖仓更接近规模数据/模型运行时；聚合商更接近跨网分发与结算。三者可以组合，但RAN快环和网元动作权限不能外移。Google Telecom Data Fabric仍为Private Preview，不能外推为已形成的跨网统一平台。"],
        size=9.5,
        color=MUTED,
        spacing=1.03,
    )
    footer(s, 6, "2.3.3 / 2.3.7", "[S18][S25] 及 P3–P5 厂商证据")


def count_cjk(prs: Presentation) -> list[int]:
    import re

    counts = []
    for slide in prs.slides:
        texts = []
        for shp in slide.shapes:
            if shp.has_text_frame:
                texts.append(shp.text_frame.text)
            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        counts.append(len(re.findall(r"[\u4e00-\u9fff]", "".join(texts))))
    return counts


def verify(prs: Presentation) -> None:
    n = len(prs.slides)
    if n != TOTAL:
        raise SystemExit(f"expected {TOTAL} slides, got {n}")
    print("cjk chars:", count_cjk(prs), "total", sum(count_cjk(prs)))

    pics = []
    tables = []
    overflow = []
    for i, slide in enumerate(prs.slides, 1):
        pic_n = 0
        table_n = 0
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_n += 1
            if shp.has_table:
                table_n += 1
            right = (shp.left + shp.width) / EMU
            bottom = (shp.top + shp.height) / EMU
            if right > W_IN + 0.03 or bottom > H_IN + 0.03 or shp.left < Emu(-0.02 * EMU) or shp.top < Emu(-0.02 * EMU):
                overflow.append((i, getattr(shp, "name", "?"), round(right, 3), round(bottom, 3)))
        pics.append(pic_n)
        tables.append(table_n)

    if pics[2] != 2:
        raise SystemExit(f"P3 must embed 2 Ericsson images, got {pics[2]}")
    if pics[3] != 1:
        raise SystemExit(f"P4 must embed 1 Huawei image, got {pics[3]}")
    if pics[4] != 1:
        raise SystemExit(f"P5 must embed 1 Nokia image, got {pics[4]}")
    if sum(tables) != 1 or tables[5] != 1:
        raise SystemExit(f"deck must contain exactly 1 table on P6, got {tables}")
    if overflow:
        raise SystemExit(f"shapes overflow canvas: {overflow}")
    print("verify ok:", {"slides": n, "pictures": pics, "tables": tables})


def build() -> Path:
    fig18 = ensure_png("ericsson-data-ingestion-architecture-2026.svg", "ericsson-data-ingestion-architecture-2026.png")
    north = ensure_png("ericsson-ai-ready-data-management-2026.svg", "ericsson-ai-ready-data-management-2026.png")
    huawei = ASSET / "huawei-6g-data-plane-architecture.jpg"
    nokia = ASSET / "nokia-autonomous-network-fabric.png"
    for path in (fig18, north, huawei, nokia):
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_p1(prs)
    slide_p2(prs)
    slide_p3(prs, fig18, north)
    slide_p4(prs, huawei)
    slide_p5(prs, nokia)
    slide_p6(prs)
    verify(prs)
    prs.save(OUT)
    print(f"wrote {OUT}")
    return OUT


if __name__ == "__main__":
    build()
