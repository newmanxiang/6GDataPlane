#!/usr/bin/env python3
"""Generate the 6-slide section 3 industry-trend deck.

Output: reports/6g-data-fabric-section-3-deck.pptx
Source: reports/plan-section3.md and chapter 3 of the full-sections HTML.
"""

from __future__ import annotations

import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "reports" / "6g-data-fabric-section-3-deck.pptx"

W_IN, H_IN = 13.333, 7.5
W, H = Inches(W_IN), Inches(H_IN)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
CORAL = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xB8, 0x86, 0x0B)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF4, 0xF7, 0xF8)
PALE = RGBColor(0xE8, 0xF0, 0xF2)
LINE = RGBColor(0xC5, 0xD0, 0xD4)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
SOFT_TEAL = RGBColor(0xD7, 0xEC, 0xEA)
SOFT_AMBER = RGBColor(0xF6, 0xED, 0xD0)
SOFT_NAVY = RGBColor(0xE6, 0xEC, 0xF2)
SOFT_CORAL = RGBColor(0xF8, 0xE6, 0xE4)
SOFT_GREEN = RGBColor(0xDE, 0xF0, 0xE4)
GREEN = RGBColor(0x1F, 0x7A, 0x4D)
DEEP_GREEN = RGBColor(0x14, 0x5A, 0x3A)
FONT = "WenQuanYi Micro Hei"
TOTAL = 6
EMU = 914400


def _set_run_east_asia(run) -> None:
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT)


def rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def lined(shape, fill: RGBColor, border: RGBColor, width: float = 0.75) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(width)


def write_box(tf, lines, *, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT, spacing=1.05):
    tf.word_wrap = True
    tf.auto_size = None
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(1)
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
        _set_run_east_asia(run)


def add_text(slide, l, t, w, h, lines, **kw):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    write_box(box.text_frame, lines, **kw)
    return box


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rgb(bg, WHITE)
    return s


def header(slide, number: str, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.07))
    rgb(bar, TEAL)
    add_text(slide, 0.38, 0.12, 12.55, 0.20, [number], size=10, color=TEAL, bold=True)
    add_text(slide, 0.38, 0.30, 12.55, 0.34, [title], size=17, color=NAVY, bold=True)


def footer(slide, page: int, chapter: str, refs: str = "") -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(7.22), Inches(12.55), Inches(0.01))
    rgb(line, LINE)
    left = f"6G 数据架构 × 数据编织 · 第三章产业趋势 · {chapter}"
    if refs:
        left = f"{left}    {refs}"
    add_text(slide, 0.38, 7.26, 10.4, 0.20, [left], size=8, color=MUTED)
    add_text(slide, 11.35, 7.26, 1.58, 0.20, [f"{page} / {TOTAL}"], size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def thesis(slide, text: str, top: float = 0.64, height: float = 0.42) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(top), Inches(12.55), Inches(height)
    )
    lined(box, PALE, TEAL)
    tf = box.text_frame
    tf.word_wrap = True
    write_box(tf, [text], size=11, color=NAVY, bold=True, spacing=1.03)
    tf.paragraphs[0].space_before = Pt(2)


def chip(slide, l, t, w, h, text, fill, color=WHITE, size=8):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    rgb(shp, fill)
    add_text(slide, l, t + 0.02, w, h - 0.02, [text], size=size, color=color, bold=True, align=PP_ALIGN.CENTER)


def evidence_bar(slide, items, y: float = 6.28, h: float = 0.86) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(y), Inches(12.55), Inches(h))
    lined(box, OFF, LINE)
    n = len(items)
    gap = 0.10
    inner = 12.55 - 0.24
    cw = (inner - gap * (n - 1)) / n
    for i, (label, text, accent) in enumerate(items):
        x = 0.50 + i * (cw + gap)
        add_text(slide, x, y + 0.04, cw, 0.16, [label], size=8, color=accent, bold=True)
        add_text(slide, x, y + 0.20, cw, h - 0.26, [text], size=8, color=SLATE, spacing=1.02)


def add_table(slide, l, t, w, h, rows, col_w=None, font=8.5):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(l), Inches(t), Inches(w), Inches(h))
    table = table_shape.table
    if col_w:
        for i, cw in enumerate(col_w):
            table.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = NAVY if r == 0 else (OFF if r % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font)
            run.font.bold = r == 0 or c == 0
            run.font.color.rgb = WHITE if r == 0 else INK
            run.font.name = FONT
            _set_run_east_asia(run)
            cell.text_frame.word_wrap = True
            cell.text_frame.margin_left = Inches(0.05)
            cell.text_frame.margin_right = Inches(0.04)
            cell.text_frame.margin_top = Inches(0.03)
            cell.text_frame.margin_bottom = Inches(0.03)
    return table


def slide_p1(prs) -> None:
    s = new_slide(prs)
    header(s, "3.1  ·  T1–T8 产业趋势总览", "从R1–R8到T1–T8：形成“可信底座、动作准入、双路径价值、标准横切”的产业结构")
    thesis(
        s,
        "四项筛选只决定哪些能力进入T1–T8；趋势之间的关系来自最小前置条件、独立价值和横切范围。可信供给同时支撑能力服务与自动化，但T4/T5只约束承担动作责任的T7-B；因此不是八条平行口号，也不是四级严格串行。",
    )

    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.16), Inches(3.72), Inches(4.96))
    lined(left, WHITE, LINE)
    add_text(s, 0.52, 1.22, 3.44, 0.22, ["趋势筛选方法"], size=12, color=NAVY, bold=True)
    gates = [
        (TEAL, "入选门槛 · 需求刚性", "问题须由6G物理分布、AI原生、跨域自治和生态开放产生；不解决会持续限制数据使用、动作安全或商业规模。只提升功能便利性、没有结构性约束的议题不进入趋势清单。"),
        (TEAL, "入选门槛 · 可证伪性", "必须能设置领先信号、可测试指标、时间窗和反证条件；否则只是愿景，不能作为趋势判断。每条趋势都要说明何时降级、停止或回到既有机制。"),
        (AMBER, "确定性 · 工程成熟度", "检查是否已有组件、原型、参考架构或部署基础，以及接口、运行时、治理和运维机制是否可实施。单点功能存在不等于跨域体系成熟。"),
        (AMBER, "确定性 · 产业牵引力", "标准讨论、厂商投入、客户需求和商业信号是否同向；证据按规范、产品、具名部署、研究和愿景分档，避免把路线图或品牌组合当成已证事实。"),
    ]
    for i, (accent, title, body) in enumerate(gates):
        y = 1.48 + i * 1.12
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.52), Inches(y), Inches(3.44), Inches(1.04))
        lined(box, SOFT_TEAL if i < 2 else SOFT_AMBER, accent)
        add_text(s, 0.64, y + 0.05, 3.20, 0.20, [title], size=10, color=NAVY, bold=True)
        add_text(s, 0.64, y + 0.26, 3.20, 0.72, [body], size=8.5, color=SLATE, spacing=1.02)

    nodes = [
        (4.28, 1.16, 4.18, 1.58, TEAL, "可信数据供给底座", "T2 元数据事件化  ·  T3 全局控制/近源执行\nT1 可运营数据产品  ·  T6 任务适用性准入", "及时可达 × 正确使用；共同提供稳定对象、版本、用途和适用性证明，工程上互为条件、并行建设"),
        (8.64, 1.16, 4.29, 1.58, GREEN, "T7 能力型服务路径", "数据/API、位置、身份、分析、网络能力\n买方保留决策权与动作权", "不必等待高风险自治成熟，可经聚合层跨网复制，按调用、会话、覆盖、用量或订阅先行收费"),
        (4.28, 2.90, 4.18, 1.58, CORAL, "自动化动作可信准入", "T4 供数—决策—执行—回证\nT5 NDT预演成为放行证据", "只约束高风险、跨域自治扩权；要求身份、策略、动作证据、风险预算与回滚，不否定只读分析和建议"),
        (8.64, 2.90, 4.29, 1.58, DEEP_GREEN, "T7 结果型服务路径", "体验、交易、风控、运营结果\n供应方承担部分效果责任", "额外依赖T4/T5、结果基线、因果归因、动作授权和回滚；按托管费、效果增益或结果分成"),
    ]
    for l, t, w, h, accent, title, mid, foot in nodes:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        lined(box, WHITE, LINE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(0.08), Inches(h))
        rgb(bar, accent)
        add_text(s, l + 0.18, t + 0.06, w - 0.28, 0.24, [title], size=12, color=NAVY, bold=True)
        add_text(s, l + 0.18, t + 0.32, w - 0.28, 0.72, mid.split("\n"), size=10, color=SLATE, spacing=1.05)
        add_text(s, l + 0.18, t + 1.08, w - 0.28, 0.42, [foot], size=8.5, color=MUTED, spacing=1.02)

    arr1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.46), Inches(1.78), Inches(0.18), Inches(0.18))
    rgb(arr1, TEAL)
    arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.22), Inches(2.74), Inches(0.18), Inches(0.16))
    rgb(arr2, CORAL)
    arr3 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.46), Inches(3.52), Inches(0.18), Inches(0.18))
    rgb(arr3, DEEP_GREEN)

    t8 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.28), Inches(4.60), Inches(8.65), Inches(0.68))
    lined(t8, SOFT_NAVY, NAVY)
    add_text(s, 4.44, 4.66, 8.33, 0.22, ["T8 横切全部趋势：标准更可能冻结能力语义、信息模型、接口行为和一致性测试"], size=11, color=NAVY, bold=True)
    add_text(s, 4.44, 4.90, 8.33, 0.30, ["不构成第五个阶段；增强既有NF、SMO服务、边缘Agent或新功能均可承载同一外部行为。验证重点是互操作、替换成本和额外状态。"], size=9, color=SLATE)

    add_table(
        s,
        4.28,
        5.38,
        8.65,
        1.74,
        [
            ["趋势", "确定性", "主要依据"],
            ["T1 / T2 / T3", "高", "需求刚性强；目录、事件、近源处理和产品化已有工程基础，主要缺跨产品复用与多厂商对账"],
            ["T4 / T5 / T6", "中高", "机制明确；通信级闭环、动作证据、适用性信封与强制准入仍需具名部署验证"],
            ["T7", "A中高 / B低—中", "A已有CAMARA/Open Gateway付费信号；B的结果基线、归因、责任转移和规模分成尚不成熟"],
            ["T8", "低—中", "能力与接口收敛方向明确；最终网元、职责边界、部署拓扑和互操作规则尚未冻结"],
        ],
        col_w=[1.70, 1.15, 5.80],
        font=8.5,
    )
    footer(s, 1, "3.1", "筛选口径：需求刚性 · 可证伪性 · 工程成熟 · 产业牵引")


def slide_p2(prs) -> None:
    s = new_slide(prs)
    header(s, "T2 / T3  ·  分布式供给基础", "元数据事件化与近源执行共同解决“变化可见、数据按时可达”")
    thesis(
        s,
        "静态目录必然落后于生产状态；高频无线与感知数据又不能默认全量集中。因此必须同时具备“快照—变化事件—对账”的状态机制，以及“全局定义—近源执行—异步回证”的部署机制；两者共同回答数据是否及时、可定位、可回退。",
    )

    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.16), Inches(6.20), Inches(5.00))
    lined(left, WHITE, LINE)
    chip(s, 0.52, 1.26, 0.70, 0.22, "T2  高", TEAL)
    add_text(s, 1.30, 1.24, 5.10, 0.26, ["元数据将从静态目录升级为“变化事件”"], size=13, color=NAVY, bold=True)
    add_text(
        s,
        0.52,
        1.54,
        5.90,
        0.46,
        ["完整机制：基线快照 + 变化事件 + 周期对账。事件携带事件时间、版本、来源、作用域、Correlation ID、变更前后值，以及受影响的产品、模型、缓存、策略和消费者；无法仅靠事件恢复全量状态。"],
        size=9.5,
        color=SLATE,
        spacing=1.03,
    )
    blocks = [
        ("事件范围", "Schema/对象关系、拓扑与配置、任务适用性与权限、血缘与成本、数据—模型—动作关联。对象覆盖小区、波束、频段、UE群、会话、切片、云资源及跨域映射；每类对象需定义Owner与版本源。"),
        ("主要作用", "自动生成影响清单；触发契约测试、缓存/特征/模型失效、策略复核和重新审批；支持跨域变更回放、延迟传播监测与责任追踪。消费者只订阅与自身作用域相关的变化。"),
        ("成立逻辑", "目录只描述扫描时刻；配置、拓扑、权限和质量持续变化。模型和Agent使用过期信息会把报表错误升级为错误动作；跨接口对象若无统一事件时间与版本，也无法判断先后和影响范围。"),
        ("适用边界", "不把全部业务数据改成事件；中央事件总线不进入毫秒快环；事件不能替代本地状态，必须保留快照和周期对账。若事件丢失、乱序或延迟超限，消费者应退回基线或停止自动动作。"),
    ]
    for i, (title, body) in enumerate(blocks):
        y = 2.06 + i * 0.98
        add_text(s, 0.52, y, 5.90, 0.18, [title], size=10, color=TEAL, bold=True)
        add_text(s, 0.52, y + 0.18, 5.90, 0.74, [body], size=9, color=SLATE, spacing=1.02)

    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.74), Inches(1.16), Inches(6.19), Inches(5.00))
    lined(right, WHITE, LINE)
    chip(s, 6.88, 1.26, 0.70, 0.22, "T3  高", TEAL)
    add_text(s, 7.66, 1.24, 5.10, 0.26, ["部署形态将收敛为“全局控制、近源执行”"], size=13, color=NAVY, bold=True)
    add_text(
        s,
        6.88,
        1.54,
        5.89,
        0.46,
        ["高频数据全量远传受带宽、时延、能耗、隐私和断链约束。快环、MEC、专网与NTN必须本地运行；完全分散又会造成语义、版本、策略和责任漂移，因此需要逻辑统一而非物理集中。"],
        size=9.5,
        color=SLATE,
        spacing=1.03,
    )
    layers = [
        (TEAL, "全局层", "定义对象/语义、任务适用性模板、产品、策略、模型、权限和版本基线，发布兼容关系与回退规则。统一的是规则、责任和验收口径，不是全部存储与处理位置。"),
        (NAVY, "本地层", "按资源预算执行过滤、聚合、特征、缓存、预编译策略、受限推理与白名单动作；维护实时状态。断链期间只运行已批准版本，超时或冲突则降级。"),
        (AMBER, "异步回流", "上传摘要、任务适用性、执行证据、漂移、冲突和资源消耗；支持对账、模型重验与版本回退。远端目录、审批或云模型不得成为快环同步依赖。"),
    ]
    for i, (accent, title, body) in enumerate(layers):
        y = 2.08 + i * 0.92
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.88), Inches(y), Inches(5.89), Inches(0.86))
        lined(box, SOFT_TEAL if i == 0 else (SOFT_NAVY if i == 1 else SOFT_AMBER), accent)
        add_text(s, 7.02, y + 0.06, 5.61, 0.20, [title], size=11, color=NAVY, bold=True)
        add_text(s, 7.02, y + 0.28, 5.61, 0.52, [body], size=9, color=SLATE, spacing=1.02)
    add_text(
        s,
        6.88,
        4.90,
        5.89,
        1.12,
        [
            "边界：不同网络仍可采用不同层级和拓扑。Ericsson强调动态放置，Nokia保留域自治，Huawei提出近源算子；它们共同指向逻辑全局控制与本地确定性执行。",
            "趋同的是时标、故障域、版本和责任约束，不是统一产品或网元名称。若中心不可达导致本地停摆，或本地长期不回证，两级架构均不成立。",
        ],
        size=9,
        color=MUTED,
        spacing=1.04,
    )

    evidence_bar(
        s,
        [
            ("领先信号", "T2：Schema/配置/任务适用性/权限变更自动生成影响清单并触发CI、缓存/模型失效或策略检查。T3：近源算子与全局目录/策略解耦，断链按批准版本运行并异步回证。", TEAL),
            ("验证指标", "T2：发现P95、事件丢失/乱序、影响覆盖、误报漏报、对账差异、失效传播。T3：P99、跨域字节减量、断链时长、资源预算、策略同步、证据完整和回退成功。", NAVY),
            ("反证条件", "T2：目录仍靠人工扫描，事故早于变化发现，事件无法关联对象、版本和消费者。T3：所有查询/审批必须经过中心，本地无降级能力，或长期不回证导致全局状态失真。", CORAL),
        ],
        y=6.26,
        h=0.88,
    )
    footer(s, 2, "T2 / T3", "[S27] OpenLineage  ·  [S9][S12][S20] 两级部署")


def slide_p3(prs) -> None:
    s = new_slide(prs)
    header(s, "T1 / T6  ·  数据运营基础", "网络数据将由接口资源升级为带契约和适用性证明的运营产品")
    thesis(
        s,
        "数据可访问不等于可复用，更不等于可进入模型或闭环。持续供给需要Owner、契约、版本、成本和退出机制；机器消费还必须证明数据对当前对象、时空、配置和任务适用，并在条件变化时重新验收。",
    )

    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.16), Inches(6.20), Inches(5.00))
    lined(left, WHITE, LINE)
    chip(s, 0.52, 1.26, 0.70, 0.22, "T1  高", TEAL)
    add_text(s, 1.30, 1.24, 5.10, 0.26, ["网络数据将被封装为可运营的数据产品"], size=13, color=NAVY, bold=True)
    add_text(
        s,
        0.52,
        1.54,
        5.90,
        0.52,
        ["最低定义：明确消费者与任务、稳定标识和版本、Owner、Schema/语义、任务适用性、SLO、用途/权限、血缘、成本/计量、保留期、降级与退出机制；同时给出变更通知和兼容策略。只有目录条目、接口包装或Marketplace页面，仍不是数据产品。"],
        size=9.5,
        color=SLATE,
        spacing=1.03,
    )
    t1_blocks = [
        ("成立逻辑", "源和消费者同时增加，点对点集成按“源×消费者”增长；同一数据若被重复采集、清洗和解释，边际成本不会下降。稳定契约是跨场景复用前提；无Owner就无法管理变更、SLO、成本、争议和退出。"),
        ("落地重点", "优先产品化故障、节能、移动性、QoE和模型特征等高复用数据族。默认交付事件、聚合指标、特征或分析服务，而非全量原始流；每项产品绑定权威源、用途、消费者和责任人。"),
        ("验收口径", "以持续消费者、第二场景/后端复用、交付周期、契约违规、重复采集下降和单位有效任务成本验收；同时核算生产、存储、跨域传输和退出成本，不以目录、接口或PB数量验收。"),
        ("适用边界", "产品化不等于上架目录。无线数据必须携带小区、频段、波束、配置、采样窗口和测量条件，不能只交付字段；长期无人消费、无版本承诺或无法退出的条目应下架。"),
    ]
    for i, (title, body) in enumerate(t1_blocks):
        y = 2.14 + i * 0.96
        add_text(s, 0.52, y, 5.90, 0.18, [title], size=10, color=TEAL, bold=True)
        add_text(s, 0.52, y + 0.18, 5.90, 0.72, [body], size=9, color=SLATE, spacing=1.02)

    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.74), Inches(1.16), Inches(6.19), Inches(5.00))
    lined(right, WHITE, LINE)
    chip(s, 6.88, 1.26, 1.05, 0.22, "T6  中高", AMBER)
    add_text(s, 8.02, 1.24, 4.75, 0.26, ["任务适用性将成为数据进入模型与闭环的准入标准"], size=13, color=NAVY, bold=True)
    add_text(
        s,
        6.88,
        1.54,
        5.89,
        0.52,
        ["任务适用性（Data Fitness for Use）不是数据的统一总分或既定3GPP指标，而是可测试的“适用性信封”：说明数据在何种对象、时空、配置、采样和误差条件下，可被哪个模型、Agent或闭环安全使用；条件变化后必须重新验证。"],
        size=9.5,
        color=SLATE,
        spacing=1.03,
    )
    envelopes = [
        (TEAL, "通用质量", "准确、完整、一致、新鲜、唯一、可追溯；明确缺失、重复、异常值和来源可信度。"),
        (NAVY, "通信质量", "采样代表性、时空覆盖、同步误差、配置一致、测量置信、suspect/incomplete；区分网络状态与采集链路误差。"),
        (AMBER, "AI质量", "标签来源、训练—服务分布偏差、特征漂移、适用对象、模型版本和禁用边界；避免历史高质量掩盖当前失配。"),
        (CORAL, "准入方式", "按用途和动作风险设阈值；未达则降级、补采、阻断或重验，并与模型卡、产品契约、策略审批和动作放行挂钩。"),
    ]
    for i, (accent, title, body) in enumerate(envelopes):
        y = 2.14 + i * 0.72
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.88), Inches(y), Inches(5.89), Inches(0.66))
        lined(box, WHITE, accent)
        add_text(s, 7.02, y + 0.04, 5.61, 0.18, [title], size=10, color=NAVY, bold=True)
        add_text(s, 7.02, y + 0.24, 5.61, 0.36, [body], size=9, color=SLATE, spacing=1.02)
    add_text(
        s,
        6.88,
        5.06,
        5.89,
        0.96,
        [
            "边界：字段完整不等于语义可比；样本量大不等于具有代表性；离线合格不等于当前网络状态可用。不能用一个总分评价测量、告警、位置、感知和训练特征。术语上不以QoD指代数据质量；QoD仅指CAMARA Quality on Demand API。",
        ],
        size=8.5,
        color=MUTED,
        spacing=1.03,
    )

    evidence_bar(
        s,
        [
            ("领先信号", "T1：契约进入版本基线，生产者与消费者共同签署SLO/适用性并管理退出。T6：适用性进入契约、模型卡、接口、重验触发器和动作放行条件。", TEAL),
            ("验证指标", "T1：交付周期、复用消费者、第二场景、违规、重复采集下降、单位任务成本。T6：覆盖/代表性、同步误差、配置一致、漂移、误动作、降级与重验触发。", NAVY),
            ("反证 / 证据", "反证：每场景重复取数清洗、产品长期无人消费，或只报完整性/接口成功率。证据：Nokia Data Suite、Ericsson Data Mesh、Huawei Data QoS仅支持机制方向，非统一标准。", CORAL),
        ],
        y=6.26,
        h=0.88,
    )
    footer(s, 3, "T1 / T6", "[S10][S14][S24]")


def slide_p4(prs) -> None:
    s = new_slide(prs)
    header(s, "T4 / T5  ·  可信自动化趋势", "Agent协同与NDT动作验证共同构成高风险自动化的可信准入机制")
    thesis(
        s,
        "数据进入真实网络控制后，错误会转化为资源、SLA和安全风险。可信自动化必须把供数、决策、验证、执行和结果串成可审计责任链，同时分离数据、模型、Agent和域控制器职责；NDT提供反事实证据，但不能替代授权和责任审批。",
    )

    steps = [
        ("供数", "版本化产品、上下文、任务适用性、用途与策略"),
        ("决策", "模型/Agent记录输入、推理、工具和候选动作"),
        ("验证", "NDT预演、策略冲突、风险预算与审批"),
        ("执行", "身份、作用域、租约、白名单、灰度与熔断"),
        ("回证", "结果、KPI、副作用、漂移、接管与回滚"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.38 + i * 2.59
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.16), Inches(2.45), Inches(0.78))
        rgb(box, NAVY)
        add_text(s, x + 0.08, 1.20, 2.29, 0.22, [f"{i + 1}  {title}"], size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.08, 1.44, 2.29, 0.44, [body], size=9, color=WHITE, align=PP_ALIGN.CENTER, spacing=1.02)
        if i < 4:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.45), Inches(1.44), Inches(0.14), Inches(0.16))
            rgb(arr, TEAL)

    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(2.06), Inches(6.20), Inches(4.10))
    lined(left, WHITE, LINE)
    chip(s, 0.52, 2.16, 1.05, 0.22, "T4  中高", AMBER)
    add_text(s, 1.66, 2.14, 4.76, 0.26, ["形成“供数—决策—执行—回证”协同闭环"], size=13, color=NAVY, bold=True)
    add_text(
        s,
        0.52,
        2.44,
        5.90,
        0.46,
        ["Agent若处理真实网络任务，必须持续获得当前对象、配置、任务适用性、权限和作用域，而不是一次性把文档放进向量库。一旦能调用工具，数据错误、模型漂移和过期策略都会转化为动作风险。"],
        size=9.5,
        color=SLATE,
        spacing=1.03,
    )
    t4 = [
        ("融合内容", "目录检索、订阅、血缘、任务适用性和策略成为受控工具；数据、模型、提示、工具、动作和结果共享版本与Correlation ID；审计证据跨工具、域和人工接管过程传播。"),
        ("责任分离", "数据Owner管用途、质量、保留和删除；模型Owner管性能、偏差和漂移；Agent管身份、权限和动作；域控制器保留实时状态与执行责任。Reasoner与Actor隔离，无边界“大脑”不可审计。"),
        ("落地顺序", "只读检索与调查 → 决策建议与工单编排 → 历史回放/影子运行 → 白名单低风险动作 → 高风险动作受控准入。每级明确证据门槛、超时降级和人工接管；快环保留确定性执行。"),
    ]
    for i, (title, body) in enumerate(t4):
        y = 2.96 + i * 1.00
        add_text(s, 0.52, y, 5.90, 0.18, [title], size=10, color=CORAL, bold=True)
        add_text(s, 0.52, y + 0.18, 5.90, 0.76, [body], size=9, color=SLATE, spacing=1.02)

    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.74), Inches(2.06), Inches(6.19), Inches(4.10))
    lined(right, WHITE, LINE)
    chip(s, 6.88, 2.16, 1.05, 0.22, "T5  中高", AMBER)
    add_text(s, 8.02, 2.14, 4.75, 0.26, ["NDT预演成为高风险动作的重要放行证据"], size=13, color=NAVY, bold=True)
    add_text(
        s,
        6.88,
        2.44,
        5.89,
        0.52,
        ["Data Fabric向NDT提供带时间、配置、任务适用性、策略和血缘的版本化状态快照；NDT比较执行、不执行及替代动作，输出收益、冲突、副作用和置信边界；策略引擎据此放行、灰度或拒绝，生产结果再回流校准场景与误差。"],
        size=9.5,
        color=SLATE,
        spacing=1.03,
    )
    t5 = [
        ("有效证据条件", "按动作类别定义保真度、最大同步滞后、场景覆盖、可接受误差和未知状态；保存输入快照、拓扑/配置、孪生与模型版本、候选动作、预演结果、放行理由和生产结果。"),
        ("超限处理", "误差、滞后或场景覆盖超门槛时，自动缩小作用域、降低风险预算或退回建议/人工模式；不得用“孪生通过”替代责任审批，也不能把离线准确率直接当生产安全。"),
        ("适用边界", "NDT是证据链中的消费者和生产者，不是Fabric同义词，也不是绝对安全证明。必须与策略、灰度、熔断、回滚、人工接管和生产校准组合；低风险动作可采用更轻量证据。"),
    ]
    for i, (title, body) in enumerate(t5):
        y = 3.02 + i * 0.98
        add_text(s, 6.88, y, 5.89, 0.18, [title], size=10, color=CORAL, bold=True)
        add_text(s, 6.88, y + 0.18, 5.89, 0.74, [body], size=9, color=SLATE, spacing=1.02)

    evidence_bar(
        s,
        [
            ("领先信号", "Agent工具调用携带身份、作用域、租约、产品/模型/策略版本和Correlation ID；NDT进入变更审批、影子运行和动作验证，而不是独立演示平台。", TEAL),
            ("验证指标", "证据完整、越权拦截、错误上下文阻断、人工接管、回滚成功；同步滞后、预测误差、场景覆盖、回放一致、灰度失败、副作用与校准偏差。", NAVY),
            ("反证 / 证据", "反证：只有聊天/RAG或工单建议；工具调用无法追责；NDT只有拓扑可视化、无误差门槛和生产反馈。证据：3GPP AI生命周期、SA5 NDT管理及厂商自治平台仅证明积木正在形成。", CORAL),
        ],
        y=6.26,
        h=0.88,
    )
    footer(s, 4, "T4 / T5", "[S4][S23][S28][S29]")


def slide_p5(prs) -> None:
    s = new_slide(prs)
    header(s, "T7  ·  责任边界与商业分化", "商业化将沿责任边界分化：能力服务出售可组合输入，结果服务承诺可归因输出")
    thesis(
        s,
        "分叉点不是数据形态，而是买方还是供应方承担决策、动作与结果责任。稳定产品、语义、任务适用性和计量只决定服务能否履约，并不解释为什么分化；客户集成能力、风险偏好和责任转移需求决定交易结构。",
    )

    shared = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.16), Inches(12.55), Inches(0.78))
    lined(shared, SOFT_TEAL, TEAL)
    add_text(s, 0.52, 1.20, 12.27, 0.20, ["分化机制  ·  决策权 × 动作权 × 结果责任"], size=11, color=NAVY, bold=True)
    add_text(
        s,
        0.52,
        1.42,
        12.27,
        0.44,
        ["有集成能力且保留控制权的买方购买可组合能力，形成A；希望转移持续运营与结果责任的客户购买托管结果，形成B。接口标准化推动A跨网复制，也降低单项能力稀缺性；归因、授权和责任风险限制B规模化，因此两者长期并存而非先后替代。"],
        size=9.5,
        color=SLATE,
        spacing=1.03,
    )

    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(2.06), Inches(6.20), Inches(4.10))
    lined(left, WHITE, LINE)
    chip(s, 0.52, 2.16, 1.80, 0.22, "路径A  ·  中高确定性", GREEN)
    add_text(s, 2.42, 2.14, 4.00, 0.26, ["数据/API与能力服务"], size=14, color=NAVY, bold=True)
    a_rows = [
        ("责任与控制", "买方保留业务决策、能力组合和网络动作权；供应方只保证接口行为、可用性、SLA、用途与撤销，不承担最终业务结果。故障时按能力SLA赔付，不对下游业务全责。"),
        ("交易对象 / 买方", "Quality on Demand（QoD）、位置、身份、状态、分析或网络能力；面向有集成能力的行业平台、开发者、CPaaS与聚合渠道。原始感知流不是默认商品。"),
        ("扩张与计价", "聚合层解决跨网发现、认证、覆盖、行为一致、用量和清结算；按调用、会话、覆盖、用量或订阅收费。接口越标准，复制越快，但差异化和单价可能下降。"),
        ("验证 / 反证", "验证跨网一致、开发者复用、活跃付费客户、收入、履约成本、失败率和续费。只有API数量、一次性试用或补贴调用而无持续收入即不成立。"),
    ]
    for i, (title, body) in enumerate(a_rows):
        y = 2.48 + i * 0.88
        add_text(s, 0.52, y, 5.90, 0.18, [title], size=10, color=GREEN, bold=True)
        add_text(s, 0.52, y + 0.18, 5.90, 0.66, [body], size=9.5, color=SLATE, spacing=1.02)

    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.74), Inches(2.06), Inches(6.19), Inches(4.10))
    lined(right, WHITE, LINE)
    chip(s, 6.88, 2.16, 1.90, 0.22, "路径B  ·  低—中确定性", DEEP_GREEN)
    add_text(s, 8.86, 2.14, 3.91, 0.26, ["结果型托管/自治服务"], size=14, color=NAVY, bold=True)
    b_rows = [
        ("责任与控制", "供应方参与数据选择、策略决策和网络动作，并对约定结果负责；客户购买效果与运营责任转移，而不是单项接口。合同必须限定作用域、可控因素和责任上限。"),
        ("交易对象 / 买方", "体验、交易成功率、节能、故障率、运营成本或风险改善；面向缺少持续集成/运营能力、愿意转移部分责任的金融、车联、工业和媒体客户。"),
        ("扩张与计价", "需T4/T5、结果基线、对照组/观测窗口、因果归因、动作授权、风险上限和回滚；按托管费、业务量、效果增益、风险降低或结果分成。"),
        ("验证 / 反证", "验证归因置信、净效果、误动作、风险暴露、责任争议、人工接管和回滚成本。无法定义基线、排除外因、定责或回退即不能规模化。"),
    ]
    for i, (title, body) in enumerate(b_rows):
        y = 2.48 + i * 0.88
        add_text(s, 6.88, y, 5.89, 0.18, [title], size=10, color=DEEP_GREEN, bold=True)
        add_text(s, 6.88, y + 0.18, 5.89, 0.66, [body], size=9.5, color=SLATE, spacing=1.02)

    evidence_bar(
        s,
        [
            ("长期共存", "有集成能力的客户购买能力，希望转移运营责任的客户购买结果；A和B不是成熟度先后。A可成为B的输入，B也可把部分能力重新开放，但收入、责任、风险和成本必须分账。", TEAL),
            ("证据状态", "A：CAMARA/Open Gateway、聚合渠道和中国移动Quality on Demand（QoD）已有具名商业信号。B：stc/KDDI/Indosat证明域级自动化有效，但公开结果定价、责任转移和规模分成合同不足。", NAVY),
            ("共同前提 / 边界", "内部产品稳定性、语义、适用性、权限和计量只决定能否履约，不是分化原因。原始感知数据出售不是默认模式；域级性能改善不能自动外推结果型合同成立。", CORAL),
        ],
        y=6.26,
        h=0.88,
    )
    footer(s, 5, "T7", "[S16][S18][S55]")


def slide_p6(prs) -> None:
    s = new_slide(prs)
    header(s, "T8  ·  标准承载趋势", "标准更可能冻结“能力与接口”，而非唯一专用数据面网元")
    thesis(
        s,
        "T1–T7的需求会留下，但承载它们的拓扑、组件边界和名称未必统一。标准应先冻结跨厂商必须一致的角色、对象、能力语义、生命周期、接口行为和测试；实现可继续由既有功能增强、管理服务、边缘Agent或新NF组合。",
    )

    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(1.16), Inches(6.20), Inches(3.62))
    lined(left, WHITE, LINE)
    add_text(s, 0.52, 1.24, 5.90, 0.24, ["标准优先固化"], size=13, color=NAVY, bold=True)
    freeze = [
        "生产者、消费者、处理者角色及其责任与错误返回",
        "稳定对象标识、权威源、对象关系、版本和信息模型",
        "任务适用性公共属性、单位、作用域、置信和表达方式",
        "发现、处理、存储、订阅、交付和暴露的外部行为",
        "生命周期、授权、撤销、审计和证据引用关系",
        "跨厂商一致性、故障、降级、迁移和替换测试",
    ]
    for i, item in enumerate(freeze):
        y = 1.54 + i * 0.50
        chip(s, 0.52, y, 0.28, 0.22, str(i + 1), TEAL, WHITE, 8)
        add_text(s, 0.90, y, 5.52, 0.42, [item], size=11, color=SLATE)

    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.74), Inches(1.16), Inches(6.19), Inches(3.62))
    lined(right, WHITE, LINE)
    add_text(s, 6.88, 1.24, 5.89, 0.24, ["实现保持开放"], size=13, color=NAVY, bold=True)
    open_items = [
        (TEAL, "增强既有功能", "扩展DCCF、ADRF、NWDAF、MDA等SBI与管理服务，复用既有鉴权、容灾和运维体系；适合渐进演进。"),
        (NAVY, "管理域/边缘承载", "以SMO服务、边缘Agent或近源运行时组合部署；按时标和故障域拆分，不把远端平台放进快环同步路径。"),
        (AMBER, "可插拔多后端", "事件、流、文件、特征、模型和派生数据可选择不同运行时；核心元数据、契约与证据应可导出迁移。"),
        (CORAL, "必要时新增NF", "仅当复用造成隔离、吞吐、移动性、状态或职责不可分时引入专用功能，并证明收益覆盖新增复杂度。"),
    ]
    for i, (accent, title, body) in enumerate(open_items):
        y = 1.56 + i * 0.76
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.88), Inches(y), Inches(5.89), Inches(0.70))
        lined(box, WHITE, accent)
        add_text(s, 7.02, y + 0.05, 5.61, 0.20, [title], size=11, color=NAVY, bold=True)
        add_text(s, 7.02, y + 0.28, 5.61, 0.36, [body], size=10, color=SLATE)

    logic = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(4.90), Inches(12.55), Inches(1.26))
    lined(logic, SOFT_NAVY, NAVY)
    add_text(s, 0.52, 4.96, 12.27, 0.20, ["成立逻辑与产品策略"], size=11, color=NAVY, bold=True)
    add_text(
        s,
        0.52,
        5.18,
        12.27,
        0.88,
        [
            "发现、采集、处理、存储、暴露和治理跨越不同域、时标、故障域与主权边界，不同用例不可能共享完全相同的在线流程。标准应优先冻结可观察的外部行为；内部拓扑写得越死，越难适配存量、边缘资源和多厂商演进。",
            "产品跟踪能力、语义、任务适用性、接口和测试，不押注厂商网元名。若独立平面增加时延、状态、容灾和重复治理，应退回能力增强；若复用导致职责、隔离或性能不可分，再考虑新功能并设置退出门。",
        ],
        size=10,
        color=SLATE,
        spacing=1.04,
    )

    evidence_bar(
        s,
        [
            ("领先信号 / 指标", "Rel-21收敛角色、能力、信息模型、接口和一致性测试，同时允许多种部署。指标：跨厂商互操作、复用、故障降级、额外时延/状态、迁移、后端与运行时替换成本。", TEAL),
            ("反证 / 证据", "反证：3GPP冻结独立统一数据平面及必选NF/流程，或既有功能在相同约束下无法满足。现有SA2 KI#21、SA5 DMFW及Huawei原型只支持能力研究，未证明唯一拓扑。", CORAL),
            ("第三章收束", "长期控制点集中于语义、任务适用性、近源运行、动作证据、责任合同和跨网结算。平台地位取决于外部契约与互操作，而不是单一名称或数据集中规模。", NAVY),
        ],
        y=6.26,
        h=0.88,
    )
    footer(s, 6, "T8", "[S5][S12][S26]")


def count_cjk(prs: Presentation) -> list[int]:
    counts = []
    for slide in prs.slides:
        texts = []
        for shp in slide.shapes:
            if shp.has_text_frame:
                texts.append(shp.text_frame.text)
            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        counts.append(len(re.findall(r"[\u4e00-\u9fff]", "".join(texts))))
    return counts


def verify(prs: Presentation) -> None:
    n = len(prs.slides)
    if n != TOTAL:
        raise SystemExit(f"expected {TOTAL} slides, got {n}")
    counts = count_cjk(prs)
    print("cjk chars:", counts, "total", sum(counts))
    if min(counts) < 280:
        raise SystemExit(f"slide too sparse: {counts}")

    overflow = []
    tables = []
    for i, slide in enumerate(prs.slides, 1):
        table_n = 0
        for shp in slide.shapes:
            if shp.has_table:
                table_n += 1
            right = (shp.left + shp.width) / EMU
            bottom = (shp.top + shp.height) / EMU
            if right > W_IN + 0.03 or bottom > H_IN + 0.03 or shp.left < Emu(-0.02 * EMU) or shp.top < Emu(-0.02 * EMU):
                overflow.append((i, getattr(shp, "name", "?"), round(right, 3), round(bottom, 3)))
        tables.append(table_n)
    if overflow:
        raise SystemExit(f"shapes overflow canvas: {overflow}")
    if tables[0] != 1:
        raise SystemExit(f"P1 must contain 1 table, got {tables}")
    required = [
        "T1",
        "T2",
        "T3",
        "T4",
        "T5",
        "T6",
        "T7",
        "T8",
        "可运营",
        "变化事件",
        "近源执行",
        "任务适用性",
        "供数",
        "NDT",
        "能力",
        "结果",
        "能力与接口",
    ]
    all_text = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                all_text.append(shp.text_frame.text)
            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        all_text.append(cell.text)
    blob = "\n".join(all_text)
    missing = [item for item in required if item not in blob]
    if missing:
        raise SystemExit(f"missing required terms: {missing}")
    print("verify ok:", {"slides": n, "tables": tables})


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_p1(prs)
    slide_p2(prs)
    slide_p3(prs)
    slide_p4(prs)
    slide_p5(prs)
    slide_p6(prs)
    verify(prs)
    prs.save(OUT)
    print(f"wrote {OUT}")
    return OUT


if __name__ == "__main__":
    build()
