#!/usr/bin/env python3
"""Generate the 6-slide sections 5-7 decision deck.

Output: reports/6g-data-fabric-sections-5-7-deck.pptx
Source: reports/plan-sections5-7.md and chapters 5-7 of the full-sections HTML.

All visible copy is written into AutoShape / table-cell text frames.
Standalone text boxes are not used, so cards remain directly editable in PowerPoint.
"""

from __future__ import annotations

import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "reports" / "6g-data-fabric-sections-5-7-deck.pptx"

W_IN, H_IN = 13.333, 7.5
W, H = Inches(W_IN), Inches(H_IN)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
CORAL = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xB8, 0x86, 0x0B)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF4, 0xF7, 0xF8)
PALE = RGBColor(0xE8, 0xF0, 0xF2)
LINE = RGBColor(0xC5, 0xD0, 0xD4)
SOFT_TEAL = RGBColor(0xD7, 0xEC, 0xEA)
SOFT_AMBER = RGBColor(0xF6, 0xED, 0xD0)
SOFT_NAVY = RGBColor(0xE6, 0xEC, 0xF2)
SOFT_CORAL = RGBColor(0xF8, 0xE6, 0xE4)
FONT = "WenQuanYi Micro Hei"
TOTAL = 6
EMU = 914400
CJK_MIN = [480, 500, 520, 600, 520, 650]
CJK_TOTAL_MIN = 3500
P4_CELL_CJK_MAX = 75


def _set_run_east_asia(run) -> None:
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT)


def rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def lined(shape, fill: RGBColor, border: RGBColor, width: float = 0.75) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(width)


def _set_anchor(tf, anchor) -> None:
    mapping = {
        MSO_ANCHOR.TOP: "t",
        MSO_ANCHOR.MIDDLE: "ctr",
        MSO_ANCHOR.BOTTOM: "b",
    }
    try:
        tf._txBody.bodyPr.set("anchor", mapping.get(anchor, "t"))
    except Exception:
        pass


def fill_shape(
    shape,
    lines,
    *,
    default_size=9.5,
    default_color=SLATE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.04,
    ml=0.10,
    mr=0.08,
    mt=0.06,
    mb=0.05,
) -> None:
    """Write all copy into a shape's own text frame (never a detached text box)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    _set_anchor(tf, anchor)
    tf.clear()
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color, palign, after, pspacing = (
                item,
                default_size,
                False,
                default_color,
                align,
                2,
                spacing,
            )
        elif isinstance(item, dict):
            text = item["text"]
            size = item.get("size", default_size)
            bold = item.get("bold", False)
            color = item.get("color", default_color)
            palign = item.get("align", align)
            after = item.get("after", 2)
            pspacing = item.get("spacing", spacing)
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else default_color
            palign = item[4] if len(item) > 4 else align
            after = item[5] if len(item) > 5 else 2
            pspacing = spacing
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = palign
        p.space_before = Pt(0)
        p.space_after = Pt(after)
        p.line_spacing = pspacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
        _set_run_east_asia(run)


def box(
    slide,
    l,
    t,
    w,
    h,
    lines,
    *,
    fill=WHITE,
    border=LINE,
    rounded=True,
    radius=0.08,
    **fill_kw,
):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l),
        Inches(t),
        Inches(w),
        Inches(h),
    )
    if border is None:
        rgb(shp, fill)
    else:
        lined(shp, fill, border)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if lines:
        fill_shape(shp, lines, **fill_kw)
    return shp


def arrow_down(slide, l, t, w=0.16, h=0.14, fill=TEAL):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    rgb(shp, fill)
    return shp


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rgb(bg, WHITE)
    return s


def header(slide, number: str, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.07))
    rgb(bar, TEAL)
    box(
        slide,
        0.38,
        0.10,
        12.55,
        0.54,
        [
            {"text": number, "size": 10, "bold": True, "color": TEAL, "after": 1},
            {"text": title, "size": 20, "bold": True, "color": NAVY, "after": 0, "spacing": 1.02},
        ],
        fill=WHITE,
        border=None,
        rounded=False,
        ml=0.0,
        mr=0.0,
        mt=0.02,
        mb=0.0,
        spacing=1.02,
    )


def footer(slide, page: int, chapter: str, refs: str = "") -> None:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(7.18), Inches(12.55), Inches(0.012)
    )
    rgb(line, LINE)
    left = f"6G 数据架构 × 数据编织 · 第5–7章决策胶片 · {chapter}"
    if refs:
        left = f"{left}    {refs}"
    box(
        slide,
        0.38,
        7.22,
        10.55,
        0.22,
        [{"text": left, "size": 8, "color": MUTED, "after": 0}],
        fill=WHITE,
        border=None,
        rounded=False,
        ml=0.0,
        mr=0.04,
        mt=0.0,
        mb=0.0,
        default_color=MUTED,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    box(
        slide,
        11.10,
        7.22,
        1.83,
        0.22,
        [{"text": f"{page} / {TOTAL}", "size": 8, "bold": True, "color": MUTED, "align": PP_ALIGN.RIGHT, "after": 0}],
        fill=WHITE,
        border=None,
        rounded=False,
        ml=0.0,
        mr=0.0,
        mt=0.0,
        mb=0.0,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def thesis(slide, text: str, top: float = 0.66, height: float = 0.46) -> None:
    box(
        slide,
        0.38,
        top,
        12.55,
        height,
        [{"text": text, "size": 12, "bold": True, "color": NAVY, "after": 0, "spacing": 1.04}],
        fill=PALE,
        border=TEAL,
        radius=0.06,
        default_color=NAVY,
        ml=0.14,
        mr=0.12,
        mt=0.06,
        mb=0.05,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_table(slide, l, t, w, h, rows, col_w=None, font=8.5):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(l), Inches(t), Inches(w), Inches(h))
    table = table_shape.table
    if col_w:
        for i, cw in enumerate(col_w):
            table.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.TOP if r else MSO_ANCHOR.MIDDLE
            fill = NAVY if r == 0 else (OFF if r % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.04)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.03
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font)
            run.font.bold = r == 0 or c == 0
            run.font.color.rgb = WHITE if r == 0 else INK
            run.font.name = FONT
            _set_run_east_asia(run)
    return table


def slide_p1(prs) -> None:
    s = new_slide(prs)
    header(s, "5.1  ·  CURRENT POSITION", "中兴已覆盖网络执行、数据智能和Agent应用，但跨产品共用能力仍未闭合")
    thesis(
        s,
        "中兴不是只在网络执行与智能体之间提供数据，而是同时拥有底层网络触点、中间数据/模型/孪生资产和上层自有Agent应用。真正缺少的是跨产品共用的数据契约、分布式供数和动作证据基线。",
    )

    layers = [
        (
            1.16,
            SOFT_NAVY,
            NAVY,
            "自有智能体应用",
            "Fault Agent · Co-Sight · Copilot · 场景Agent",
            "已有域级应用、插件能力和量化试点；尚未证明不同Agent共享同一数据产品、任务适用性、策略版本和动作证据。中兴不仅为第三方供数，也提供自有智能体应用。",
            "判断：已有域级应用和量化案例，尚未证明跨Agent共享同一契约与证据。自有Agent既是产品，也是公共能力的首批消费者。",
        ),
        (
            2.52,
            SOFT_TEAL,
            TEAL,
            "数据与智能基础",
            "AIR Net数据引擎 · 模型引擎 · 数字孪生引擎 · KG",
            "公开材料已给出数据资产、治理、模型与孪生组件。建议沿可信数据产品、语义治理、质量洞察、跨系统共享与分布式处理演进，但尚不足以证明跨产品统一版本与责任基线。",
            "判断：组件存在，尚未形成跨产品统一版本与责任基线。公共能力应升级现有引擎，而不是另建通用湖仓或独立IT平台。",
        ),
        (
            3.88,
            SOFT_AMBER,
            AMBER,
            "网络执行与权威数据",
            "RAN · Core · Transport · OSS/SMO · NEF · CUDR",
            "掌握无线对象、配置、测量、会话、网管状态、网络API入口和域控制触点。ZXUN CUDR、NEF及主航道产品证明电信数据管理与能力暴露基础，但不能外推为跨域统一数据平台。",
            "判断：掌握对象、状态和动作触点，是相对云/IT平台最难替代的优势。近源执行和动作回证必须落在这一层，而不是远端目录。",
        ),
    ]
    for y, fill, accent, title, assets, body, judge in layers:
        box(
            s,
            0.38,
            y,
            8.72,
            1.22,
            [
                {"text": title, "size": 13, "bold": True, "color": NAVY, "after": 2},
                {"text": assets, "size": 10.5, "bold": True, "color": accent, "after": 3},
                {"text": body, "size": 9.5, "color": SLATE, "after": 0, "spacing": 1.03},
            ],
            fill=fill,
            border=accent,
            radius=0.06,
            mt=0.07,
            mb=0.05,
        )
        box(
            s,
            9.22,
            y,
            3.71,
            1.22,
            [
                {"text": "右侧判断", "size": 9, "bold": True, "color": accent, "after": 3},
                {"text": judge, "size": 9, "color": SLATE, "after": 0, "spacing": 1.03},
            ],
            fill=WHITE,
            border=accent,
            radius=0.06,
            mt=0.07,
        )
    arrow_down(s, 4.55, 2.38, 0.18, 0.12, TEAL)
    arrow_down(s, 4.55, 3.74, 0.18, 0.12, AMBER)

    bottom = [
        (
            0.38,
            TEAL,
            SOFT_TEAL,
            "强项",
            "网络触点、自智应用、闭环场景和API入口均有公开锚点。AIR Net三引擎、Fault Agent案例、NEF/CAMARA入口构成可升级基础，而不是从零建设数据平台。",
        ),
        (
            4.58,
            AMBER,
            SOFT_AMBER,
            "缺口",
            "跨产品契约、任务适用性、变化事件、近源执行/供数和统一动作证据尚未形成一套可复用基线。同一对象仍可能被无线、核心网、OSS和Agent重复解释。",
        ),
        (
            8.78,
            CORAL,
            SOFT_CORAL,
            "未证明",
            "多个组件存在，不等于完整Fabric、公司级统一中枢、跨厂商控制面或规模结果服务已经成立。也不能把域级技术效果直接写成结果型商业合同。",
        ),
    ]
    for x, accent, fill, title, body in bottom:
        box(
            s,
            x,
            5.18,
            4.15,
            0.90,
            [
                {"text": title, "size": 11, "bold": True, "color": accent, "after": 2},
                {"text": body, "size": 9, "color": SLATE, "after": 0, "spacing": 1.02},
            ],
            fill=fill,
            border=accent,
            radius=0.06,
            mt=0.05,
            mb=0.04,
        )

    questions = [
        ("诊断 1", "无线、核心网、OSS和Agent是否引用同一对象ID、权威源与配置版本？"),
        ("诊断 2", "同一数据产品能否被第二产品、第二Agent和第二后端复用，而无需大规模重写？"),
        ("诊断 3", "Agent使用的数据、模型、策略、动作和结果能否关联回放，并支持人工接管？"),
        ("诊断 4", "现有能力是进入正式版本和持续使用，还是仍依赖项目级定制与一次性对接？"),
    ]
    for i, (title, body) in enumerate(questions):
        x = 0.38 + i * 3.24
        box(
            s,
            x,
            6.16,
            3.14,
            0.94,
            [
                {"text": title, "size": 9, "bold": True, "color": NAVY, "after": 2},
                {"text": body, "size": 9, "color": SLATE, "after": 0, "spacing": 1.02},
            ],
            fill=OFF,
            border=LINE,
            radius=0.06,
            mt=0.05,
            mb=0.04,
        )
    footer(s, 1, "5.1 现状", "[S40][S41][S42][S48][S51][S53][S55]")


def slide_p2(prs) -> None:
    s = new_slide(prs)
    header(s, "5.2  ·  TARGET POSITION", "把AIR Net数据引擎及相邻资产升级为主航道共用的“电信数据与动作可信能力”")
    thesis(
        s,
        "公共能力不是独立IT平台，而是嵌入无线、核心网、OSS/SMO和Agent应用的横向组件。向下连接权威数据与域控制器，向上同时服务中兴自有Agent、第三方Agent、行业应用和开放API。",
        height=0.44,
    )

    chain = [
        (NAVY, "1  网络事实与动作触点", "RAN / Core / Transport / OSS/SMO 提供对象、状态、配置和动作入口；公共能力连接这些触点，不替代域控制器。"),
        (TEAL, "2  可信数据契约", "对象语义、数据产品、Owner/SLO、任务适用性、变化事件与退出，使跨产品消费同一责任基线。"),
        (TEAL, "3  分布式供数运行", "全局定义规则与版本；近源执行过滤、特征与受限处理；断链运行、版本同步和异步回证。"),
        (AMBER, "4  动作证据机制", "Agent身份与作用域、策略与风险预算、NDT预演、灰度、结果和回滚，形成可审计责任链。"),
        (NAVY, "5  自有Agent / 第三方Agent / 开放API", "先服务中兴自有智能体和主航道产品，再向第三方Agent、NEF/CAMARA和行业API开放。"),
    ]
    y = 1.16
    for i, (accent, title, body) in enumerate(chain):
        fill = SOFT_NAVY if accent == NAVY else (SOFT_AMBER if accent == AMBER else SOFT_TEAL)
        box(
            s,
            0.38,
            y,
            4.42,
            0.64,
            [
                {"text": title, "size": 11, "bold": True, "color": NAVY, "after": 1},
                {"text": body, "size": 9, "color": SLATE, "after": 0, "spacing": 1.02},
            ],
            fill=fill,
            border=accent,
            radius=0.06,
            mt=0.04,
            mb=0.03,
            ml=0.10,
        )
        if i < 4:
            arrow_down(s, 2.48, y + 0.64, 0.16, 0.10, accent)
        y += 0.76

    caps = [
        (
            TEAL,
            SOFT_TEAL,
            "数据契约",
            "对象语义、数据产品、Owner/SLO、任务适用性、变化事件与退出。无线、核心网、OSS与Agent围绕同一责任基线消费数据；产品线可扩展域字段，但不得复制核心ID。",
        ),
        (
            TEAL,
            SOFT_TEAL,
            "分布式供数",
            "全局定义、近源过滤/特征、断链运行、版本同步和异步回证。近源执行按时标和故障域就近完成处理；公共服务不得成为毫秒快环同步单点。",
        ),
        (
            AMBER,
            SOFT_AMBER,
            "动作证据",
            "Agent身份与作用域、策略与风险预算、NDT预演、灰度、结果和回滚。公共层定义证据Schema，域控制器保留具体动作权限，失败必须可接管、可回放。",
        ),
    ]
    for i, (accent, fill, title, body) in enumerate(caps):
        box(
            s,
            4.96,
            1.16 + i * 1.26,
            8.00,
            1.18,
            [
                {"text": title, "size": 13, "bold": True, "color": NAVY, "after": 3},
                {"text": body, "size": 10, "color": SLATE, "after": 0, "spacing": 1.03},
            ],
            fill=fill,
            border=accent,
            radius=0.06,
            mt=0.07,
        )

    duties = [
        ("公共能力责任", "核心对象、契约、公共接口、版本兼容和证据Schema；识别重复建设，维护互操作测试。", TEAL),
        ("产品线责任", "域语义扩展、实时状态、资源预算、具体动作和客户SLO；不得绕过公共对象与证据基线。", NAVY),
        ("Agent责任", "模型/工具选择、身份权限、动作申请和结果回证；无门控不得改生产契约或高风险策略。", AMBER),
        ("验证门", "跨两产品/Agent复用、两种后端、断链回退、动作证据完整，以及开放接口可替换。", CORAL),
    ]
    for i, (title, body, accent) in enumerate(duties):
        x = 0.38 + i * 3.24
        box(
            s,
            x,
            5.04,
            3.14,
            0.92,
            [
                {"text": title, "size": 11, "bold": True, "color": accent, "after": 3},
                {"text": body, "size": 9.5, "color": SLATE, "after": 0, "spacing": 1.03},
            ],
            fill=WHITE,
            border=accent,
            radius=0.06,
            mt=0.07,
        )

    box(
        s,
        0.38,
        6.04,
        12.55,
        1.06,
        [
            {"text": "目标边界：不把建议写成现有SKU", "size": 11, "bold": True, "color": CORAL, "after": 3},
            {
                "text": "不建设通用湖仓/Catalog；不替代域控制器；不进入毫秒快环同步单点；不绑定单一模型、后端或渠道；不把目标架构写成现有SKU或客户承诺。电信数据与动作可信能力是战略工作定位，不是正式产品名或标准网元。向上出口包括中兴自有Agent、第三方Agent和开放API，自有Agent不是第三方生态的附属验证工具。",
                "size": 10,
                "color": SLATE,
                "after": 0,
                "spacing": 1.03,
            },
        ],
        fill=SOFT_CORAL,
        border=CORAL,
        radius=0.06,
        mt=0.07,
    )
    footer(s, 2, "5.2 目标定位", "[S40][S41][S43][S48][S54]")


def slide_p3(prs) -> None:
    s = new_slide(prs)
    header(s, "MODE ANALYSIS  ·  ONE HORIZONTAL / MANY VERTICALS", "横向共用三项可信能力，纵向嵌入主航道产品和Agent应用")
    thesis(
        s,
        "中兴不需要另建一个统一平台，而应让无线、核心网、OSS/SMO和自有Agent复用同一数据契约、供数规则和动作证据接口；各产品继续保留场景、实时状态、域内执行和客户结果责任。",
        height=0.42,
    )

    box(
        s,
        0.38,
        1.14,
        12.55,
        0.24,
        [{"text": "横向公共能力：统一理解、版本、接口和证据，而不是统一全部存储、运行时与动作", "size": 10.5, "bold": True, "color": WHITE, "after": 0, "align": PP_ALIGN.CENTER}],
        fill=TEAL,
        border=None,
        radius=0.04,
        mt=0.02,
        mb=0.01,
        ml=0.08,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    commons = [
        (
            "可信数据契约",
            "语义 / 产品 / 任务适用性",
            "核心对象、权威源、语义和对象关系；数据产品、Owner/SLO、任务适用性；版本、变化事件、用途、保留和退出。产品线允许域扩展，但不能复制核心ID和责任口径。",
        ),
        (
            "分布式供数规则",
            "订阅 / 版本 / 近源执行",
            "全局发布、订阅、缓存和版本基线；近源执行过滤、聚合、特征和资源预算；断链运行、异步回证、对账和回退。公共服务不得成为毫秒快环同步单点。",
        ),
        (
            "动作证据接口",
            "身份 / 策略 / 结果 / 回滚",
            "Agent身份、委托、作用域和租约；数据/模型/策略版本与候选动作；审批、执行结果、副作用、人工接管和回滚。公共层定义证据Schema，域控制器保留具体动作权限。",
        ),
    ]
    for i, (title, sub, body) in enumerate(commons):
        x = 0.38 + i * 4.20
        box(
            s,
            x,
            1.42,
            4.08,
            1.58,
            [
                {"text": title, "size": 13, "bold": True, "color": NAVY, "after": 1},
                {"text": sub, "size": 9.5, "bold": True, "color": TEAL, "after": 3},
                {"text": body, "size": 9, "color": SLATE, "after": 0, "spacing": 1.02},
            ],
            fill=SOFT_TEAL,
            border=TEAL,
            radius=0.06,
            mt=0.06,
            mb=0.04,
        )
        arrow_down(s, x + 1.96, 3.02, 0.16, 0.14, TEAL)

    box(
        s,
        0.38,
        3.18,
        12.55,
        0.24,
        [{"text": "纵向嵌入：主航道产品与中兴自有Agent保留场景责任，第三方Agent/API按成熟度接入同一接口", "size": 10.5, "bold": True, "color": WHITE, "after": 0, "align": PP_ALIGN.CENTER}],
        fill=NAVY,
        border=None,
        radius=0.04,
        mt=0.02,
        mb=0.01,
        ml=0.08,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    verticals = [
        (SOFT_NAVY, NAVY, "无线产品", "无线对象、近源数据、实时状态和域动作"),
        (SOFT_NAVY, NAVY, "核心网产品", "会话、策略、NWDAF/NEF和核心网动作"),
        (SOFT_NAVY, NAVY, "OSS / SMO", "配置、拓扑、Non-RT编排和跨域保证"),
        (SOFT_TEAL, TEAL, "中兴自有Agent", "Fault、节能、优化、运维等正式场景产品"),
        (SOFT_AMBER, AMBER, "第三方Agent / API", "受控数据产品、工具接口和开放能力"),
    ]
    for i, (fill, accent, title, body) in enumerate(verticals):
        x = 0.38 + i * 2.52
        box(
            s,
            x,
            3.46,
            2.42,
            0.72,
            [
                {"text": title, "size": 11, "bold": True, "color": NAVY, "after": 1},
                {"text": body, "size": 9, "color": SLATE, "after": 0, "spacing": 1.02},
            ],
            fill=fill,
            border=accent,
            radius=0.06,
            mt=0.06,
            mb=0.03,
        )

    owners = [
        (
            TEAL,
            SOFT_TEAL,
            "公共能力Owner",
            "管理核心对象、契约、版本、公共接口和互操作测试；识别重复建设；维护证据Schema和工具调用约束；统计跨产品复用和公共能力生命周期。公共团队不承担所有客户场景和网络动作责任。",
        ),
        (
            NAVY,
            SOFT_NAVY,
            "产品与Agent Owner",
            "管理域语义扩展、实时状态和资源预算；对场景数据产品、模型/Agent性能和客户SLO负责；管理具体动作授权、灰度、人工接管和回滚；将质量、漂移、异常和动作结果反馈到公共基线。",
        ),
    ]
    for i, (accent, fill, title, body) in enumerate(owners):
        x = 0.38 + i * 6.30
        box(
            s,
            x,
            4.26,
            6.18,
            1.08,
            [
                {"text": title, "size": 12, "bold": True, "color": NAVY, "after": 2},
                {"text": body, "size": 9.5, "color": SLATE, "after": 0, "spacing": 1.03},
            ],
            fill=fill,
            border=accent,
            radius=0.06,
            mt=0.06,
        )

    checks = [
        ("复用", "同一契约、工具和证据接口进入两条主航道或两类Agent，第二场景无需大规模重写。", TEAL),
        ("运行", "中心不可达时域内继续运行，恢复后完成对账、证据回传和冲突处理。", NAVY),
        ("责任", "公共团队不承担所有客户场景，产品线不得绕过公共对象、契约和证据基线。", AMBER),
        ("反馈", "结果只触发影响分析、重验、审批和版本发布，不允许Agent无门控改规则。", CORAL),
    ]
    for i, (title, body, accent) in enumerate(checks):
        x = 0.38 + i * 3.24
        box(
            s,
            x,
            5.42,
            3.14,
            0.78,
            [
                {"text": title, "size": 11, "bold": True, "color": accent, "after": 1},
                {"text": body, "size": 9, "color": SLATE, "after": 0, "spacing": 1.02},
            ],
            fill=WHITE,
            border=accent,
            radius=0.06,
            mt=0.05,
            mb=0.03,
        )

    box(
        s,
        0.38,
        6.28,
        12.55,
        0.82,
        [
            {"text": "模式停止线", "size": 11, "bold": True, "color": CORAL, "after": 2},
            {
                "text": "如果公共能力不能跨第二产品或第二Agent复用，或引入后增加快环依赖、模糊动作责任、阻碍产品版本演进，应退回域内能力；不能为了统一强迫所有产品共享同一存储、模型或在线处理流程。商业路径不是本页主轴：能力API或结果型服务只在后续扩展中按成熟度选择。",
                "size": 9.5,
                "color": SLATE,
                "after": 0,
                "spacing": 1.02,
            },
        ],
        fill=SOFT_CORAL,
        border=CORAL,
        radius=0.06,
        mt=0.06,
        mb=0.04,
    )
    footer(s, 3, "6.1 模式分析", "[S40][S41][S42][S48][S53]")


def slide_p4(prs) -> None:
    s = new_slide(prs)
    header(s, "6.2  ·  OPPORTUNITY PACKAGES", "机会集中在三项公共能力和一类智能体应用产品")
    thesis(
        s,
        "三项公共能力解决跨产品重复建设，自有Agent把公共能力转成真实场景和客户结果。四类机会共用对象模型、契约、身份、版本和证据Schema，不分别建设平台。",
        height=0.40,
    )
    add_table(
        s,
        0.38,
        1.14,
        12.55,
        5.04,
        [
            ["机会", "客户/产品任务", "最小交付", "中兴优势", "验证门", "停止线"],
            [
                "可信数据契约",
                "Agent与产品须知道拿到什么、由谁负责、当前任务是否适用，否则无法复用或放行。",
                "对象语义、数据产品、Owner/SLO、任务适用性、变化事件、用途、版本与退出。",
                "网元、配置、测量与场景知识；AIR Net数据引擎与知识图谱资产。",
                "两产品线、两类Agent或两种后端复用；版本兼容、退出可执行。",
                "第二场景仍需大规模重写，或数据产品长期没有真实消费者。",
            ],
            [
                "分布式供数运行",
                "高频数据少搬运、近源执行过滤与特征，断链时仍运行并向全局异步回证。",
                "全局规则、近源过滤/特征、缓存、版本、资源预算、对账、断链与异步证据。",
                "RAN/Edge触点、网管运行工程和域内实时状态，具备近源处理条件。",
                "P99、资源预算、字节减量、断链可用、同步、证据完整和回退成功。",
                "收益覆盖不了站点开销，或公共服务成为快环远端同步单点。",
            ],
            [
                "动作证据机制",
                "Agent获权后须解释为何放行、谁负责、失败如何接管，动作可回放定责。",
                "身份、作用域、租约、策略、NDT、灰度、结果、熔断、人工接管和回滚。",
                "自有Agent、数字孪生引擎和域控制器相邻资产，能把动作落到真实网元。",
                "越权拦截、误差门槛、证据完整、人工接管成功和回滚可执行。",
                "只有聊天/RAG或孪生可视化，动作无法回放、定责和退回。",
            ],
            [
                "自有Agent产品",
                "把公共能力转化为故障、节能、优化和运维场景，并形成持续使用与结果反馈。",
                "Fault Agent加第二类Agent；复用同一数据产品、工具、策略和证据接口。",
                "已有Agent产品、故障知识、网络动作触点和公开量化案例。",
                "第二Agent复用、交付周期下降、场景效果、证据完整和版本回流。",
                "每个Agent仍独立取数治理，或效果无法跨场景复制改进。",
            ],
        ],
        col_w=[1.58, 2.18, 2.42, 2.18, 2.18, 2.01],
        font=8.5,
    )
    notes = [
        ("公共能力不是独立销售单元", "三项能力是主航道和Agent产品的共用组件，不分别建设平台，也不以单独品牌替代现有产品版本体系。"),
        ("自有Agent是正式产品", "中兴自有Agent不是第三方生态的附属验证，而是正式场景产品与首批消费者，负责把契约和证据用进真实闭环。"),
        ("第三方与API是开放扩展", "第三方Agent和API不改变中兴必须掌握的电信约束和责任接口；Quality on Demand（QoD）等能力API属于后续出口。"),
    ]
    for i, (title, body) in enumerate(notes):
        x = 0.38 + i * 4.20
        box(
            s,
            x,
            6.26,
            4.08,
            0.84,
            [
                {"text": title, "size": 10.5, "bold": True, "color": NAVY, "after": 2},
                {"text": body, "size": 9, "color": SLATE, "after": 0, "spacing": 1.02},
            ],
            fill=OFF,
            border=LINE,
            radius=0.06,
            mt=0.05,
            mb=0.04,
        )
    footer(s, 4, "6.2 机会分析", "[S40][S41][S42][S43][S48][S53]")


def slide_p5(prs) -> None:
    s = new_slide(prs)
    header(s, "6.3  ·  PORTFOLIO CHOICE", "优先补齐共用基线和自有Agent复用，近源与动作证据按场景推进")
    thesis(
        s,
        "机会优先级首先看跨产品复用和网络差异化，其次看能否进入真实Agent场景并形成可运行证据，最后再选择第三方生态、能力API或结果型服务等外部扩展方式。",
        height=0.44,
    )

    cards = [
        (
            TEAL,
            SOFT_TEAL,
            "优先建设",
            "公共数据契约 + 自有Agent复用",
            [
                "统一核心对象、语义、数据产品、任务适用性、变化事件、Agent供数工具和动作证据接口，使无线、核心网、OSS与Agent不再各建一套解释。",
                "Fault Agent加一类非故障Agent成为真实消费者；公共能力进入至少两条主航道产品，而不是停留在演示或单项目对接。",
                "验收第二场景重写比例、两种后端可替换、Owner、真实消费者、交付周期和契约违规；无消费者的目录条目不计入成功。",
                "这一步决定后续所有扩展是否有可复制基线：先闭合横向公共能力，再讨论近源规模和外部接口。",
            ],
        ),
        (
            AMBER,
            SOFT_AMBER,
            "时间盒验证",
            "近源供数 + 动作证据",
            [
                "选择一种高频无线数据和一个真实Agent场景，验证近源执行中的过滤/特征、资源预算、断链、同步、证据与回退，而不是先铺全网。",
                "同步验证身份、策略、NDT误差、灰度、人工接管和回滚；只有聊天/RAG或孪生可视化，不能算动作证据成立。",
                "收益不覆盖站点开销、无法标定误差、或必须把远端目录/Agent放入快环关键路径时停止，退回域内分析或辅助能力。",
                "近源与证据按场景推进，不与契约基线抢同一时间的全面铺开；验证的是差异化，不是新平台品牌。",
            ],
        ),
        (
            NAVY,
            SOFT_NAVY,
            "后续扩展",
            "第三方Agent、开放API与行业场景",
            [
                "共用能力和自有Agent验证成立后，再开放第三方Agent工具；外部Agent使用同一契约、身份和证据接口，保持可替换。",
                "将成熟能力映射到NEF、CAMARA Quality on Demand（QoD）、行业API和聚合渠道；能力API验证跨网一致、持续使用和履约成本。",
                "结果型服务只有在业务Owner、基线、授权、归因和回滚清晰时承诺；域级技术效果不能直接外推结果合同。",
                "商业路径只是成熟后的出口选择，不回写为本页之前的模式分析主轴，也不改变中兴必须掌握的电信约束。",
            ],
        ),
    ]
    for i, (accent, fill, kicker, title, bullets) in enumerate(cards):
        x = 0.38 + i * 4.20
        lines = [
            {"text": kicker, "size": 10, "bold": True, "color": accent, "after": 1},
            {"text": title, "size": 13, "bold": True, "color": NAVY, "after": 6},
        ]
        for j, item in enumerate(bullets):
            lines.append({"text": "· " + item, "size": 9.5, "color": SLATE, "after": 4 if j < len(bullets) - 1 else 0, "spacing": 1.03})
        box(
            s,
            x,
            1.16,
            4.08,
            4.18,
            lines,
            fill=fill,
            border=accent,
            radius=0.06,
            mt=0.08,
            mb=0.06,
            ml=0.12,
            mr=0.10,
        )

    gates = [
        ("战略控制门", "核心电信对象、任务适用性、近源约束和动作证据不得外包；通用湖仓、模型和渠道保持开放可替换。", TEAL),
        ("复用与场景门", "跨第二产品/Agent、两种后端、有真实消费者和版本承诺；自有Agent必须实际消费公共能力。", NAVY),
        ("工程停止门", "不成为快环远端单点；无回退、无责任、无消费者或长期项目化定制时停止扩张。阶段门决定加码或收缩。", CORAL),
    ]
    for i, (title, body, accent) in enumerate(gates):
        x = 0.38 + i * 4.20
        box(
            s,
            x,
            5.44,
            4.08,
            1.66,
            [
                {"text": title, "size": 12, "bold": True, "color": accent, "after": 3},
                {"text": body, "size": 10, "color": SLATE, "after": 0, "spacing": 1.03},
            ],
            fill=WHITE,
            border=accent,
            radius=0.06,
            mt=0.08,
        )
    footer(s, 5, "6.3 机会选择")


def slide_p6(prs) -> None:
    s = new_slide(prs)
    header(s, "7.1  ·  STRATEGY TABLE", "围绕共用、嵌入、Agent、工程验证和开放接口形成公司级策略")
    thesis(
        s,
        "本页只保留公司级选择、执行时序和停止条件。详细AP、组织、风险和治理进入第8章，不在决策胶片重复展开。",
        height=0.38,
    )
    add_table(
        s,
        0.38,
        1.10,
        12.55,
        5.12,
        [
            ["我司选择", "具体策略（带时间 / 触发条件）", "不做的风险", "备注"],
            [
                "建立横向公共能力",
                "当前阶段启动：基于AIR Net数据引擎及相邻资产，统一核心对象、语义、数据产品契约、任务适用性、变化事件、Agent工具和动作证据接口；不另建覆盖所有存储与模型的大平台。",
                "产品线重复建设目录、管道、语义和策略；客户看到多个平台孤岛，跨场景交付持续项目化。",
                "进入门：第二场景无需大规模重写、两种后端可替换、有公共Owner和版本基线；内部项目只能作为候选载体。",
            ],
            [
                "嵌入主航道产品",
                "公共基线形成后：公共契约、供数和证据能力进入无线、核心网、OSS/SMO至少两条产品线；产品线保留域扩展、实时状态、资源预算、动作和客户SLO责任。",
                "公共能力脱离网络产品和版本体系成为独立平台；产品线绕开公共基线，重复建设和对象冲突持续。",
                "验收跨产品对象一致、版本兼容、断链运行和证据回传；公共层不得成为毫秒快环依赖。",
            ],
            [
                "做强自有Agent产品",
                "与主航道嵌入并行：Fault Agent加一类非故障Agent复用同一数据产品、任务适用性、工具、策略和证据链；结果反哺质量规则、模型重验和版本发布。",
                "自有Agent继续按项目取数治理；公共能力缺少真实消费者；Agent效果无法跨场景复制和持续改进。",
                "自有Agent是正式产品和首批消费者；验收第二Agent复用、交付周期、动作证据和结果反馈。",
            ],
            [
                "时间盒验证近源与动作证据",
                "真实场景和资源预算具备后：选择高频无线数据与Agent场景，验证近源处理、断链、回退、身份、NDT误差、灰度、人工接管和动作回证。",
                "失去相对云/IT平台的网络差异化；或过度平台化，把远端目录和Agent放入快环关键路径。",
                "停止线：收益不覆盖开销、误差不可标定、动作无法定责回退或必须成为快环单点。",
            ],
            [
                "建立开放接口和阶段门",
                "能力成熟及每个版本门持续：开放第三方Agent、NEF/CAMARA Quality on Demand（QoD）和行业API；标准主张对应代码和测试；按复用、运行、Agent场景和外部使用决定加码或停止。",
                "形成封闭栈，标准与产品脱节；PoC长期化、资源平均分配；API扩张但内部供给和责任不稳定。",
                "能力服务与结果型服务只是成熟后的不同出口；详细AP、Owner和治理进入第8章。阶段门用于加码或停止，而不是并行铺开。",
            ],
        ],
        col_w=[2.26, 5.02, 3.01, 2.26],
        font=8.5,
    )
    box(
        s,
        0.38,
        6.30,
        12.55,
        0.80,
        [
            {
                "text": "中兴应以网络数据与执行触点为基础，建设跨产品共用的电信数据与动作可信能力，并将其嵌入无线、核心网、OSS/SMO和自有智能体应用，形成可信供数、分布式处理、智能决策与动作回证相互贯通的产品体系。",
                "size": 11.5,
                "bold": True,
                "color": NAVY,
                "after": 0,
                "spacing": 1.04,
            }
        ],
        fill=PALE,
        border=TEAL,
        radius=0.06,
        mt=0.08,
        mb=0.06,
        ml=0.14,
        mr=0.12,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s, 6, "7.1 具体策略")


def count_cjk(prs: Presentation) -> list[int]:
    counts = []
    for slide in prs.slides:
        texts = []
        for shp in slide.shapes:
            if shp.has_text_frame:
                texts.append(shp.text_frame.text)
            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        counts.append(len(re.findall(r"[\u4e00-\u9fff]", "".join(texts))))
    return counts


def slide_blob(slide) -> str:
    texts = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            texts.append(shp.text_frame.text)
        if shp.has_table:
            for row in shp.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return "\n".join(texts)


def verify(prs: Presentation) -> None:
    n = len(prs.slides)
    if n != TOTAL:
        raise SystemExit(f"expected {TOTAL} slides, got {n}")

    textboxes = []
    overflow = []
    tables = []
    for i, slide in enumerate(prs.slides, 1):
        table_n = 0
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                textboxes.append((i, getattr(shp, "name", "?")))
            if shp.has_table:
                table_n += 1
            right = (shp.left + shp.width) / EMU
            bottom = (shp.top + shp.height) / EMU
            if right > W_IN + 0.03 or bottom > H_IN + 0.03 or shp.left < Emu(-0.02 * EMU) or shp.top < Emu(-0.02 * EMU):
                overflow.append((i, getattr(shp, "name", "?"), round(right, 3), round(bottom, 3)))
        tables.append(table_n)
    if textboxes:
        raise SystemExit(f"standalone text boxes are not allowed: {textboxes}")
    if overflow:
        raise SystemExit(f"shapes overflow canvas: {overflow}")
    if tables != [0, 0, 0, 1, 0, 1]:
        raise SystemExit(f"expected tables [0,0,0,1,0,1], got {tables}")

    counts = count_cjk(prs)
    print("cjk chars:", counts, "total", sum(counts))
    for i, (got, need) in enumerate(zip(counts, CJK_MIN), 1):
        if got < need:
            raise SystemExit(f"P{i} CJK {got} < {need}")
    if sum(counts) < CJK_TOTAL_MIN:
        raise SystemExit(f"total CJK {sum(counts)} < {CJK_TOTAL_MIN}")

    p4 = prs.slides[3]
    for shp in p4.shapes:
        if not shp.has_table:
            continue
        for r, row in enumerate(shp.table.rows):
            if r == 0:
                continue
            for c, cell in enumerate(row.cells):
                n_cjk = len(re.findall(r"[\u4e00-\u9fff]", cell.text))
                if n_cjk > P4_CELL_CJK_MAX:
                    raise SystemExit(f"P4 cell r{r}c{c} CJK {n_cjk} > {P4_CELL_CJK_MAX}: {cell.text[:40]}")

    required = [
        "AIR Net",
        "Fault Agent",
        "第三方Agent",
        "任务适用性",
        "近源执行",
        "动作证据",
        "横向公共能力",
        "中兴自有Agent",
        "Quality on Demand",
        "阶段门",
    ]
    blob = "\n".join(slide_blob(slide) for slide in prs.slides)
    missing = [item for item in required if item not in blob]
    if missing:
        raise SystemExit(f"missing required terms: {missing}")
    ambiguous_qod = [line for line in blob.splitlines() if "QoD" in line and "Quality on Demand" not in line]
    if ambiguous_qod:
        raise SystemExit(f"QoD must only mean Quality on Demand: {ambiguous_qod}")

    forbidden = [
        "smart data hub",
        "Smart Data Hub",
        "已商用6G Data Fabric",
        "统一控制所有网元",
        "Quality of Data",
        "Build / Partner / Buy",
        "O1",
        "O8",
    ]
    hit = [item for item in forbidden if item in blob]
    if hit:
        raise SystemExit(f"forbidden terms present: {hit}")
    print("verify ok:", {"slides": n, "tables": tables, "textboxes": 0})


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_p1(prs)
    slide_p2(prs)
    slide_p3(prs)
    slide_p4(prs)
    slide_p5(prs)
    slide_p6(prs)
    verify(prs)
    prs.save(OUT)
    print(f"wrote {OUT}")
    return OUT


if __name__ == "__main__":
    build()
