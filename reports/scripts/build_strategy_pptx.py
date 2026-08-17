#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate detailed strategy PPT from the 6G Data Fabric HTML report (~34 pages)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

W, H = Inches(13.333), Inches(7.5)
ASSETS = Path("/workspace/reports/data-fabric-assets")

NAVY = RGBColor(0x07, 0x1F, 0x3D)
INK = RGBColor(0x0A, 0x27, 0x44)
TEAL = RGBColor(0x0F, 0x9F, 0x9A)
CYAN = RGBColor(0x5E, 0xB6, 0xE4)
LIME = RGBColor(0xCA, 0xE8, 0x5E)
CORAL = RGBColor(0xFF, 0x8A, 0x4C)
SAND = RGBColor(0xF4, 0xF7, 0xFB)
MIST = RGBColor(0xDC, 0xE7, 0xF2)
MUTED = RGBColor(0x5A, 0x73, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xE8, 0xF1, 0xF6)
WARN = RGBColor(0xFF, 0xF0, 0xE7)
ROW_ALT = RGBColor(0xF0, 0xF5, 0xF8)


def rgb_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_run(run, text, size=12, bold=False, color=INK, font="Microsoft YaHei"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    east = rPr.find(qn("a:eastAsian"))
    if east is None:
        east = rPr.makeelement(qn("a:eastAsian"), {})
        rPr.append(east)
    east.set("typeface", font)


def write_box(shape, lines, *, default_size=11, default_color=INK):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", "t")
    except Exception:
        pass
    p0 = tf.paragraphs[0]
    p0.clear()
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color, after = item, default_size, False, default_color, 3
        elif len(item) == 2:
            text, size = item
            bold, color, after = False, default_color, 3
        elif len(item) == 3:
            text, size, bold = item
            color, after = default_color, 3
        elif len(item) == 4:
            text, size, bold, color = item
            after = 3
        else:
            text, size, bold, color, after = item
        p = p0 if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(after)
        run = p.add_run()
        set_run(run, text, size=size, bold=bold, color=color)


def rect(slide, left, top, width, height, fill, *, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    rgb_fill(shape, fill)
    no_line(shape)
    if radius:
        try:
            shape.adjustments[0] = 0.1
        except Exception:
            pass
    return shape


def footer(slide, page, total, chapter):
    rect(slide, Inches(0.4), Inches(7.15), Inches(12.5), Emu(12700), MIST)
    left = slide.shapes.add_textbox(Inches(0.4), Inches(7.2), Inches(9.6), Inches(0.25))
    write_box(left, [(f"6G 数据架构 × 数据编织 · 完整论证版 · {chapter}", 8, False, MUTED)])
    right = slide.shapes.add_textbox(Inches(10.2), Inches(7.2), Inches(2.7), Inches(0.25))
    write_box(right, [(f"{page:02d} / {total:02d}", 8, True, NAVY)])
    for p in right.text_frame.paragraphs:
        p.alignment = PP_ALIGN.RIGHT


def topbar(slide, kicker, title):
    rect(slide, Inches(0), Inches(0), W, Inches(0.07), TEAL)
    kb = slide.shapes.add_textbox(Inches(0.4), Inches(0.14), Inches(12.5), Inches(0.24))
    write_box(kb, [(kicker, 9, True, TEAL)])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.34), Inches(12.5), Inches(0.38))
    write_box(tb, [(title, 18, True, NAVY)])


def thesis_box(slide, text, top=Inches(0.78), height=Inches(0.68)):
    box = rect(slide, Inches(0.4), top, Inches(12.5), height, NAVY, radius=True)
    write_box(box, [("本页论点", 8, True, LIME, 1), (text, 11, False, WHITE, 0)])
    tf = box.text_frame
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05)
    return box


def card(slide, left, top, width, height, lines, *, fill=WHITE, accent=None, pad=True):
    shape = rect(slide, left, top, width, height, fill, radius=True)
    if accent is not None:
        rect(slide, left, top, Emu(50800), height, accent)
    write_box(shape, lines)
    if pad:
        tf = shape.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.05)
    return shape


def add_table(slide, left, top, width, height, headers, rows, *, col_w=None, font_size=8):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    if col_w:
        for i, w in enumerate(col_w):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        set_run(run, h, size=font_size, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            set_run(run, str(val), size=max(font_size - 0.5, 7), bold=False, color=INK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else ROW_ALT
            cell.vertical_anchor = MSO_ANCHOR.TOP
            cell.text_frame.word_wrap = True
    return table_shape


def bg(slide, dark=False):
    rect(slide, 0, 0, W, H, NAVY if dark else SAND)
    if not dark:
        rect(slide, 0, 0, Inches(0.1), H, TEAL)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def maybe_pic(slide, path, left, top, width=None, height=None):
    if path.exists():
        kwargs = {}
        if width is not None:
            kwargs["width"] = width
        if height is not None:
            kwargs["height"] = height
        return slide.shapes.add_picture(str(path), left, top, **kwargs)
    return None


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = 34
    meta = []
    page = 0

    def finish(s, chapter, name):
        nonlocal page
        page += 1
        footer(s, page, total, chapter)
        meta.append(name)

    # ===== 01 Cover =====
    s = new_slide(prs)
    bg(s, dark=True)
    rect(s, 0, 0, W, Inches(0.16), TEAL)
    rect(s, 0, Inches(7.34), W, Inches(0.16), CYAN)
    for i, c in enumerate([TEAL, CYAN, LIME, CORAL]):
        rect(s, Inches(0.5 + i * 0.32), Inches(1.35), Inches(0.24), Inches(0.07), c)
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(12), Inches(0.35))
    write_box(box, [("产业战略研究 · 中兴通讯视角 · 非公司正式立场 · 研究截点 2026-08", 12, False, CYAN)])
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(12.2), Inches(1.5))
    write_box(box, [
        ("6G 数据架构 × 数据编织", 32, True, WHITE, 6),
        ("行业趋势与中兴战略定位 · 完整论证胶片", 24, True, LIME, 0),
    ])
    box = s.shapes.add_textbox(Inches(0.5), Inches(3.85), Inches(12), Inches(1.6))
    write_box(box, [
        ("核心问题：当6G同时成为数据生产者、AI载体、感知系统与自治执行体时，如何建立跨RAN/Core/管理域/边云的可信数据能力，并判断哪些能力将成为真正控制点。", 13, False, MIST, 8),
        ("章节对齐 HTML：01需求 → 02行业竞争 → 03趋势T1–T8 → 04三/五年 → 05中兴定位 → 06机会 → 07策略 → 08关键AP", 12, False, CYAN, 6),
        ("本胶片为完整论证加深版（非预览摘要）：保留反证条件、证据分层、准入止损、阶段门、厂商深论证与验收指标。", 12, False, LIME, 0),
    ])
    box = s.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12), Inches(0.5))
    write_box(box, [("锚点：ITU · 3GPP SA2/SA5 · O-RAN · ETSI ZSM 029　｜　公司：AIR Net · Fault Agent · CUDR · Co-Sight · Open Gateway QoD", 11, False, MUTED)])
    page += 1
    meta.append("封面")
    # cover has no standard footer look - add page mark
    mark = s.shapes.add_textbox(Inches(11.5), Inches(6.9), Inches(1.4), Inches(0.3))
    write_box(mark, [("01 / 34", 10, True, LIME)])

    # ===== 02 Logic =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "00  READING LOGIC", "阅读逻辑：八章是一条因果链，不是八个平行清单")
    thesis_box(s, "先回答「为何需要数据管理层、行业在抢什么」，再推出「趋势不变量与时间窗」，最后才落到中兴「定位—机会—策略—可验收AP」。跳过前四章，后四章会变成口号。")
    steps = [
        ("01", "需求和目标", "定义问题：可信供数与控制点；输出七个可交付判断与 R1–R8"),
        ("02", "行业与竞争", "范式转变、Fabric适配、成熟度、标准职责与 Ericsson/Huawei/Nokia 抢位"),
        ("03", "产业趋势", "T1–T8 因果链：产品化→事件→两级部署→有界智能→证据→QoD→商业→能力集合"),
        ("04", "3年/5年", "看透不变量；看清控制点复利；十条决策压缩判断"),
        ("05", "中兴定位", "证据 E0–E4 分层；C1–C4 控制点与开放边界；做/不做"),
        ("06", "机会选择", "六维加权；主攻 O2/O3/O1/O7/O8；配合项触发与资源配比"),
        ("07", "策略建议", "策略1–11；Build/Partner/Buy；节奏、G0–G4、度量与风险台账"),
        ("08", "关键 AP", "AP0–AP10；90天六输出；治理层级；D1–D6 公司裁决"),
    ]
    for i, (n, t, d) in enumerate(steps):
        col, row = i % 4, i // 4
        card(s, Inches(0.4 + col * 3.2), Inches(1.65 + row * 2.5), Inches(3.05), Inches(2.3), [
            (n, 16, True, TEAL, 2),
            (t, 14, True, NAVY, 6),
            (d, 11, False, MUTED, 8),
            ("↓ 输出喂给下一章" if i < 7 else "→ 形成可验收行动包", 10, True, CORAL if i < 7 else TEAL, 0),
        ], accent=TEAL if row == 0 else CORAL)
    finish(s, "阅读逻辑", "阅读逻辑")

    # ===== 03 Need =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "01  REQUIREMENT & OBJECTIVE", "需求和目标：不是再建平台，而是建立可信数据能力与控制点判断")
    thesis_box(s, "本报告要解决的不是「是否再建一个数据平台」，而是：跨域可信数据能力如何建立，以及哪些能力会在标准与产业竞争中成为真正控制点。")
    needs = [
        ("需求1 · 识别架构范式变化", "回答6G为何不能只沿用「会话连接 + 中央分析」；数据管理为何从后台支撑上升为 AI 原生、通感一体、数字孪生与自治闭环的共同底座。"),
        ("需求2 · 建立统一能力地图", "用数据编织把目录、语义、质量、集成、编排、数据产品、安全和可观察组织成同一逻辑控制面；明确快环/慢环、控制面/执行面边界。"),
        ("需求3 · 看清标准与竞争窗口", "区分3GPP/O-RAN既定积木、Rel-20研究、Rel-21规范窗口、厂商候选架构与商业产品；避免把研究主张误判为标准结论，也避免错过语义与QoD话语权。"),
    ]
    for i, (t, b) in enumerate(needs):
        card(s, Inches(0.4 + i * 4.2), Inches(1.65), Inches(4.05), Inches(2.2), [
            (t, 13, True, NAVY, 6),
            (b, 11, False, MUTED, 0),
        ], accent=TEAL)
    note = s.shapes.add_textbox(Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.28))
    write_box(note, [("报告目标：形成八个可交付判断（七项产出 + 边界声明）", 11, True, TEAL)])
    dels = [
        ("1张范式图", "解释连接架构如何演进为数据—智能—行动架构"),
        ("1套能力栈", "统一数据编织与6G数据生命周期，不重复造体系"),
        ("1个竞争矩阵", "比较标准域与 Ericsson/Huawei/Nokia 等主张"),
        ("2个时间窗", "3年看标准收敛与工程底座；5年看控制点与生态"),
        ("1个中兴定位", "在「编织×6G」坐标系明确强弱与主攻位"),
        ("1组机会选择", "按可达性×契合度筛主攻、中期押注与路标"),
        ("1张策略表", "先做什么、后做什么、不做风险与组织抓手"),
        ("边界声明", "公开证据与内部路标建议分层；非正式产品承诺"),
    ]
    for i, (t, d) in enumerate(dels):
        col, row = i % 4, i // 4
        card(s, Inches(0.4 + col * 3.2), Inches(4.35 + row * 1.2), Inches(3.05), Inches(1.1), [
            (t, 12, True, NAVY, 2),
            (d, 10, False, MUTED, 0),
        ], fill=SOFT)
    finish(s, "01 需求和目标", "01需求与目标")

    # ===== 04 Evidence chain =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "01  EVIDENCE CHAIN", "研究对象是一条「事实—数据—智能—行动」证据链")
    thesis_box(s, "目标不是搬运更多原始数据，而是让每次数据使用与网络动作都可解释、可约束、可验证、可回退。核心需求一句话：把数据从「网络运行副产物」升级为「可发现、可理解、可组合、可约束、可证明的生产要素」。")
    chain = [
        ("① 网络与环境事实", "RAN · Core · 云\n终端 · 感知 · 业务", WHITE, NAVY),
        ("② 近源处理", "过滤 · 聚合 · 特征\n匿名 · 时间对齐", WHITE, NAVY),
        ("③ 数据编织逻辑控制面", "目录·语义·QoD·血缘\n契约·策略·版本·证据", NAVY, LIME),
        ("④ 智能与孪生", "训练 · 推理 · NDT\n意图 · 策略 · Agent", WHITE, NAVY),
        ("⑤ 行动与反馈", "控制 · 开放 · 变现\n效果 · 回滚 · 审计", WHITE, NAVY),
    ]
    for i, (t, d, fill, tc) in enumerate(chain):
        card(s, Inches(0.35 + i * 2.58), Inches(1.65), Inches(2.45), Inches(1.7), [
            (t, 12, True, tc if fill != NAVY else LIME, 4),
            (d, 11, False, MUTED if fill != NAVY else MIST, 0),
        ], fill=fill, accent=CORAL if i != 2 else LIME)
        if i < 4:
            a = s.shapes.add_textbox(Inches(2.7 + i * 2.58), Inches(2.25), Inches(0.28), Inches(0.3))
            write_box(a, [("→", 14, True, CORAL)])
    card(s, Inches(0.4), Inches(3.55), Inches(12.5), Inches(0.7), [
        ("控制面统一的是理解与约束，不是所有存储与时延。近源执行与快环保持分布式；否则 AI 原生只会把数据孤岛与错误动作自动化。", 12, False, NAVY, 0),
    ], fill=SOFT)
    drivers = [
        ("AI原生", "训练/验证/部署/监控/回滚依赖版本一致；AI错误升级为网络动作风险"),
        ("通感一体", "ISAC产生高敏感数据；近源最小化、用途绑定、置信度成产品属性"),
        ("边云一体", "算力与数据位置共同决定架构；快环不能依赖远端目录"),
        ("数字孪生", "NDT需关联拓扑/配置/遥测/模型；无事件时间与血缘则不可作证据"),
        ("跨域自治", "多Agent目标冲突；需意图分解、风险预算、租约、熔断与接管"),
        ("生态开放", "Open Gateway推动API；变现取决于跨网一致语义、计量与工作流"),
    ]
    for i, (t, b) in enumerate(drivers):
        col, row = i % 3, i // 3
        card(s, Inches(0.4 + col * 4.2), Inches(4.45 + row * 1.2), Inches(4.05), Inches(1.1), [
            (t, 12, True, TEAL, 2),
            (b, 10, False, MUTED, 0),
        ])
    finish(s, "01 需求和目标", "01证据链与驱动")

    # ===== 05 R1-R8 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "01  R1–R8 REQUIREMENTS", "六类场景反推：数据管理层八项刚性需求（含验收口径）")
    thesis_box(s, "数据管理层目标由此明确：不是建更大的数据池，而是建立跨域、跨时标的「数据供应与行动证据控制系统」——允许执行分布在RAN/Core/边云，同时统一语义、质量、策略与责任。")
    add_table(
        s, Inches(0.35), Inches(1.6), Inches(12.6), Inches(5.2),
        ["需求", "刚性内容", "验收口径"],
        [
            ["R1 统一发现与语义寻址", "消费者要知道对象、时间、空间、配置版本、业务含义与权威源；RAN小区/E2/O1 DN/Core会话/云资源可关联", "跨域对象映射率、语义版本覆盖率"],
            ["R2 按需采集与近源减量", "ISAC/UE测量/高频遥测不能默认全量上送；按SLO编排采样、过滤、聚合、特征、匿名；决定数据移动或模型下沉", "重复采集下降率、跨域字节/有效任务"],
            ["R3 跨时标多模式交付", "亚10ms内环、Near-RT、Non-RT/OAM与月度训练不能共享同一在线依赖；支持本地状态/流/CDC/批/联邦/缓存/异步制品", "分时标P99、新鲜度与截止期达标率"],
            ["R4 通信级数据质量QoD", "除完整性与新鲜度外：采样代表性、对象/频段/波束覆盖、时间同步误差、配置一致性、测量置信度、训练—服务偏差", "质量信封覆盖、suspect/incomplete比例"],
            ["R5 数据—模型—策略—动作血缘", "每个结论与动作可追溯到输入、特征、模型版本、策略、作用域与结果；反馈反向修正采集/质量/模型", "端到端证据链覆盖率、回放成功率"],
            ["R6 用途、安全与主权强制", "身份、同意、目的限制、驻留、保留、删除贯穿查询/API/导出/特征库/向量索引/模型训练，而非只在目录挂标签", "跨路径策略一致率、删除传播时长"],
            ["R7 有界智能与冲突治理", "多Agent/模型/闭环可能争夺资源；提供身份、委托、风险预算、动作租约、冲突仲裁、仿真、熔断与人工接管证据", "越权拦截、无效动作率、回滚耗时"],
            ["R8 数据产品与跨组织运营", "输出需Owner、契约、QoD、用途、SLO、成本与退出；衔接CAMARA/Operate API、数据空间合同与计量", "复用消费者数、契约合规与单位收入"],
        ],
        col_w=[Inches(2.6), Inches(7.0), Inches(3.0)],
        font_size=9,
    )
    finish(s, "01 需求和目标", "01 R1-R8")

    # ===== 06 Paradigm =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  INDUSTRY WINDOW · PARADIGM", "行业窗口：愿景已定、研究展开、规范未冻、厂商抢位")
    thesis_box(s, "ITU已给IMT-2030框架；Rel-20推进6G研究，Rel-21将成为首个规范性6G Release。6G变化不是「多一个数据湖」，而是几乎每个域同时生产/消费数据、运行模型并执行动作——从连接中心转向数据理解、策略与证据层。")
    card(s, Inches(0.4), Inches(1.65), Inches(6.1), Inches(5.0), [
        ("5G / 5G-A 主导范式", 15, True, NAVY, 8),
        ("中心叙事", 12, True, TEAL, 3),
        ("连接、会话、网络功能与暴露；分析集中在 NWDAF / MDA / 湖仓。", 11, False, MUTED, 6),
        ("数据角色", 12, True, TEAL, 3),
        ("RAN/Core 主要上送测量与流量；管理域做报表与优化建议。", 11, False, MUTED, 6),
        ("局限", 12, True, CORAL, 3),
        ("中央分析难覆盖多时标、多主体与动作风险；权威对象与语义分散；数据多为运行副产物，而非可运营产品。", 11, False, INK, 8),
        ("范式变化不是取消网络功能，而是给分布式网络增加统一的数据理解、策略与证据层。", 11, True, NAVY, 0),
    ], fill=SOFT)
    card(s, Inches(6.75), Inches(1.65), Inches(6.1), Inches(5.0), [
        ("6G 数据—智能原生范式", 15, True, LIME, 8),
        ("域角色升级为 prosumer", 12, True, CYAN, 3),
        ("RAN：无线事实/ISAC/本地AI　｜　Core：意图/会话/暴露\n边云：训练/推理/孪生　｜　SMO/RIC：跨时标策略与闭环", 11, False, MIST, 8),
        ("可信数据管理层", 12, True, LIME, 3),
        ("语义 · QoD · 契约 · 策略 · 证据\n统一理解与约束，不统一所有存储与时延", 11, False, WHITE, 8),
        ("Nokia 亦把 data prosumer、目录、抽象、产品与生命周期列为 AI-native 6G 使能能力。右侧为本文综合，不是 3GPP 既定统一架构。", 10, False, MUTED, 0),
    ], fill=NAVY)
    finish(s, "02 行业环境", "02范式转变")

    # ===== 07 Fabric fit =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  DATA FABRIC FIT", "Data Fabric 能否满足 R1–R8？总体适配，必须通信级增强")
    thesis_box(s, "数据编织不是6G新增网元，而是跨异构分布式环境的数据管理与集成设计：以主动元数据、语义、质量、策略和编排构成逻辑控制面。它可做6G数据管理层设计骨架，但不能原样进入通信网，更不能成为毫秒快环的远端在线依赖。")
    card(s, Inches(0.4), Inches(1.6), Inches(6.15), Inches(1.35), [
        ("结论：架构机制总体适配", 13, True, TEAL, 3),
        ("数据分布但元数据统一、控制与执行分离、多模式集成、主动元数据、数据产品与贯穿治理 → 对应 R1/R2/R3、R5、R8、R6；尤其适合 Non-RT/管理域、跨域供给及快环的环外支撑。", 10, False, INK, 0),
    ], fill=SOFT, accent=TEAL)
    card(s, Inches(6.75), Inches(1.6), Inches(6.15), Inches(1.35), [
        ("但必须通信级增强", 13, True, CORAL, 3),
        ("通用Fabric缺RAN对象与时空语义、无线QoD、亚秒时标隔离、动作安全、NDT证据门槛和多Agent冲突治理；不能替代RAN协议、SA2/SA5功能或O-RAN接口。", 10, False, INK, 0),
    ], fill=WARN, accent=CORAL)
    add_table(
        s, Inches(0.35), Inches(3.15), Inches(12.6), Inches(3.65),
        ["6G需求", "Fabric原生机制", "适配", "必须补齐的通信扩展"],
        [
            ["R1 发现与语义", "目录、KG、血缘、主动元数据", "高", "O1/E2/Core/云对象统一、事件时间、配置版本"],
            ["R2 近源减量", "联邦查询、策略编排、计算下推", "高", "DU/CU/MEC可移植算子、实时资源预算、断链策略"],
            ["R3 跨时标交付", "批/流/CDC/API/缓存多模式", "中高", "按闭环时标隔离；本地确定性路径不依赖远端控制面"],
            ["R4 通信QoD", "质量规则、Profiling、SLO", "部分", "时空代表性、无线误差、同步/配置/采样质量信封"],
            ["R5 证据链", "数据/模型血缘、元数据事件", "高", "模型—策略—动作—网络效果闭环血缘与回放"],
            ["R6 安全主权", "身份、策略即代码、用途审计", "中高", "跨路径强制、ISAC隐私、撤销/删除到缓存与模型"],
            ["R7 有界智能", "策略、工作流、推荐与自动化", "不足", "Agent身份、动作租约、冲突仲裁、NDT预演、熔断回滚"],
            ["R8 产品运营", "数据产品、契约、Marketplace", "中等", "网络QoD计量、CAMARA/TMF运营接口、结果定价与责任"],
        ],
        col_w=[Inches(2.2), Inches(3.3), Inches(1.0), Inches(6.1)],
        font_size=8,
    )
    finish(s, "02 行业环境", "02 Fabric适配映射")

    # note: DDAA figure appears on challenges/context via Nokia slide imagery; keep table full-bleed here

    # ===== 08 Stack maturity =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  CAPABILITY STACK & MATURITY", "能力栈与成熟度：基础组件成熟，通信级控制闭环未成熟")
    thesis_box(s, "不能因为目录、湖仓或流平台已商用，就宣称「6G Data Fabric成熟」。L1–L3大多可买可建；真正决定Fabric是否成立的L4主动元数据、L5产品化、R7有界智能和跨路径横切治理仍处试点或研究阶段——总体 G1—G2、局部 G3。")
    layers = [
        ("L6", "业务、生态与网络行动", "网络API·行业数据空间·RAG/Agent·运营闭环·自动控制", "传统消费高；受治Agent低"),
        ("L5", "数据产品与可信自助", "产品目录·契约·语义指标·QoD/SLO·计量与价值", "试点增多；跨域语义偏弱"),
        ("L4", "主动元数据与智能控制面", "目录·血缘·RAN/Core语义·策略·推荐·影响分析·自动化", "被动目录成熟；主动闭环中等"),
        ("L3", "集成编排与跨时标平面", "批/CDC/流·联邦·O1/E2/A1/R1适配·缓存·工作流", "中高；多接口统一语义不足"),
        ("L2", "近源处理、运行时与存储", "DU/CU/MEC·SMO/Non-RT·湖仓·特征/向量/图·NDT", "组件高；跨域编排中等"),
        ("L1", "分布式权威数据源", "UE·RU/DU/CU·RIC·Core NF·O-Cloud·OSS/BSS·ISAC/NTN", "数据存在；权威语义分散"),
    ]
    for i, (lid, name, body, st) in enumerate(layers):
        card(s, Inches(0.4), Inches(1.55 + i * 0.72), Inches(8.3), Inches(0.68), [
            (f"{lid}  {name}　｜　{body}", 10, True, NAVY, 1),
            (st, 9, False, MUTED, 0),
        ], fill=WHITE if i % 2 == 0 else SOFT, accent=TEAL)
    card(s, Inches(8.9), Inches(1.55), Inches(4.0), Inches(5.05), [
        ("组织落地闸门", 14, True, NAVY, 8),
        ("G0 边界与语义", 12, True, TEAL, 2),
        ("场景、Owner、对象、时间、配置、权威源与用途明确；否则自动化只传播歧义。", 10, False, MUTED, 6),
        ("G1 只读可观察", 12, True, TEAL, 2),
        ("目录/Schema/QoD/血缘可查询，先证明能回答「哪份数据对应哪个网络状态」。", 10, False, MUTED, 6),
        ("G2 产品与契约", 12, True, TEAL, 2),
        ("高频输出产品化；契约进CI；消费者、成本、复用与退出可度量。", 10, False, MUTED, 6),
        ("G3 影子闭环", 12, True, TEAL, 2),
        ("只给建议，对照人工与真实结果；高风险先在孪生/沙箱演练。", 10, False, MUTED, 6),
        ("G4–G5 有界自治与跨域", 12, True, TEAL, 2),
        ("白名单低风险自动执行；内部证据稳定后再向Core/行业/伙伴供给。", 10, False, MUTED, 8),
        ("结合原则：Fabric不应进入所有实时关键路径。", 10, True, CORAL, 0),
    ])
    finish(s, "02 行业环境", "02能力栈成熟度")

    # ===== 09 Challenges =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  KEY CHALLENGES", "适配性论证后的剩余挑战：从通用架构走向通信级控制系统")
    thesis_box(s, "挑战不是对 Data Fabric 适配性的否定，而是通信级控制系统必须补齐的工程与治理门槛。能力演进不是「平台越大越先进」，而是控制面逐步主动化：被动采集→可观察→可产品化→影子闭环→有界自治。")
    add_table(
        s, Inches(0.3), Inches(1.6), Inches(12.7), Inches(5.15),
        ["关键技术挑战", "为何在6G更难", "关键突破点", "错误路线"],
        [
            ["跨域语义与身份", "O1 DN、E2、Core会话、切片、云与业务实体标识体系不同", "版本化对象模型、Correlation ID、语义映射与权威源", "把接口互通等同语义统一"],
            ["RAN数据QoD", "采样代表性、时空粒度、配置版本直接影响模型与闭环", "质量信封：覆盖、时效、同步误差、suspect、配置一致", "只统计记录数或PB规模"],
            ["快环/慢环协同", "从亚10ms到月度训练跨越多个数量级", "本地快环、异步元数据、分层缓存、模型/策略制品同步", "让远端目录成为在线依赖"],
            ["模型—数据—动作证据链", "模型输出可直接改网；训练服务偏差扩大故障半径", "版本绑定、NDT预演与效果回流", "只管模型仓库，不管数据与动作"],
            ["Agent治理与冲突", "多域Agent目标、权限与时限不同，局部最优伤害全局", "IPOE、租约、冲突仲裁、人工接管", "把自然语言界面当作自治"],
            ["数据主权与隐私", "ISAC/位置/终端/行业数据敏感；跨组织用途难持续强制", "近源最小化、用途策略、可信执行、删除传播", "以接口加密替代全生命周期治理"],
            ["多厂商互操作", "接口可标准化，目录/QoD/策略/模型/Agent语义仍可能锁定", "开放元数据事件、数据契约、可移植运行时、一致性测试", "把单厂商端到端等同开放"],
            ["单位经济性", "高频遥测、跨云传输、训练与孪生能耗可能吞噬收益", "近源减量、按消费者SLO采集、成本/能耗纳入编排", "默认采集越多越好"],
        ],
        col_w=[Inches(2.3), Inches(3.6), Inches(3.7), Inches(3.1)],
        font_size=8,
    )
    finish(s, "02 行业环境", "02关键挑战")

    # ===== 10 Standards =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  STANDARDS LANDSCAPE", "标准竞争：谁定义来源、语义、质量、发现、处理、暴露、安全与行动边界")
    thesis_box(s, "当前职责分散且互补：RAN定义无线事实与接口；SA5管理数据与AI/NDT生命周期；SA2研究系统级6G数据框架；O-RAN开放RAN管理与智能控制；TM Forum/GSMA/CAMARA连接自治运营与商业暴露。竞争本质不在新网元名称，而在跨域语义、QoD、可信执行和生态分发能否形成复利。")
    add_table(
        s, Inches(0.3), Inches(1.6), Inches(12.7), Inches(4.0),
        ["标准域", "当前可确认进展", "对6G Data Fabric贡献", "不能扩大化为"],
        [
            ["ITU-R IMT-2030", "M.2160已生效；六类场景及泛在智能/通感/孪生趋势", "给出6G目标与能力边界", "未定义「6G数据编织」网元或接口"],
            ["3GPP Rel-20 / Rel-21", "Rel-20以研究为主；Rel-21首个规范性6G Release，Stage2计划2028-06冻结", "2026–2028是架构收敛与标准输入窗口", "Rel-20研究文本不是最终规范"],
            ["3GPP SA2", "TR 23.801-01仍为Draft；KI#21=6G data framework", "研究系统级发现、采集、处理、存储、暴露等候选能力", "不预判统一框架/网元，不主导RAN内部协议"],
            ["3GPP RAN / SA5", "SA5已有PM/CM/FM/Trace、MDA、AI/ML、NDT管理；Rel-20 Draft TR研究DMFW", "无线事实、管理数据、模型/孪生生命周期主锚点", "Draft≠规范性TS；SA5不是毫秒级RAN总线"],
            ["O-RAN", "WG2 Non-RT/A1/R1；WG3 Near-RT/E2；WG10 O1/OAM", "按时标提供RAN数据、策略、模型与动作通道", "不是跨O1/O2/A1/E2/R1统一语义总线"],
            ["TMF / GSMA / CAMARA", "AN等级、AI/数据治理、Open Gateway与Operate API认证推进", "把受治理网络能力包装为可运营、可计量、可聚合产品", "不定义无线快环或内部数据处理实现"],
        ],
        col_w=[Inches(2.2), Inches(4.0), Inches(3.5), Inches(3.0)],
        font_size=8,
    )
    card(s, Inches(0.4), Inches(5.8), Inches(12.5), Inches(0.95), [
        ("证据分层纪律", 12, True, CORAL, 2),
        ("「已有运营商部署」≠「已形成跨域统一Fabric」；「厂商已发布产品」≠「已广泛商用」；「标准贡献进Draft TR」≠「形成规范性TS」。产品、参考架构、标准输入、原型与愿景必须分别判断。", 10, False, INK, 0),
    ], fill=WARN)
    finish(s, "02 行业环境", "02标准格局")

    # ===== 11 Vendor Ericsson =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  VENDOR · ERICSSON", "Ericsson：以联邦数据管理连接网络自动化与 API 生态")
    thesis_box(s, "路线特征是演进式、平台组合式：从「采集一次、授权复用」扩展到 AI-ready data mesh 与 6G common data plane。不会押注单一新网元，更可能用平台组合占据数据—自动化—API链路；但尚无中立证据证明已组合成跨域统一 Fabric。")
    card(s, Inches(0.4), Inches(1.6), Inches(7.5), Inches(3.7), [
        ("架构主线：数据摄取 → 联邦网格 → 智能平台", 13, True, NAVY, 6),
        ("· 2021：源数据进EDCA；DRG跨边界中继；DDC与Global Data Catalog协同形成联邦数据湖，解决重复采集。", 10, False, MUTED, 4),
        ("· 2026 AI-ready data mesh：数据源→摄取→refinement & governance→consumer support；强调计算靠近数据、按产品与授权复用。", 10, False, MUTED, 4),
        ("· 6G愿景 common data plane：采集/传输/潜在存储与治理；模型编排器选择「数据到智能」或「智能到数据」——是厂商愿景，不是3GPP标准。", 10, False, MUTED, 6),
        ("成熟度分层（关键）", 12, True, TEAL, 3),
        ("Swisscom等合同证明EIAP进入运营网络；Mediation/Telco DataOps有成熟采集加工（如Grameenphone日超60亿条记录）。这些数字不证明 mesh/EFDL/common data plane 已同规模统一产品化。", 10, False, INK, 0),
    ])
    card(s, Inches(8.1), Inches(1.6), Inches(4.8), Inches(3.7), [
        ("锚点与判断", 13, True, NAVY, 6),
        ("商用：NWDAF、EIAP/rApps、CCES、2025 Telco DataOps", 10, False, MUTED, 4),
        ("架构：AI-ready mesh、EFDL、摄取架构", 10, False, MUTED, 4),
        ("愿景：Common Data Plane / intelligent fabric", 10, False, MUTED, 4),
        ("生态：Aduna聚合CAMARA（独立合资）", 10, False, MUTED, 8),
        ("优势控制点", 11, True, TEAL, 2),
        ("运营商关系 + SMO/RIC应用 + Core/API + 集成平台", 10, False, MUTED, 6),
        ("关键短板", 11, True, CORAL, 2),
        ("名称多；缺跨域统一Fabric中立证据", 10, False, MUTED, 0),
    ], fill=SOFT)
    rect(s, Inches(0.4), Inches(5.45), Inches(12.5), Inches(1.3), WHITE, radius=True)
    maybe_pic(s, ASSETS / "ericsson-data-ingestion-architecture.png", Inches(0.55), Inches(5.55), height=Inches(1.1))
    cap = s.shapes.add_textbox(Inches(7.0), Inches(5.7), Inches(5.7), Inches(0.9))
    write_box(cap, [
        ("图｜Ericsson Data Ingestion Architecture", 11, True, NAVY, 2),
        ("EDCA汇集源数据，经DRG跨域转发；DDC与Global Data Catalog支持外部应用与联邦湖。研究引用图，版权归Ericsson。", 9, False, MUTED, 0),
    ])
    finish(s, "02 行业环境", "02 Ericsson")

    # ===== 12 Vendor Huawei =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  VENDOR · HUAWEI", "Huawei：以专用数据面研究连接全栈智能与 Agent 闭环")
    thesis_box(s, "路线特征是研究架构激进、产品桥梁广：DO/DA/DCP提出独立数据协作逻辑，AUTINOps提供当前部署证据。不宜复制命名，应识别其背后的可组合算子、数据QoD与异步分发需求；研究架构≠现网产品≠标准共识。")
    card(s, Inches(0.4), Inches(1.6), Inches(7.5), Inches(3.55), [
        ("架构主线：服务需求 → 数据拓扑 → 近源执行 → 异步分发", 13, True, NAVY, 5),
        ("· DO将服务需求拆成采集/过滤/聚合/转换/分析，依据DA/DPF能力形成任意拓扑；DA既可作源也可承载模型处理；DCP以消息代理/队列/PubSub分发非连接类数据。", 10, False, MUTED, 4),
        ("· 直接回应AI/ISAC多输入多处理多输出；含与RocketMQ的原型比较。技术价值与标准状态必须分开：不证明必须引入独立Data Plane或名称会进3GPP。", 10, False, MUTED, 4),
        ("· DCP实验：单机转发约2ms、相对RocketMQ平均降>65%；但移除持久化而对照开启持久化，非同功能公平比较；顺序/容灾/多租户等未达运营商级证据→成熟度=研究原型。", 10, False, MUTED, 5),
        ("商用桥梁 AUTINOps / IOH", 12, True, TEAL, 2),
        ("跨无线/微波/IP/传输/能源；故障一次闭环率80%、MTTR-15%、年Agent调用>200万——证明部署，不证明6G专用Data Plane商用。", 10, False, INK, 0),
    ])
    card(s, Inches(8.1), Inches(1.6), Inches(4.8), Inches(3.55), [
        ("分层锚点", 13, True, NAVY, 5),
        ("部署：AUTINOps/IOH跨域案例", 10, False, MUTED, 3),
        ("产品桥：AI Core/AgenticCore、Digital Map、Open Gateway", 10, False, MUTED, 3),
        ("研究：DO/DA/DCP、NET4AI", 10, False, MUTED, 3),
        ("标准输入：DMFW、Data QoS、Agent Fabric、NDT sandbox", 10, False, MUTED, 6),
        ("优势", 11, True, TEAL, 2),
        ("全栈覆盖：RAN—Core—IP—Cloud—AI—终端", 10, False, MUTED, 4),
        ("短板", 11, True, CORAL, 2),
        ("联合披露 vs 独立审计；研究/产品/标准仍有距离", 10, False, MUTED, 0),
    ], fill=SOFT)
    rect(s, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.4), WHITE, radius=True)
    maybe_pic(s, ASSETS / "huawei-6g-data-plane-architecture.jpg", Inches(0.55), Inches(5.45), height=Inches(1.2))
    cap = s.shapes.add_textbox(Inches(6.3), Inches(5.55), Inches(6.3), Inches(1.05))
    write_box(cap, [
        ("图｜Huawei提出的6G Data Plane（厂商研究架构，非3GPP既定方案）", 11, True, NAVY, 2),
        ("DO编排数据拓扑；DA/DPF近源执行；DCP/Data Spine以Pub/Sub解耦生产者与消费者。", 9, False, MUTED, 0),
    ])
    finish(s, "02 行业环境", "02 Huawei")

    # ===== 13 Vendor Nokia + adjacent + summary =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "02  VENDOR · NOKIA & ADJACENT", "Nokia 产品化最清晰；云/湖仓/API 控制外层入口；尚无绝对领先")
    thesis_box(s, "Nokia证明 Data Fabric 已从白皮书进入可售产品，不能再仅作旁注；但命名客户与广泛跨域生产闭环证据仍有限。云与数据平台争夺Catalog/模型/工具，难独代RAN语义与动作安全。现阶段最稳健断言：产业确认「需要哪些能力」，仍在竞争「以什么架构组合、由谁控制、怎样证明价值」。")
    card(s, Inches(0.4), Inches(1.55), Inches(4.4), Inches(3.15), [
        ("Nokia：域自治 → 公共Fabric → 数据产品 → 解释性自动化", 11, True, NAVY, 3),
        ("AN Fabric + Data Suite：统一数据管理、可观察、telco AI、安全与自动化；TM Forum语义与多厂商适配。", 9, False, MUTED, 3),
        ("域证据：KDDI/stc/Indosat证明域应用，不能替代跨域统一生产闭环证据。", 9, False, MUTED, 3),
        ("规模信号为匿名厂商数据，非独立审计。", 9, False, CORAL, 0),
    ])
    rect(s, Inches(4.95), Inches(1.55), Inches(3.5), Inches(3.15), WHITE, radius=True)
    maybe_pic(s, ASSETS / "nokia-autonomous-network-fabric.png", Inches(5.05), Inches(1.7), width=Inches(3.3), height=Inches(2.85))
    card(s, Inches(8.6), Inches(1.55), Inches(4.3), Inches(3.15), [
        ("相邻竞合者", 11, True, NAVY, 3),
        ("云：控制运行时/模型/开发者；Telecom Data Fabric多为Private Preview。", 9, False, MUTED, 3),
        ("湖仓：可替换后端，缺RAN快环与动作安全。", 9, False, MUTED, 3),
        ("API渠道：聚合/计费倒逼内部产品化与行为一致。", 9, False, MUTED, 0),
    ], fill=SOFT)
    items = [
        ("共识已形成", "分布式接入、近源、统一语义/QoD、产品化，并与AI/Agent和闭环证据关联"),
        ("标准未冻结", "独立数据面/新NF/增强既有/混合均未定型"),
        ("产品先于6G", "EIAP/AUTINOps/AN Fabric用5G-A/OAM验证部分能力"),
        ("成熟度不均", "采集湖仓流目录成熟；主动元数据/QoD/Agent/高风险闭环仍G1—G3"),
        ("竞争跨层", "设备商争语义执行；云争运行时；聚合商争渠道；运营商争主权"),
        ("尚无绝对领先", "无中立证据证明任一家多运营商多厂商广泛部署完整跨域Fabric"),
    ]
    for i, (t, b) in enumerate(items):
        card(s, Inches(0.35 + (i % 6) * 2.15), Inches(5.0), Inches(2.08), Inches(1.7), [
            (t, 11, True, NAVY, 3),
            (b, 9, False, MUTED, 0),
        ], accent=TEAL)
    finish(s, "02 行业环境", "02 Nokia与现状小结")

    # ===== 14 Trends overview =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "03  TREND CAUSAL CHAIN", "产业趋势：T1–T8 是因果链，不是八个平行产品")
    thesis_box(s, "数据越分布，越需要元数据事件；机器消费越多，越需要显式语义与契约；自动化权限越大，越需要策略、证明和可回滚性；跨组织流通越多，越需要数据空间式信任与商业运营。产业断言：不会以大一统平台胜出，而以「统一控制语义 + 分布式执行 + 场景化数据产品 + 有界智能」进入网络。")
    chain = [
        ("分布式数据生产", "RAN/Core/Edge/ISAC"),
        ("元数据事件化", "发现变化与影响"),
        ("语义与产品契约", "机器可理解可组合"),
        ("智能/Agent行动", "建议→低风险执行"),
        ("可信控制与生态", "策略·证明·回滚·计量"),
    ]
    for i, (t, d) in enumerate(chain):
        card(s, Inches(0.35 + i * 2.58), Inches(1.65), Inches(2.45), Inches(1.2), [
            (t, 12, True, WHITE, 3),
            (d, 10, False, LIME, 0),
        ], fill=NAVY)
        if i < 4:
            a = s.shapes.add_textbox(Inches(2.7 + i * 2.58), Inches(2.05), Inches(0.28), Inches(0.3))
            write_box(a, [("→", 14, True, CORAL)])
    mains = [
        ("工程主线", "从管道到事件化控制面：批/流/虚拟化继续存在，差异化转向元数据激活、跨接口语义、QoD、影响分析与策略强制。"),
        ("标准主线", "从研究问题到规范能力包：Rel-20验证候选，Rel-21选最小可标准化集合；统一网元概率低于统一语义/接口/流程概率。"),
        ("产品主线", "从平台功能到场景闭环：客户不为「数据湖规模」买单，而为故障定位、能耗、移动性、QoE、模型泛化与开放收入的可测结果买单。"),
        ("竞争主线", "从设备份额到控制点份额：谁掌握数据语义、QoD、近源运行时、模型/策略证据链和API渠道，谁更可能获得软件与生态复利。"),
    ]
    for i, (t, b) in enumerate(mains):
        col, row = i % 2, i // 2
        card(s, Inches(0.4 + col * 6.45), Inches(3.1 + row * 1.75), Inches(6.25), Inches(1.6), [
            (t, 13, True, TEAL, 4),
            (b, 11, False, INK, 0),
        ])
    finish(s, "03 产业趋势", "03趋势主线")

    # ===== 15 T1-T2 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "03  T1–T2", "产品化与事件化控制面：数据如何被消费与感知变化")
    thesis_box(s, "T1把输出从原始流升级为带消费者与SLO的数据产品；T2把目录从定期盘点升级为驱动影响分析、契约测试与策略复核的事件网络。二者共同构成可信控制面的输入侧。")
    for i, block in enumerate([
        ("T1 · 高确定性", "网络数据从「采集对象」变为「可运营产品」",
         "价值计量单位将从表、Topic和PB，转为满足明确消费者与SLO的数据产品。",
         "为何成立：AI、NDT和跨域自动化需要稳定契约，而非一次性取数。Nokia Data Suite已把电信语义、血缘、质量和多厂商数据包装为可复用产品；Ericsson AI-ready mesh强调按产品与授权复用；Open Gateway从消费侧要求跨运营商一致行为。",
         "因此：跨域默认输出将转向事件、聚合指标、特征、模型输入与分析服务；每个产品必须携带对象范围、事件时间、语义版本、QoD、用途、保留期、成本和Owner。",
         "落地：先在移动性、节能、故障、QoE和模型训练形成少量高复用产品；用消费者数、契约合规、重复采集减少和单位有效任务成本评价。",
         "领先信号：厂商Data Suite、开放契约、产品目录进入平台基线　｜　反证：只有Marketplace页面无Owner/SLO/退出；长期无人消费仍重复采集。"),
        ("T2 · 高确定性", "元数据事件化成为可信控制面的输入",
         "目录将从「定期盘点」升级为感知结构、配置、质量、权限、成本和模型变化的事件网络。",
         "为何成立：模型与Agent直接消费网络数据时，目录延迟数小时意味着错误结构、过期配置或权限变化进入生产决策。只有把Schema、拓扑、配置、QoD、血缘和使用行为形成事件，才能触发影响分析、契约测试、策略复核和发布闸门。OpenLineage等已验证工程可行性。",
         "6G特殊性：元数据事件还需包含小区/波束/频段、配置版本、事件时间与模型/动作作用域——比企业IT表字段血缘更接近「网络状态证据」。",
         "落地：跨O1/E2/A1/R1和Core统一事件标识与Correlation ID；高风险事件先阻断/审批，低风险自动生成测试或建议；事件总线不能成为快环同步依赖。",
         "领先信号：血缘/配置/QoD变更自动触发CI与策略检查　｜　反证：目录更新仍靠人工扫描，变更后才由事故发现。"),
    ]):
        tag, title, claim, why, therefore, how, proof = block
        card(s, Inches(0.4 + i * 6.45), Inches(1.6), Inches(6.25), Inches(5.1), [
            (tag, 10, True, TEAL, 2),
            (title, 13, True, NAVY, 4),
            (claim, 11, False, INK, 5),
            (why, 9, False, MUTED, 4),
            (therefore, 9, False, MUTED, 4),
            (how, 9, False, INK, 4),
            (proof, 9, False, CORAL, 0),
        ])
    finish(s, "03 产业趋势", "03 T1-T2")

    # ===== 16 T3-T4 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "03  T3–T4", "两级部署与有界智能：如何执行、如何与 Agent 耦合而不失控")
    thesis_box(s, "T3回答部署形态：集中协调与分布执行并存。T4回答智能边界：数据/特征/模型/工具/Agent/策略目录互相引用，但责任不混成无边界「大脑」；快环保留确定性执行。")
    for i, block in enumerate([
        ("T3 · 高确定性", "边缘子Fabric与两级控制将成为主流部署形态",
         "全局控制面管理语义和策略，本地子Fabric保证时限、自治和断链生存。",
         "为何：RAN快环、MEC、专网和NTN无法依赖远端联邦查询；跨域训练、策略与治理又需全局一致。Ericsson动态放置、Nokia域湖仓+公共Fabric、Huawei近源算子——三条路线部署原则趋同。",
         "形态：更可能是「SMO/Non-RT全局控制 + RAN/Edge局部执行」：全局下发版本化模型、策略和数据产品定义，本地执行过滤、特征、推理与受限动作，并异步回传证据。",
         "落地：明确<10ms、Near-RT、Non-RT和训练时标的在线依赖清单；子Fabric支持断链运行、缓存保留、冲突合并和版本回退。",
         "风险：过度中央化→时延/单点；过度分散→语义与策略漂移。反证：所有查询/审批必须经中心，或各域无法共享语义。"),
        ("T4 · 中高确定性", "Data Fabric与AI/Agent Fabric耦合，但不合并成无边界大脑",
         "编排对象从数据管道扩展为数据—模型—算力—动作组合；耦合核心是共同证据图。",
         "为何：AI原生要求训练/验证/部署/推理共同访问可信数据；Agent执行还需工具目录、身份和委托权限。3GPP分别研究Network for AI、AI for Network、数据框架与计算支持。",
         "责任边界：数据有用途/保留责任，模型有性能/漂移责任，Agent有身份/权限/动作责任——不能混成一个系统。",
         "落地：关联数据产品、特征、模型、提示、工具、策略、动作与结果版本；采用IPOE约束Agent运行时；先在OAM/Non-RT辅助与低风险动作落地。",
         "领先信号：数据/模型/Agent目录可关联，权限与审计跨工具传播。反证：只有聊天入口，无身份、作用域和动作证据。"),
    ]):
        tag, title, claim, why, form, how, proof = block
        card(s, Inches(0.4 + i * 6.45), Inches(1.6), Inches(6.25), Inches(5.1), [
            (tag, 10, True, TEAL, 2),
            (title, 13, True, NAVY, 4),
            (claim, 11, False, INK, 5),
            (why, 9, False, MUTED, 4),
            (form, 9, False, MUTED, 4),
            (how, 9, False, INK, 4),
            (proof, 9, False, CORAL, 0),
        ])
    finish(s, "03 产业趋势", "03 T3-T4")

    # ===== 17 T5-T6 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "03  T5–T6", "证据闸门与 QoD：如何证明动作安全、如何计量无线数据可用性")
    thesis_box(s, "T5把NDT从规划可视化推向高风险闭环的强制证据；T6把竞争从连接器数量转向通信级数据可信度。二者决定自动化能否扩大权限、模型与闭环能否被验收。")
    for i, block in enumerate([
        ("T5 · 中高确定性", "NDT由规划工具走向高风险闭环的证据闸门",
         "竞争点从可视化与仿真规模，转向能否为生产动作提供可量化的反事实证据。",
         "为何：网络自治提高后，直接在现网试错不可接受。NDT可在动作前复现拓扑、配置、流量、故障与策略，验证候选动作收益、冲突和副作用。SA5已有NDT管理规范与Rel-20增强研究。",
         "关键边界：孪生不是数学意义上的安全证明。只有输入快照、模型版本、同步滞后、覆盖场景、预测误差和生产回放都可量化，孪生结果才可成为放行动作的证据。",
         "落地：为每类闭环定义保真度、最大同步滞后和可接受误差；「孪生通过」须与灰度、风险预算、熔断和人工接管组合；超过误差门槛退回建议模式。",
         "验证指标：同步滞后、预测误差、场景覆盖、生产回放一致。反证：只有可视化模型，无误差门槛与生产反馈。"),
        ("T6 · 中高确定性", "QoD将成为6G数据标准与产品竞争的核心标尺",
         "无线数据可用性由「在哪个网络状态、以何种采样和误差产生」决定，不能只用完整性和新鲜度描述。",
         "为何：同一KPI在不同小区、频段、波束、UE群、配置版本和采样窗口下不可直接比较；训练数据量很大也可能因覆盖偏差导致模型失效。Huawei提出Data QoS，Nokia强调有效性与电信语义，Ericsson强调质量与生命周期。",
         "演进预期：QoD会先以厂商契约和场景指标出现，再逐步收敛公共属性；短期难形成覆盖所有数据族的单一评分。",
         "落地：按数据族定义质量信封而非一个总分；把采样覆盖、时空代表性、同步、配置、标签来源和训练—服务偏差纳入；与产品SLO、模型验收和动作风险绑定。",
         "反证：只统计记录量、接口成功率或单一完整性分数。"),
    ]):
        tag, title, claim, why, boundary, how, proof = block
        card(s, Inches(0.4 + i * 6.45), Inches(1.6), Inches(6.25), Inches(5.1), [
            (tag, 10, True, TEAL, 2),
            (title, 13, True, NAVY, 4),
            (claim, 11, False, INK, 5),
            (why, 9, False, MUTED, 4),
            (boundary, 9, False, MUTED, 4),
            (how, 9, False, INK, 4),
            (proof, 9, False, CORAL, 0),
        ])
    finish(s, "03 产业趋势", "03 T5-T6")

    # ===== 18 T7-T8 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "03  T7–T8", "商业分层与标准化形态：如何变现、标准会冻结什么")
    thesis_box(s, "T7说明外部API成功取决于内部产品化；T8说明需求会留下，但拓扑与网元名称未必留下——应跟踪能力语义与接口，不押注厂商网元命名。")
    for i, block in enumerate([
        ("T7 · 中等确定性", "内部Data Fabric、数据空间与网络API将形成分层商业链",
         "内部数据治理解决「可信生产」，网络API解决「标准消费」，Operate API和聚合渠道解决「规模运营」。",
         "为何：CAMARA把QoD/位置验证等包装为开发者API；TM Forum Operate API支持伙伴入驻、目录和用量；Aduna/CPaaS/云市场负责跨运营商分发。外部一致性最终依赖内部对象语义、QoD、权限、血缘和计量。",
         "数据空间提供可验证身份、合同协商和用途控制，可补上跨组织交换信任层，但不适合直接进入RAN快环。",
         "落地：先产品化内部数据/分析服务，再映射CAMARA或行业API；计价从每次调用走向交易、会话或业务结果；成功取决于多运营商覆盖、行为一致、责任与撤销机制。",
         "反证：API数量增加但无跨网一致、活跃开发者和收入。"),
        ("T8 · 低—中确定性", "专用数据面概念可能被标准化为「能力集合」，而非固定新网元",
         "需求留下，拓扑和名称未必留下：发现、编排、处理、Pub/Sub、存储、暴露与治理可能分布到新旧功能中。",
         "为何：SA2 KI#21明确研究数据发现/采集/处理/存储/暴露，但不要求所有用例共享同一组NF/服务/流程；SA5研究管理侧框架；Huawei DO/DA/DCP只是候选原型。3GPP通常在复用既有能力、引入新功能与复杂度之间折中。",
         "因此：更可能收敛的是能力、信息模型、生命周期和接口行为；部署上可由增强DCCF/ADRF/NWDAF/MDA、SMO服务、边缘Agent或新NF组合实现。",
         "落地：标准跟踪聚焦能力语义与接口，不押注某个厂商网元名称；产品采用可插拔运行时与多后端；若独立平面增加时延、状态和治理重复，应退回能力增强方案。",
         "反证：3GPP明确冻结独立统一数据平面及必选NF。"),
    ]):
        tag, title, claim, why, mid, how, proof = block
        card(s, Inches(0.4 + i * 6.45), Inches(1.6), Inches(6.25), Inches(5.1), [
            (tag, 10, True, TEAL, 2),
            (title, 13, True, NAVY, 4),
            (claim, 11, False, INK, 5),
            (why, 9, False, MUTED, 4),
            (mid, 9, False, MUTED, 4),
            (how, 9, False, INK, 4),
            (proof, 9, False, CORAL, 0),
        ])
    finish(s, "03 产业趋势", "03 T7-T8")

    # ===== 19 Horizon =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "04  3-YEAR / 5-YEAR", "3年看透不变量，5年看清控制点复利")
    thesis_box(s, "「看透」指在标准冻结前看透不会改变的底层需求与可验证能力；「看清」指到首批6G商用准备期，看清真正形成产品控制点、生态分工和商业闭环的路线。最终验证点：Rel-21文本、跨厂商互操作、现网闭环收益与可持续收入。")
    card(s, Inches(0.4), Inches(1.6), Inches(6.2), Inches(5.1), [
        ("未来 1–3 年 · 看透", 15, True, NAVY, 8),
        ("标准窗口", 12, True, TEAL, 2),
        ("Rel-21于2027确定包与Stage 1，2028完成Stage 2/3主要冻结；数据框架从候选功能向最小规范集合收敛。", 10, False, MUTED, 5),
        ("工程底座", 12, True, TEAL, 2),
        ("元数据事件、跨O1/E2/A1/R1对象映射、RAN QoD、数据契约和模型/动作血缘成为领先试点基线。", 10, False, MUTED, 5),
        ("产品形态", 12, True, TEAL, 2),
        ("SMO/Non-RT数据能力、近源处理运行时、NDT预演、受治理RAG/Agent先以外挂或增强组件落地，不会一次性重构全网。", 10, False, MUTED, 5),
        ("竞争焦点", 12, True, TEAL, 2),
        ("厂商继续提交统一框架/数据面/Agent Fabric主张，但市场将用现网闭环收益和多厂商互操作过滤营销。", 10, False, MUTED, 5),
        ("商业现实", 12, True, TEAL, 2),
        ("Open Gateway供给侧与认证趋于成熟，瓶颈转向行业工作流、跨运营商行为一致和结果定价。", 10, False, MUTED, 0),
    ], fill=SOFT)
    card(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(5.1), [
        ("未来 3–5 年 · 看清", 15, True, LIME, 8),
        ("架构落位", 12, True, CYAN, 2),
        ("控制面集中协调、数据/算力分布执行成为主形态；独立数据面是否存在不再重要，能力是否可组合和可证明更重要。", 10, False, MIST, 5),
        ("智能边界", 12, True, CYAN, 2),
        ("低风险网络动作进入有界自治；高风险动作形成「策略—孪生—灰度—回滚—证据」强制路径；全面无人治理仍不可信。", 10, False, MIST, 5),
        ("产业分层", 12, True, CYAN, 2),
        ("设备商控制网络语义与近源执行；云/数据平台控制模型与开发工具；聚合商控制API分发；运营商要求数据主权和可替换性。", 10, False, MIST, 5),
        ("价值产品", 12, True, CYAN, 2),
        ("移动性、节能、故障、QoE、模型生命周期及ISAC行业数据形成可计量产品；原始数据出售不成为主模式。", 10, False, MIST, 5),
        ("胜负手", 12, True, CYAN, 2),
        ("多厂商语义映射、QoD标准、Agent可信运行时、NDT证据门槛和跨域产品运营能力形成长期壁垒。", 10, False, MIST, 0),
    ], fill=NAVY)
    finish(s, "04 3年/5年", "04时间窗")

    # ===== 20 Ten judgements =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "04  TEN JUDGEMENTS", "十条趋势判断：对 T1–T8 的决策压缩（不是另起三条）")
    thesis_box(s, "十条用于决策筛选：高确定性应进入工程与标准基线；中高/中等确定性进入时间盒验证。它们与第三章八条趋势是「展开论证→压缩决策」关系，不是数量矛盾。")
    judgements = [
        ("01", "高", "6G一定需要通用数据能力，但未必需要一个独立「数据面」网元——采集/处理/发现/暴露/治理明确；实现可能增强既有NF、引入新功能或混合。"),
        ("02", "高", "RAN数据语义与QoD将比原始数据规模更稀缺——可用性取决于时空粒度、采样代表性、配置版本和测量误差。"),
        ("03", "高", "Data Fabric价值重心从「集成」上移到「可信智能控制」——管道和存储商品化，元数据激活、策略、契约、证据与回滚形成差异。"),
        ("04", "高", "快环不会被中央Fabric接管——亚秒与毫秒路径继续本地执行；Fabric异步供应模型/策略/摘要并回收证据。"),
        ("05", "中高", "Agent将先辅助、后有界执行，不会直接跨越治理成熟度——生产动作必须绑定身份、授权、作用域、租约与回滚。"),
        ("06", "中高", "NDT会成为高风险闭环重要证据源，但不是绝对安全证明——保真度、同步滞后与场景覆盖必须量化。"),
        ("07", "高", "Mesh、Fabric、Lakehouse将组合而非决出唯一赢家——Mesh管领域责任，Fabric管跨域控制，Lakehouse承载存储计算，Data Space管组织间信任。"),
        ("08", "中高", "网络API商业化将倒逼内部数据产品化——外部一致性、安全与计量无法建立在内部语义/QoD/血缘不稳定之上。"),
        ("09", "中", "云与数据平台将成为设备商之外的关键竞合方——控制Catalog/模型/Agent工具链，但难独自替代通信实时性与RAN语义资产。"),
        ("10", "中", "跨企业自主签约与完全无人治理五年内仍不会成为普遍现实——法律责任、语义协商、用途强制和异常接管仍要求组织承担责任。"),
    ]
    for i, (n, cert, text) in enumerate(judgements):
        col, row = i % 2, i // 2
        color = TEAL if cert == "高" else (CORAL if cert == "中高" else MUTED)
        card(s, Inches(0.35 + col * 6.45), Inches(1.55 + row * 1.0), Inches(6.3), Inches(0.92), [
            (f"{n}  [{cert}确定性]", 9, True, color, 1),
            (text, 9, False, INK, 0),
        ], accent=color)
    finish(s, "04 3年/5年", "04十条判断")

    # ===== 21 ZTE assets =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "05  ZTE POSITION · ASSETS", "中兴定位起点：强于网络智能产品化，弱于通用 Fabric 完整产品")
    thesis_box(s, "正确战略不是再造对标IBM/Informatica/Nokia Data Suite的通用编织平台，而是把AIR Net数据引擎升级为主航道共用的横向数据控制能力。争夺「RAN/跨域语义 + QoD + 近源执行 + 有界Agent供数」，而非独立「大一统数据面」品牌。", height=Inches(0.72))
    anchors = [
        ("产品锚点 · AIR Net三引擎", "数据引擎、大模型引擎与数字孪生引擎协同，支撑Copilot/Agent与跨域闭环；数据引擎被描述为提供数据资产、AI Ready治理与智能分析。"),
        ("场景锚点 · Fault Agent", "GSMA浙江移动：大/小模型+数字孪生+中兴方案；根因识别>90%，诊断约30→<5分钟，工单-20%。TM Forum披露中国电信多省试点。案例口径，非独立审计。"),
        ("标准锚点 · ETSI ZSM 029", "GS ZSM 029《Data Management Agent for AN》由中国电信、中兴、CAICT、亚信等支持；2026-04工作项采纳；覆盖注册/发现/资产/认证/采集传输/工作流编排——与控制语义高度重合。"),
    ]
    for i, (t, b) in enumerate(anchors):
        card(s, Inches(0.4 + i * 4.2), Inches(1.65), Inches(4.05), Inches(2.35), [
            (t, 12, True, NAVY, 5),
            (b, 10, False, MUTED, 0),
        ], accent=TEAL)
    evid = [
        ("5G数据底座", "ZXUN CUDR正式产品，统一UDR/UDSF及多代网络/IMS数据；前代vSDM有Banglalink具名商用——证明电信数据管理基础，不是通用Fabric。"),
        ("开放与互操作", "TM Forum目录将AIR Net/U-Smart-Net列为Ready for ODA——证明架构对齐，不等于所有组件均完成互操作认证。"),
        ("Agent框架", "Co-Sight公开发布低代码、模块化、多Agent协同及第三方插件；Reasoner/Actor分离与MCP/A2A主张对齐有界智能。"),
        ("生态与变现", "GSMA中国移动QoD：CAMARA+中兴NEF；截至2024-10有25家企业客户、月调用超3400万次及收入——最强商业锚点。"),
    ]
    for i, (t, b) in enumerate(evid):
        col, row = i % 2, i // 2
        card(s, Inches(0.4 + col * 6.45), Inches(4.2 + row * 1.25), Inches(6.25), Inches(1.15), [
            (t, 11, True, CORAL, 2),
            (b, 9, False, MUTED, 0),
        ])
    finish(s, "05 公司定位", "05资产锚点")

    # ===== 22 Evidence maturity =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "05  EVIDENCE MATURITY E0–E4", "证据纪律：把「有什么」与「证明了什么」分开")
    thesis_box(s, "后续路线图只把E3/E4视为可复用资产，把E1/E2视为放大器，把E0视为待验证假设。任何对外材料都不应把「AIR Net有数据引擎」直接等价为「中兴已有完整6G Data Fabric」。能力连续 ≠ 已商用完整产品。")
    add_table(
        s, Inches(0.25), Inches(1.55), Inches(12.8), Inches(5.2),
        ["证据层", "中兴公开锚点", "可以支持", "不能外推", "下一步补证"],
        [
            ["E4 API", "中国移动QoD：CAMARA+中兴NEF；25客户、月调用>3400万、收入", "API产品化与高频商业运营有强证据", "不能归因中兴独占；不证明内部Fabric统一", "跨网一致、上游契约、成本利润"],
            ["E4 故障", "浙江移动：根因>90%、诊断30→<5分钟、工单-20%", "具名生产场景可量化价值", "单场景≠全网L4或完整Fabric", "动作证据、回滚、独立审计"],
            ["E3 产品", "CUDR、AIR Net、Co-Sight、Ready for ODA", "有产品载体", "不证明同一控制语义/跨厂商数据面", "边界、兼容矩阵、互操作报告"],
            ["E2 试点", "中国电信多省Fault Agent（piloted）", "故障域方法可跨区域验证", "试点≠规模商用，难迁RAN快环", "持续使用、复用、单位成本"],
            ["E1 标准", "ZSM029 Supporting Org；SA5 Unified data framework输入", "具备正式标准入口", "工作项/提案≠最终规范或主导权", "参考实现、联署、进工作文本"],
            ["E0 假设", "RAN QoD、近源Data Agent、跨接口元数据事件、可信供数控制面", "由R1–R8与T1–T8推导的合理主攻", "无足够公开证据称现有产品", "内部盘点、PoC、付费意愿三重验证"],
        ],
        col_w=[Inches(1.2), Inches(3.8), Inches(2.5), Inches(2.7), Inches(2.6)],
        font_size=8,
    )
    finish(s, "05 公司定位", "05证据分层")

    # ===== 23 Competition coords + C1-C4 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "05  CONTROL POINTS C1–C4", "相对坐标与目标定位：四个必须掌握的控制点，三个开放边界")
    thesis_box(s, "不站队Mesh，也不复制Data Plane命名。工作定位（研究命名）：Telecom Trusted Data Control Fabric——统一理解、质量、契约与证据，不统一所有存储、模型、协议与毫秒级执行。")
    add_table(
        s, Inches(0.3), Inches(1.55), Inches(12.7), Inches(1.7),
        ["维度", "Ericsson", "Huawei", "Nokia", "对中兴含义"],
        [
            ["叙事", "联邦+EIAP/DataOps+API", "DO/DA/DCP+AUTINOps", "AN Fabric/Data Suite最清晰", "讲「数据引擎外溢」，避网元命名战"],
            ["短板", "独立网元弱；组合复杂", "研究与商用易混读", "跨厂商统一缺中立证明", "短板：通用编织完整度、6G框架话语权、RAN QoD benchmark"],
        ],
        col_w=[Inches(1.1), Inches(2.6), Inches(2.6), Inches(2.8), Inches(3.6)],
        font_size=8,
    )
    ctrls = [
        ("C1 语义与QoD", "对象映射·质量信封·配置/时空/采样", "必须主导", TEAL),
        ("C2 产品/契约/主权", "Owner·SLO·用途·驻留·保留·删除传播", "电信模板", CYAN),
        ("C3 近源Data Agent", "过滤·特征·缓存·本地策略·断链", "差异化主场", NAVY),
        ("C4 Agent/NDT证据", "身份·策略·仿真·动作·回滚", "与生态联合", CORAL),
    ]
    for i, (t, d, role, c) in enumerate(ctrls):
        card(s, Inches(0.35 + i * 3.2), Inches(3.45), Inches(3.05), Inches(1.7), [
            (t, 12, True, WHITE, 4),
            (d, 10, False, MIST, 5),
            (role, 11, True, LIME, 0),
        ], fill=c)
    opens = [
        ("开放A 湖仓/Catalog", "开放格式、可替换后端，不重建通用栈"),
        ("开放B 模型/Agent", "多模型多Agent；统一身份与证据"),
        ("开放C API/渠道", "对接CAMARA/聚合商，不垄断消费入口"),
    ]
    for i, (t, b) in enumerate(opens):
        card(s, Inches(0.4 + i * 4.2), Inches(5.35), Inches(4.05), Inches(1.3), [
            (t, 12, True, CORAL, 3),
            (b, 11, False, MUTED, 0),
        ], fill=WARN)
    finish(s, "05 公司定位", "05控制点")

    # ===== 24 Capability + conclusion =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "05  CAPABILITY & CONCLUSION", "能力自评与定位结论：成为横向组件，不当通用挑战者")
    thesis_box(s, "用AIR Net数据引擎抢占「电信语义—QoD—近源执行—有界供数」四位一体控制点；以ZSM029/SA2·SA5把接口与语义写成标准事实，以场景闭环证明价值——而不是先宣称拥有完整6G Data Fabric。")
    scores = [
        ("自智/Agent场景", "4.0", "AIR Net+Fault Agent有试点与认证"),
        ("语义/KG生产化", "3.5", "故障域强；缺跨RAN/Core标准语义"),
        ("标准接口话语权", "3.0", "ZSM029在手；SA2/SA5需加码"),
        ("通用Fabric产品", "1.5", "不宜与IT巨头/Nokia比完整度"),
    ]
    for i, (t, v, n) in enumerate(scores):
        card(s, Inches(0.4 + i * 3.2), Inches(1.6), Inches(3.05), Inches(1.25), [
            (t, 10, False, MUTED, 1),
            (v, 20, True, TEAL, 1),
            (n, 9, False, MUTED, 0),
        ])
    add_table(
        s, Inches(0.3), Inches(3.0), Inches(12.7), Inches(2.5),
        ["能力议题", "位置", "映射", "含义"],
        [
            ["数据产品化与契约", "邻近", "R8·T1/T7", "从资产目录升级为Owner/SLO/QoD/退出"],
            ["元数据事件化", "探索", "R1/R5·T2", "优先变更→影响分析与血缘闸门"],
            ["RAN/边缘子Fabric", "有载体", "R2/R3·T3", "优势×空白重叠，适合时间盒PoC"],
            ["Agent有界执行", "相对亮点", "R7·T4", "最短延长线：用数Agent→供数工具"],
            ["NDT证据闸门", "邻近", "R5/R7·T5", "与高风险变更/节能/故障闭环绑定"],
            ["RAN QoD", "关键缺口", "R4·T6", "差异化必争标尺"],
            ["用途主权强制", "相邻", "R6·T4/T7", "目录标签→策略即代码与删除传播"],
            ["API/数据空间商业", "入口在手", "R8·T7/T8", "外部API成功取决于内部语义/QoD/计量"],
        ],
        col_w=[Inches(2.6), Inches(1.4), Inches(1.8), Inches(6.9)],
        font_size=8,
    )
    card(s, Inches(0.4), Inches(5.7), Inches(6.2), Inches(1.05), [
        ("应成为：主航道共用横向数据能力组件", 11, True, TEAL, 2),
        ("服务无线/核网/承载/SMO/行业；统一控制语义+分布式执行；数据引擎升级为产品/QoD/证据/供数底座。", 9, False, MUTED, 0),
    ], fill=SOFT)
    card(s, Inches(6.8), Inches(5.7), Inches(6.1), Inches(1.05), [
        ("不应成为：通用Data Fabric平台挑战者", 11, True, CORAL, 2),
        ("不打IT湖仓目录；不复制华为网元命名/不机械站队mesh；不承诺编织进入亚毫秒快环。", 9, False, MUTED, 0),
    ], fill=WARN)
    finish(s, "05 公司定位", "05能力结论")

    # ===== 25 Opportunity scoring =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "06  OPPORTUNITY PORTFOLIO", "机会组合：六维加权筛选，避免因技术新颖立项")
    thesis_box(s, "排序维度：战略契合25%·存量复用20%·客户牵引20%·标准窗口15%·差异化10%·执行可控10%。分数表达相对优先级，不表达市场规模；O3必须服务O2/O1/O7/O8，O8是跨域供数合规前提，O6是验证场而非先建平台的理由。")
    add_table(
        s, Inches(0.25), Inches(1.55), Inches(12.8), Inches(5.2),
        ["机会", "一句话", "契合", "复用", "牵引", "标准", "差异", "可控", "加权", "角色"],
        [
            ["O2", "数据引擎→6G数据能力组件外溢", "5", "5", "4", "4", "4", "5", "4.55", "P0主航道"],
            ["O3", "ZSM029+SA2/SA5标准双入口", "5", "4", "3", "5", "4", "5", "4.30", "P0放大器"],
            ["O7", "NDT+Agent高风险动作证据闸门", "5", "4", "4", "3", "4", "4", "4.10", "P1差异化"],
            ["O1", "RAN近源流式元数据/子Fabric PoC", "5", "4", "3", "4", "5", "3", "4.05", "P1技术期权"],
            ["O6", "国内运营商显式数据产品试点", "4", "4", "4", "3", "3", "4", "3.75", "P1商业验证"],
            ["O8", "跨路径数据主权与策略即代码", "4", "3", "4", "3", "3", "3", "3.45", "P1合规底座"],
            ["O4", "承载SLA与数据策略语义映射", "4", "3", "3", "3", "4", "3", "3.35", "P2路标接口"],
            ["O5", "终端/行业侧轻量数据代理", "3", "2", "2", "2", "4", "2", "2.45", "P3观察期权"],
        ],
        col_w=[Inches(0.7), Inches(4.0), Inches(0.7), Inches(0.7), Inches(0.7), Inches(0.7), Inches(0.7), Inches(0.7), Inches(0.9), Inches(1.9)],
        font_size=9,
    )
    finish(s, "06 机会选择", "06机会评分")

    # ===== 26 O2 O3 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "06  MAIN · O2 / O3", "首选主攻：数据引擎外溢 + 标准双入口（共享核心）")
    thesis_box(s, "O2是增量成本最低、技术连续性最好的主攻；O3是增量投入中杠杆最高的一类。二者必须由同一核心团队负责：标准文稿来自可运行代码，核心代码对齐标准对象。")
    card(s, Inches(0.4), Inches(1.6), Inches(6.25), Inches(5.1), [
        ("O2 · 数据引擎外溢为6G数据能力组件 · 加权4.55", 13, True, NAVY, 5),
        ("断言：若不做，数据引擎易固化在运维单场景，6G竞标中「数据与智能层」将缺少可复用底座。", 10, False, CORAL, 5),
        ("场景：①网络域AI/Agent可信数据底座；②SMO/Non-RT统一目录与契约；③NDT/故障/节能闭环供数。先服务主航道，再对外可售。", 10, False, MUTED, 4),
        ("技术：LLM辅助语义层（规范/MIB/YANG→约束抽取→图谱）；五类产品（移动性/故障/节能/QoE/模型特征）；目录/血缘/订阅/策略封装为MCP/OpenAPI工具。", 10, False, MUTED, 4),
        ("买方任务：无线跨O1/E2对齐与质量信封；OSS故障/节能/投诉复用；核网/NEF/API上游可信供给。", 10, False, MUTED, 4),
        ("最小交付：语义Registry+五类模板；四类工具接口；两种异构后端、两条产品线可移植验证。", 10, False, INK, 4),
        ("准入：至少两条产品线指定Owner并提供真实数据与消费者。", 10, False, TEAL, 3),
        ("退出：9个月仍不能跨第二场景复用，或每次接入仍需项目级定制→降级为OSS内部能力。", 10, False, CORAL, 3),
        ("指标：取数周期、重复采集下降、复用率、映射准确率、契约覆盖、血缘回放、进入主航道版本清单。", 10, False, MUTED, 0),
    ])
    card(s, Inches(6.85), Inches(1.6), Inches(6.05), Inches(5.1), [
        ("O3 · ZSM029 + SA2/SA5 标准双入口 · 加权4.30", 13, True, NAVY, 5),
        ("断言：错过Rel-20/21收敛窗口，后续只能适配他人接口；做对了，O2/O1/O7组件可直接获得标准对齐规格。", 10, False, CORAL, 5),
        ("场景：ZSM029已采纳、规范早期起草；SA2 KI#21与SA5 DMFW仍Draft。价值是用参考实现影响互操作语义，而非预设最终文本。", 10, False, MUTED, 4),
        ("技术：ZSM029贡献DMA参考实现（注册/发现/编排）；SA2/SA5主攻编织接口双栈、语义本体映射、秒级近实时时效与QoD属性；不纠缠独立数据面网元是否必选。", 10, False, MUTED, 4),
        ("提案包：语义/QoD信息模型；DMA服务与生命周期；DCP类总线/DCCF-NWDAF类增强双栈映射。", 10, False, INK, 4),
        ("联合机制：每项提案绑定「运营商痛点+可运行参考实现+两家异构后端/域验证」；优先与共同支持组织及国内运营商联合。", 10, False, MUTED, 4),
        ("准入：产品Owner承诺实现提案对象。", 10, False, TEAL, 3),
        ("退出：连续两个会议周期无法获得运营商联署或工作组讨论→收缩议题，不继续铺开术语。", 10, False, CORAL, 3),
        ("KPI：有效提案、联合署名、进入基线文本、互操作测试项；不以篇数代替条款采纳和产品落地。", 10, False, MUTED, 0),
    ], fill=SOFT)
    finish(s, "06 机会选择", "06 O2-O3")

    # ===== 27 O1 O7 O8 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "06  MAIN · O1 / O7 / O8", "并行验证：近源差异、动作证据与主权强制（共用O2内核）")
    thesis_box(s, "三者分别验证「分布式近源执行」「高风险动作证据」「跨路径主权强制」，不另造目录与策略体系。中兴最优解：数据引擎外溢+标准卡位+RAN近源QoD+动作/主权双证据。")
    ops = [
        ("O1 · 4.05 · 中期押注", "RAN近源流式元数据与子Fabric PoC",
         "不做亚毫秒控制替换，做站点侧轻量标注+秒级流上QoD/标签+异步回传证据。",
         "边界：一种高频遥测+一种消费者；不进<10ms同步；不建中心化原始湖；不承诺联邦查询替代本地状态。",
         "验收：标注时延与CPU/内存；字节减量、新鲜度、断链可用、版本回退；QoD对模型准确率/闭环误动作改善。",
         "准入：无线线提供真实链路、资源预算和消费者。退出：收益不覆盖站点开销或必须侵入快环→停止产品化。",
         "护城河：电信对象上下文+站点运行触点+可量化QoD。"),
        ("O7 · 4.10 · 并行主攻", "NDT+Agent高风险动作证据闸门",
         "将孪生、Fault Agent/Co-Sight与动作接口连成「建议—仿真—灰度—执行—回滚—学习」强制链。",
         "买方买的是降低错误变更、缩短审批、可证明回滚，不是「一个孪生」。先Non-RT/OAM与低风险动作。",
         "最小产品：数据/模型/策略/动作快照；保真度/同步滞后/场景覆盖门槛；身份委托、风险预算、灰度与熔断。",
         "准入：有可回放历史和明确损失函数。退出：误差长期不可标定或收益不超人工基线→保留分析辅助。",
         "指标：预测误差、同步滞后、越权拦截、审批周期、灰度失败率、回滚成功、证据完整率。"),
        ("O8 · 3.45 · 必备底座", "跨路径数据主权与策略即代码",
         "用途/同意/驻留/保留/删除从目录标签变成运行时策略，贯穿查询、API、导出、特征/向量与训练。",
         "是O2/O6/Open Gateway准入条件，不单独包装「合规大平台」。模型只辅助解释，规则须可测试、可审批、可回滚。",
         "最小产品：策略DSL；查询/API/导出/训练四类执行点适配器；决策、执行结果与例外审批证据。",
         "准入：一类高敏数据+至少两条不同消费路径。退出：无法脱离项目定制或开销超预算→收缩为管理域参考架构。",
         "指标：跨路径一致率、拒绝/越权覆盖、删除传播时长、驻留违规、例外闭环、策略执行P99。"),
    ]
    for i, (tag, title, a, b, c, d, e) in enumerate(ops):
        card(s, Inches(0.35 + i * 4.3), Inches(1.6), Inches(4.15), Inches(5.1), [
            (tag, 10, True, TEAL, 2),
            (title, 12, True, NAVY, 4),
            (a, 9, False, INK, 4),
            (b, 9, False, MUTED, 4),
            (c, 9, False, MUTED, 4),
            (d, 9, False, CORAL, 4),
            (e, 9, False, TEAL, 0),
        ])
    finish(s, "06 机会选择", "06 O1-O7-O8")

    # ===== 28 Supporting + resources =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "06  SUPPORTING & DEPENDENCY", "配合项有触发条件；资源不是八支平行队伍")
    thesis_box(s, "依赖顺序：O2共用内核→O3放大并约束O2→O1/O7/O8分别提供近源、动作与主权证据→O6验证客户是否为结果付费；O4/O5只能消费前述成果，不能反向拉出两套新平台。")
    add_table(
        s, Inches(0.3), Inches(1.55), Inches(12.7), Inches(2.3),
        ["机会", "当前动作", "升级触发", "保持/降级", "禁止项"],
        [
            ["O6试点", "围绕故障/节能/QoE/API选一明确业务问题；先定基线、对照与收益归属", "业务Owner、多厂商数据、结果计量、联合标准/案例", "只有「建平台」预算则退回O2内部自用", "不得以单供应商演示替代跨厂商适配与ROI审计"],
            ["O4承载", "定义数据产品时效/可靠/主权属性到承载策略的最小语义，不自建承载编排", "承载线有接口Owner，且O1/O6需端到端保证", "标准未收敛或不愿为差异SLA付费→仅保留草案", "不得让管理层成快环在线单点"],
            ["O5终端", "跟踪ISAC/AI终端/隐私计算；复用C1/C2契约，不开发独立平台", "明确终端数据产品、端侧预算、隐私责任、可付费客户", "仅有数据量故事则保持观察", "不得默认采集原始感知/个人数据"],
        ],
        col_w=[Inches(1.1), Inches(3.5), Inches(2.9), Inches(2.7), Inches(2.5)],
        font_size=8,
    )
    card(s, Inches(0.4), Inches(4.05), Inches(4.05), Inches(2.6), [
        ("共享核心 55–65%", 13, True, WHITE, 4),
        ("O2 + O3", 15, True, LIME, 4),
        ("语义、QoD、契约、工具接口与标准参考实现同一团队；标准↔代码互相喂养。", 11, False, MIST, 0),
    ], fill=NAVY)
    card(s, Inches(4.65), Inches(4.05), Inches(4.05), Inches(2.6), [
        ("差异验证 25–35%", 13, True, NAVY, 4),
        ("O1 + O7 + O8", 15, True, TEAL, 4),
        ("分别验证近源、动作证据与主权强制；共用O2语义/契约，不另造目录体系。", 11, False, MUTED, 0),
    ], fill=SOFT)
    card(s, Inches(8.9), Inches(4.05), Inches(4.0), Inches(2.6), [
        ("客户与期权 10–15%", 13, True, NAVY, 4),
        ("O6 + O4 + O5", 15, True, CORAL, 4),
        ("O6商业证伪；O4/O5仅接口与观察。明确不选：通用企业Fabric、独立数据面胜负手、ZB全量湖、无边界Agent改网。", 11, False, MUTED, 0),
    ], fill=WARN)
    finish(s, "06 机会选择", "06配合与资源")

    # ===== 29 Strategy 1-11 =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "07  STRATEGY 1–11", "具体策略：先卡位叙事与标准，再加深原型，后铺开近源与试点")
    thesis_box(s, "策略总断言：先把数据引擎做成主航道共用的「可信供数控制面」——用标准写下语义与接口，用近源守住RAN QoD，用动作证据与主权策略约束自治和跨组织流通；其余一律后置。【D】条线可推，【R】需公司决策。")
    add_table(
        s, Inches(0.25), Inches(1.55), Inches(12.8), Inches(5.2),
        ["#", "选择", "具体策略", "不做的风险", "备注"],
        [
            ["1", "【D】", "形成「数据引擎→6G数据能力组件」对内白皮书，对齐R/T/友商/主攻位", "引擎停在故障单点，竞标缺横向叙事", "P0·90天"],
            ["2", "【D】", "ZSM029参考实现；SA2/SA5提交语义映射、双栈、时效与QoD（≥2篇有效输入）", "接口术语被友商单极定义", "P0·Rel窗口"],
            ["3", "【D】", "以Fault Agent KG为种子建语义抽取与五类高复用产品（Owner/SLO/QoD/血缘）", "只有平台能力、无可运营产品", "P0/P1"],
            ["4", "【D】", "目录/血缘/订阅/策略封装为MCP/OpenAPI；IPOE；Reasoner/Actor分离", "Agent只能聊天取数，难进生产", "P1"],
            ["5", "【D+R】", "站点/边缘轻量代理：流内标注、质量信封、本地缓存、异步证据；不进<10ms", "丢失相对云不可替代战场，缺QoD实证", "P1·12–18月"],
            ["6", "【R】", "设立「6G数据能力」虚拟团队，统一对象/QoD/北向，禁各线重复造目录", "横向碎片化，客户看到多套平台", "组织前提"],
            ["7", "【D】", "孪生接入变更/Agent强制路径：快照、版本、滞后、误差、灰度、回滚成放行条件", "孪生停在可视化，难支撑有界自治", "P1/P2"],
            ["8", "【R】", "联合运营商显式数据产品试点+独立ROI；承载映射与行业API写入6G路标", "国内首发与路标位次被占；Gateway有入口无供给", "P1"],
            ["9", "【D+R】", "自研语义/QoD/近源/证据；合作湖仓Catalog模型隐私计算；采购连接器扫描可观察", "要么重复造通用底座，要么交出电信控制点", "立项前门"],
            ["10", "【R】", "对O1/O2/O7/O8建G0–G4；每季按复用/价值/标准/开销/证据调资源，允许停止", "PoC长期化、资源平均分配、标准产品脱节", "组合治理"],
            ["11", "【D+R】", "用途/同意/驻留/保留/删除策略DSL与跨查询/API/导出/向量/训练执行点", "治理停在目录标签；进入缓存/向量/模型后无法持续约束", "P1·R6/O8"],
        ],
        col_w=[Inches(0.45), Inches(0.9), Inches(5.6), Inches(3.7), Inches(1.15)],
        font_size=7,
    )
    finish(s, "07 策略建议", "07策略表")

    # ===== 30 BPB rhythm gates =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "07  BUILD / RHYTHM / GATES", "把钱花在不可替代处；每一步都能加码、转向或停止")
    thesis_box(s, "自研电信语义、QoD、近源运行时和动作证据；合作湖仓/Catalog/模型/隐私计算；采购通用扫描、连接器与可观察工具。全部策略遵循「需求驱动+证据分层+时间盒」。")
    add_table(
        s, Inches(0.3), Inches(1.55), Inches(12.7), Inches(2.0),
        ["能力", "选择", "原因", "架构要求"],
        [
            ["电信对象语义与QoD", "Build", "网元/配置/测量/动作上下文是设备商核心知识资产", "版本化语义包；跨厂商映射可插拔；质量信封可测试"],
            ["近源Data Agent", "Build", "需站点/边缘资源、升级、断链和时标工程能力", "轻量、沙箱、可远程回退；不成快环单点"],
            ["数据产品/主权/动作证据", "Build", "决定复用、责任、持续用途控制与有界自治", "开放schema/API；跨路径执行；证据可导出"],
            ["湖仓/Catalog/流批", "Partner", "产业成熟、客户存量多，自研缺规模优势", "开放表格式、多云多后端；核心元数据不被绑死"],
            ["连接器/扫描/可观察", "Buy", "差异低且维护面广", "统一插件SDK；替换成本可量化"],
            ["大模型/Agent/隐私计算", "Partner+选建", "演进快；中兴掌握电信约束、评测与运行证据", "多模型可替换；身份/策略/审计不随模型绑定"],
        ],
        col_w=[Inches(2.3), Inches(1.3), Inches(4.3), Inches(4.8)],
        font_size=8,
    )
    timeline = [
        ("0–3月", "看齐与卡位", "叙事；冻结O2/O3；ZSM029与首批文稿；五类产品Owner"),
        ("3–9月", "原型与证据", "语义/供数MVP；非故障闭环扩展；QoD与主权DSL v0"),
        ("9–18月", "近源与试点", "RAN PoC；1个运营商显式试点；证据与策略进基线"),
        ("18–36月", "产品化收敛", "随Rel-21沉淀进主航道；API只映射已产品化能力"),
    ]
    for i, (y, t, d) in enumerate(timeline):
        card(s, Inches(0.35 + i * 3.2), Inches(3.75), Inches(3.05), Inches(1.25), [
            (y, 10, True, TEAL, 1),
            (t, 12, True, NAVY, 2),
            (d, 9, False, MUTED, 0),
        ])
    gates = [
        ("G0", "资产核验", "五清单；无真实资产不进G1"),
        ("G1", "双场景原型", "两类场景两种后端"),
        ("G2", "工程基线", "性能/权限/证据/回退"),
        ("G3", "客户验证", "Owner+对照；可归因"),
        ("G4", "产品/标准收敛", "两主航道+互操作"),
    ]
    for i, (g, t, d) in enumerate(gates):
        card(s, Inches(0.35 + i * 2.55), Inches(5.2), Inches(2.45), Inches(1.4), [
            (g, 12, True, CORAL, 1),
            (t, 11, True, NAVY, 2),
            (d, 9, False, MUTED, 0),
        ], fill=SOFT)
    finish(s, "07 策略建议", "07节奏阶段门")

    # ===== 31 Metrics + risks =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "07  METRICS & RISK LEDGER", "统一度量四账合一；用前置信号而不是事后解释管理风险")
    thesis_box(s, "若12个月内仍无法形成至少两类可复用数据产品与一篇有效标准输入，应复盘主攻是否过于分散。公开试点指标须可被第三方理解；标准成功≠产品成功；产品演示≠跨厂商互操作。")
    add_table(
        s, Inches(0.25), Inches(1.55), Inches(12.8), Inches(2.9),
        ["维度", "北极星问题", "领先指标", "结果指标", "红线"],
        [
            ["复用", "能力是否从项目代码变成公共组件？", "已签契约产品、接入产品线、异构后端", "复用率、交付周期、重复采集下降", "第二场景仍重写>50%核心逻辑"],
            ["QoD/语义", "数据是否真的更可用、更可解释？", "对象覆盖、质量规则、配置上下文", "映射准确率、漂移/误动作、契约违规", "只统计记录数或单一完整性分"],
            ["运行工程", "分布式能力是否通信级可运行？", "P99、资源、断链/升级/回退演练", "SLO达标、跨域字节下降、事故恢复", "成为快环同步单点或无法回退"],
            ["可信自治", "Agent/NDT能否安全扩大权限？", "身份/策略/证据覆盖、仿真误差、灰度", "越权拦截、误动作、回滚、接管", "生产动作无Owner/作用域/租约/证据"],
            ["主权用途", "限制是否贯穿所有消费路径？", "策略覆盖路径、拒绝测试、删除演练", "跨路径一致率、删除时长、违规", "只在目录挂标签"],
            ["客户价值", "是否为结果而非平台名买单？", "业务Owner、对照、活跃消费者", "MTTR/能耗/QoE/API收入、单位成本", "只有参观/奖项，无持续使用"],
            ["标准", "席位是否转化为产品可用规则？", "联署、讨论采纳、参考实现与测试项", "进入基线文本、互操作、版本需求", "只追文稿数量"],
        ],
        col_w=[Inches(1.2), Inches(3.0), Inches(2.9), Inches(2.9), Inches(2.8)],
        font_size=7,
    )
    risks = [
        ("平台化过度", "无新增消费者→冻结通用功能"),
        ("产品线碎片化", "绕过公共schema→架构否决"),
        ("标准—产品脱节", "无实现→停扩写并绑定Owner"),
        ("近源开销过大", "超预算→减代理或终止O1"),
        ("Agent越权", "证据不全→退回建议模式"),
        ("主权策略失效", "跨路径不一致→停新路径"),
        ("ROI无法归因", "无基线→不进试点/不对外宣称"),
    ]
    for i, (t, d) in enumerate(risks):
        card(s, Inches(0.3 + i * 1.85), Inches(4.65), Inches(1.78), Inches(2.0), [
            (t, 10, True, NAVY, 4),
            (d, 9, False, MUTED, 0),
        ])
    finish(s, "07 策略建议", "07度量风险")

    # ===== 32 AP table =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "08  ACTION PROGRAMS AP0–AP10", "关键AP：按共同控制点组织，而非按部门各自立项")
    thesis_box(s, "每个AP必须同时具备业务/产品Owner、技术Owner、标准接口（如适用）、可运行交付物、量化验收和停止条件。没有真实消费者、没有基线、没有退出门的任务，不进入清单。")
    add_table(
        s, Inches(0.2), Inches(1.5), Inches(12.9), Inches(5.25),
        ["AP", "内容", "牵头", "关键交付物", "验收/决策门", "优先级"],
        [
            ["AP0", "存量资产与重复建设盘点", "组合负责人/战略架构", "代码/接口/数据/标准/客户五清单", "关键资产有Owner；形成BPB决策", "P0"],
            ["AP1", "可信供数控制面章程与参考架构", "数据与智能产品", "C1–C4边界、对象模型、开放边界、版本与商业包装假设", "两主航道签字；明确不做通用湖仓；进G1", "P0"],
            ["AP2", "ZSM029+SA2/SA5联合标准包", "标准部", "DMA参考实现；语义/QoD/时效/双栈提案；互操作测试草案", "≥1运营商联署、≥2有效输入、≥1进工作文本；产品同步", "P0"],
            ["AP3", "电信语义包与QoD质量信封v1", "数据引擎研发", "核心对象ID/关系/版本；采样时空配置误差代表性规则；评测集", "跨无线+OSS/核网；准确率/覆盖/漂移达门槛；可回退", "P0"],
            ["AP4", "五类数据产品与Agent供数工具集", "数据与智能产品", "Owner/SLO/QoD/用途/血缘/成本模板；四类MCP/OpenAPI工具", "两类产品两类消费者两后端；第二场景重写<50%；进G2", "P0"],
            ["AP5", "Agent身份、策略与证据运行时", "自智/Co-Sight", "身份、委托、作用域、租约、工具白名单、证据schema、接管", "高风险100%过策略；越权可拦截；无证据默认拒绝", "P0/P1"],
            ["AP6", "NDT高风险动作证据闸门", "数字孪生/自智", "快照、保真度、滞后、误差、覆盖、灰度、熔断、回滚链", "一高风险场景历史回放+影子运行；超差退回建议", "P1"],
            ["AP7", "RAN近源Data Agent时间盒PoC", "无线研究院/平台", "流内标注、QoD、过滤/特征、本地缓存、断链、异步证据", "资源/P99达标；字节下降；QoD有可归因改善，否则停止", "P1"],
            ["AP8", "运营商灯塔场景与独立ROI", "战略客户/产品经营", "明确业务问题、对照基线、多厂商数据、成本收益模型", "连续活跃；结果可审计；转版本/合同；无Owner不启动", "P1"],
            ["AP9", "主航道产品化与跨厂商互操作", "产品组合管理", "版本清单、兼容矩阵、开放SDK、互操作报告、商业模型", "进≥两条产品线；两异构后端；一第三方Agent；达G4", "P1/P2"],
            ["AP10", "跨路径数据主权与策略即代码", "数据治理/安全", "用途/同意/驻留/保留/删除DSL；四路径执行点；例外与证据", "高敏数据跨≥两路径一致；删除传播与拒绝测试通过", "P1"],
        ],
        col_w=[Inches(0.75), Inches(2.5), Inches(1.7), Inches(3.3), Inches(3.5), Inches(0.9)],
        font_size=7,
    )
    finish(s, "08 关键AP", "08 AP总表")

    # ===== 33 90 days + D + closing  -- wait total was 32, need to fit remaining into last pages =====
    # Recalculate: we have finish calls. Let me count pages created so far...
    # Actually I set total=32 but may have more finish() calls. I'll use dynamic total at end.
    # For now continue with pages 31-32 style - wait we already did 32 as AP.
    # Need pages for 90 days, decisions, and closing. Update total to 34.

    # ===== 90 days + decisions =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "08  90 DAYS · GOVERNANCE · D1–D6", "前90天六个输出、组合治理与公司层待裁决问题")
    thesis_box(s, "AP总断言：先用90天完成「一个Owner、两类数据产品、一套标准—代码映射、两个带退出门的PoC」；随后只对跨场景可复用、客户价值可归因、风险证据完整的能力加码。")
    outs = [
        ("输出1 唯一组合Owner", "Sponsor+组合负责人+C1–C4技术Owner；消除「大家参与、无人负责」"),
        ("输出2 资产事实底稿", "逐项核验AIR Net/引擎/Co-Sight/NDT/NEF/无线边缘可复用代码、接口、许可证与部署证据"),
        ("输出3 两类种子数据产品", "从故障/节能/移动性/QoE/模型特征中选两类；必须有Owner、消费者、基线与退出"),
        ("输出4 标准—代码映射", "ZSM029/SA2/SA5每个主张对应代码模块、信息对象、接口和测试"),
        ("输出5 PoC资源与红线", "冻结O1/O7/O8资源预算、运行边界、指标与停止扳机，先有退出门再投入"),
        ("输出6 首个客户问题", "锁定愿提供多厂商数据与业务基线的运营商问题；没有真实问题则先内部验证"),
    ]
    for i, (t, d) in enumerate(outs):
        col, row = i % 3, i // 3
        card(s, Inches(0.4 + col * 4.2), Inches(1.55 + row * 1.15), Inches(4.05), Inches(1.05), [
            (t, 11, True, NAVY, 2),
            (d, 9, False, MUTED, 0),
        ], accent=TEAL)
    decisions = [
        ("D1", "数据引擎是OSS内部能力还是公司级横向组件？→建议后者，以两产品线复用为条件"),
        ("D2", "谁拥有电信语义与QoD最终定义权？→跨线架构Owner；产品线仅域扩展"),
        ("D3", "标准投入是否绑定参考实现与客户联署？→绑定；无代码与运营商问题不进P0"),
        ("D4", "Agent自动动作允许到哪一级？→低风险有界执行；高风险仿真/灰度/回滚/接管"),
        ("D5", "如何与云/数据平台分工？→掌握语义QoD近源证据；湖仓Catalog模型开放合作"),
        ("D6", "谁对跨路径数据主权执行负责？→治理与安全共有规则；引擎统一执行框架"),
    ]
    for i, (n, t) in enumerate(decisions):
        col, row = i % 2, i // 2
        card(s, Inches(0.4 + col * 6.45), Inches(4.0 + row * 0.85), Inches(6.25), Inches(0.78), [
            (f"{n}  {t}", 9, False, INK, 0),
        ], fill=SOFT if row % 2 == 0 else WHITE)
    gov = s.shapes.add_textbox(Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.35))
    write_box(gov, [("治理层：Sponsor季度裁决 → 组合月度复用与风险 → 架构双周对象/QoD/策略 → 交付两周迭代 → 价值层冻结基线与对外披露边界", 9, False, MUTED)])
    finish(s, "08 关键AP", "08九十天与裁决")

    # ===== Closing =====
    s = new_slide(prs)
    bg(s)
    topbar(s, "08  DO / DON'T / VERIFY", "行动摘要与边界：把逻辑收束为可执行纪律")
    thesis_box(s, "产业断言：6G Data Fabric不会以「大一统平台」胜出，而以「统一控制语义+分布式执行+场景化数据产品+有界智能」进入网络。中兴用控制点与证据竞争，而不是用平台名竞争。")
    add_table(
        s, Inches(0.35), Inches(1.55), Inches(12.6), Inches(2.7),
        ["做什么", "不做什么", "遗留 / 待验证"],
        [
            ["把AIR Net升级为C1–C4共用控制能力；用两类数据产品验证复用", "不另建覆盖所有湖仓/目录/模型/协议的大平台", "内部资产范围、代码权属、跨产品线Owner"],
            ["用ZSM029+SA2/SA5建立语义、QoD、接口与测试话语权", "不押注独立数据面名称；不以提案数量代替采纳", "Rel-20/21最终能力包与运营商联署意愿"],
            ["以O1/O7/O8验证近源差异、可信自治与主权强制，设明确停止线", "Fabric不进快环同步单点；Agent无证据不改网；治理不停在目录标签", "站点开销、QoD真实增益、NDT误差门槛、跨路径策略一致率"],
            ["以O6验证客户是否为场景结果付费，使用独立ROI口径", "不因客户提出「平台」就先做重资产集成项目", "首个灯塔场景、持续使用、试点转版本/合同"],
        ],
        col_w=[Inches(4.3), Inches(4.2), Inches(4.1)],
        font_size=9,
    )
    card(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.15), [
        ("最终边界", 13, True, CORAL, 4),
        ("本胶片基于截至2026年8月的标准进度、公开产品与研究主张推演，不等于3GPP已采纳方案，也不构成中兴正式技术路线或市场承诺。", 11, False, INK, 4),
        ("真正的趋势验证点：Rel-21规范文本 · 跨厂商互操作 · 现网闭环收益 · 可持续商业收入。", 11, True, NAVY, 4),
        ("完整论证、引用与图表见：reports/6g-data-fabric-industry-trends-first-four-sections.html", 10, False, MUTED, 4),
        ("再生本胶片：python3 reports/scripts/build_strategy_pptx.py", 10, False, TEAL, 0),
    ], fill=WARN)
    finish(s, "结语", "结语")

    # Fix footers if page count differs from planned total
    actual = len(prs.slides)
    if actual != total:
        # rewrite page numbers in footers by regenerating is hard; instead update total dynamically next build
        print(f"WARNING: planned {total}, actual {actual}")

    out = Path("/workspace/reports/6g-data-fabric-strategy-deck.pptx")
    prs.save(out)
    print(f"Wrote {out} with {actual} slides")
    for i, name in enumerate(meta, 1):
        print(f"  {i:02d} {name}")
    return out


if __name__ == "__main__":
    build()
