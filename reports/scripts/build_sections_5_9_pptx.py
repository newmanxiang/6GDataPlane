#!/usr/bin/env python3
"""Generate the five-slide internal strategy deck for chapters 5-8.

Output: reports/6g-data-fabric-sections-5-9-deck.pptx
Source: chapters 5-8 of the two HTML strategy reports.

All visible copy is stored in AutoShape or table-cell text frames. Standalone
text boxes are intentionally forbidden so cards remain easy to edit.
Chapter 9 is the HTML-only report conclusion and does not add a sixth slide.
"""

from __future__ import annotations

import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "reports" / "6g-data-fabric-sections-5-9-deck.pptx"

W_IN, H_IN = 13.333, 7.5
W, H = Inches(W_IN), Inches(H_IN)
EMU = 914400
TOTAL = 5

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
CORAL = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xB8, 0x86, 0x0B)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF4, 0xF7, 0xF8)
PALE = RGBColor(0xE8, 0xF0, 0xF2)
LINE = RGBColor(0xC5, 0xD0, 0xD4)
SOFT_TEAL = RGBColor(0xD7, 0xEC, 0xEA)
SOFT_AMBER = RGBColor(0xF6, 0xED, 0xD0)
SOFT_NAVY = RGBColor(0xE6, 0xEC, 0xF2)
SOFT_CORAL = RGBColor(0xF8, 0xE6, 0xE4)
FONT = "WenQuanYi Micro Hei"

CJK_MIN = [500, 500, 560, 650, 600]
CJK_TOTAL_MIN = 3000


def _set_run_east_asia(run) -> None:
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT)


def rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def lined(shape, fill: RGBColor, border: RGBColor, width: float = 0.75) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(width)


def _set_anchor(tf, anchor) -> None:
    mapping = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
    tf._txBody.bodyPr.set("anchor", mapping.get(anchor, "t"))


def fill_shape(
    shape,
    lines,
    *,
    default_size=9.2,
    default_color=SLATE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.02,
    ml=0.10,
    mr=0.08,
    mt=0.05,
    mb=0.04,
) -> None:
    """Write copy into the shape's own text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    _set_anchor(tf, anchor)
    tf.clear()
    for i, item in enumerate(lines):
        if isinstance(item, str):
            item = {"text": item}
        text = item["text"]
        size = item.get("size", default_size)
        bold = item.get("bold", False)
        color = item.get("color", default_color)
        palign = item.get("align", align)
        after = item.get("after", 2)
        pspacing = item.get("spacing", spacing)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = palign
        p.space_before = Pt(0)
        p.space_after = Pt(after)
        p.line_spacing = pspacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
        _set_run_east_asia(run)


def box(
    slide,
    l,
    t,
    w,
    h,
    lines,
    *,
    fill=WHITE,
    border=LINE,
    rounded=True,
    radius=0.07,
    **fill_kw,
):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l),
        Inches(t),
        Inches(w),
        Inches(h),
    )
    if border is None:
        rgb(shp, fill)
    else:
        lined(shp, fill, border)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if lines:
        fill_shape(shp, lines, **fill_kw)
    return shp


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rgb(bg, WHITE)
    return slide


def header(slide, number: str, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.07))
    rgb(bar, TEAL)
    box(
        slide,
        0.38,
        0.10,
        12.55,
        0.54,
        [
            {"text": number, "size": 9.5, "bold": True, "color": TEAL, "after": 1},
            {"text": title, "size": 19.5, "bold": True, "color": NAVY, "after": 0},
        ],
        fill=WHITE,
        border=None,
        rounded=False,
        ml=0,
        mr=0,
        mt=0.01,
        mb=0,
    )


def thesis(slide, text: str, height: float = 0.43) -> None:
    box(
        slide,
        0.38,
        0.68,
        12.55,
        height,
        [{"text": text, "size": 11.5, "bold": True, "color": NAVY, "after": 0}],
        fill=PALE,
        border=TEAL,
        ml=0.14,
        mr=0.12,
        mt=0.04,
        mb=0.03,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def footer(slide, page: int, chapter: str, refs: str = "") -> None:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(7.18), Inches(12.55), Inches(0.012)
    )
    rgb(line, LINE)
    left = f"6G 数据架构 × 智能数据中枢 · 第5–8章内部策略胶片 · {chapter}"
    if refs:
        left += f"    {refs}"
    box(
        slide,
        0.38,
        7.22,
        10.55,
        0.20,
        [{"text": left, "size": 7.5, "color": MUTED, "after": 0}],
        fill=WHITE,
        border=None,
        rounded=False,
        ml=0,
        mr=0,
        mt=0,
        mb=0,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    box(
        slide,
        11.10,
        7.22,
        1.83,
        0.20,
        [{"text": f"{page} / {TOTAL}", "size": 8, "bold": True, "color": MUTED, "align": PP_ALIGN.RIGHT, "after": 0}],
        fill=WHITE,
        border=None,
        rounded=False,
        ml=0,
        mr=0,
        mt=0,
        mb=0,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_table(slide, l, t, w, h, rows, col_w, font=8.2):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(l), Inches(t), Inches(w), Inches(h))
    table = shape.table
    for i, width in enumerate(col_w):
        table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE if r == 0 else MSO_ANCHOR.TOP
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (OFF if r % 2 else WHITE)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.045)
            tf.margin_right = Inches(0.04)
            tf.margin_top = Inches(0.03)
            tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(0)
            p.line_spacing = 1.01
            run = p.add_run()
            run.text = value
            run.font.size = Pt(font)
            run.font.bold = r == 0 or c == 0
            run.font.color.rgb = WHITE if r == 0 else INK
            run.font.name = FONT
            _set_run_east_asia(run)
    return table


def slide_p1(prs) -> None:
    slide = new_slide(prs)
    header(slide, "5.1 + 5.2  ·  PAIN × TREND", "痛点与趋势共同驱动：智能数据中枢不是功能扩张，而是共同底座升级")
    thesis(
        slide,
        "无线数据的成本、质量和交付痛点给出近期必答题；数据产品、事件治理、两级部署和可信Agent给出长期约束。两类驱动收敛到同一选择：把重复能力沉淀为可复用的数据与智能底座。",
    )

    pains = [
        (
            CORAL,
            SOFT_CORAL,
            "生产与接入",
            "基站数据/模型难以扩展新类型和消费场景；缺少智能应用语义；订阅不精准，网元—网管采集通道多达7类；同一数据难以一次采集、多方消费。",
        ),
        (
            AMBER,
            SOFT_AMBER,
            "存储与处理",
            "配置、工参、FM、PM、画像、MR分散在多类数据库和持久卷；内部材料显示UE级数据量约为小区级9倍、占存储约90%；流批端到端仍以小时计。",
        ),
        (
            TEAL,
            SOFT_TEAL,
            "消费与治理",
            "消费侧仍有版本定制；跨域拼接、动态编排和共享不足；缺少统一模型、Owner、数据质量观测与治理工具；语义、本体和Agent工具未形成公共基线。",
        ),
    ]
    for i, (accent, fill, title, body) in enumerate(pains):
        x = 0.38 + i * 4.20
        box(
            slide,
            x,
            1.18,
            4.08,
            1.28,
            [
                {"text": f"痛点 {i + 1}｜{title}", "size": 12, "bold": True, "color": accent, "after": 3},
                {"text": body, "size": 9.2, "color": SLATE, "after": 0},
            ],
            fill=fill,
            border=accent,
            mt=0.07,
        )

    driver_cards = [
        (
            0.38,
            3.05,
            "近期工程压力",
            "重复采集/存储/治理；项目式交付；数据质量不可观测；UE级数据增长放大成本。",
            CORAL,
            SOFT_CORAL,
        ),
        (
            3.66,
            5.99,
            "共同响应：六项横向能力",
            "统一理解：对象、指标、语义、本体\n统一责任：Owner、契约、版本、退出\n统一质量：规则、观测、诊断、治理\n统一供数：NISM、流批、联邦、编排\n分布执行：秒/分/时、近源、断链\n可信行动：身份、策略、证据、回滚",
            TEAL,
            SOFT_TEAL,
        ),
        (
            9.78,
            3.15,
            "长期产业约束",
            "T1/T2：数据产品与元数据事件；T3：全局定义、近源执行；T4/T5：Agent协同与NDT证据；T6：数据质量进入模型与闭环准入。",
            NAVY,
            SOFT_NAVY,
        ),
    ]
    for x, width, title, body, accent, fill in driver_cards:
        box(
            slide,
            x,
            2.57,
            width,
            1.35,
            [
                {"text": title, "size": 11.5, "bold": True, "color": accent, "after": 3},
                {"text": body, "size": 9.1, "color": SLATE, "after": 0, "spacing": 1.01},
            ],
            fill=fill,
            border=accent,
            mt=0.07,
        )

    implications = [
        ("数据供给", "近期：结束多通道、项目式接入。", "长期：产品有Owner/契约/版本/消费者。", "中枢：可订阅、可复用、可运营的数据产品。"),
        ("质量与语义", "近期：统一模型、质量观测和跨域对齐。", "长期：质量/语义变化触发重验与版本处置。", "中枢：数据质量与语义进入Agent准入链。"),
        ("处理与部署", "近期：按秒/分钟/批分时，降低搬运成本。", "长期：全局定义与近源执行并存、断链可运行。", "中枢：湖仓、NDS、网元计算、联邦分层协同。"),
        ("智能消费", "近期：NISM、编排、智能问数服务多Agent。", "长期：工具调用具备身份、策略、证据与回滚。", "中枢：既供数又提供受控工具，不替代域控制器。"),
    ]
    for i, (title, near, future, meaning) in enumerate(implications):
        x = 0.38 + i * 3.15
        box(
            slide,
            x,
            4.04,
            3.03,
            2.21,
            [
                {"text": title, "size": 11.5, "bold": True, "color": NAVY, "after": 4},
                {"text": near, "size": 9.1, "color": SLATE, "after": 4},
                {"text": future, "size": 9.1, "color": SLATE, "after": 4},
                {"text": meaning, "size": 9.2, "bold": True, "color": TEAL, "after": 0},
            ],
            fill=OFF,
            border=LINE,
            mt=0.08,
        )

    box(
        slide,
        0.38,
        6.36,
        12.55,
        0.70,
        [
            {"text": "决策含义", "size": 10.5, "bold": True, "color": CORAL, "after": 2},
            {"text": "近期以降本、提质和复用解决真实痛点；中期以数据产品、分布式处理和Agent消费形成产品闭环；长期以数据质量、语义、近源执行和动作证据建立6G差异化控制点。", "size": 10, "bold": True, "color": NAVY, "after": 0},
        ],
        fill=PALE,
        border=TEAL,
        mt=0.05,
    )
    footer(slide, 1, "5.1+5.2 痛点与趋势", "[S59]")


def slide_p2(prs) -> None:
    slide = new_slide(prs)
    header(slide, "5.3 + 5.4  ·  FOUNDATION × TARGET", "现实基础已存在，目标是把数据引擎升级为6G无线产品横向共用的智能数据中枢")
    thesis(
        slide,
        "智能数据中枢以RAN网管与无线大数据资产为根，完全吸收现有数据引擎；统一可信数据产品、多维数据质量、跨系统共享、分布式处理和数据语义，向上支撑自有Agent与开放生态。",
    )

    assets = [
        (
            TEAL,
            SOFT_TEAL,
            "已有实践｜湖仓、流批、降本",
            "CM/PM/工参/MR统一建模入湖，NGI湖仓已有外场实践；Zeno、Saturn和统一调度形成降本基础。部分项目披露迁移资源下降50%、后续再降30%，仍需按产品独立验收。",
        ),
        (
            NAVY,
            SOFT_NAVY,
            "在建能力｜NISM、治理、语义",
            "统一查询/目录/订阅、NL2SQL/API、质量观测、治理Agent和本体语义已有方案或实践。当前关键是统一对象、Owner、数据质量、接口和版本，而不是继续堆叠组件。",
        ),
        (
            AMBER,
            SOFT_AMBER,
            "差异化预研｜分时、分布、近源",
            "NDS近实时、联邦查询、动态编排、6G新数据通道和Saturn下沉共同指向全局定义、分布执行；必须验证资源预算、断链、版本一致和真实收益。",
        ),
    ]
    for i, (accent, fill, title, body) in enumerate(assets):
        x = 0.38 + i * 4.20
        box(
            slide,
            x,
            1.18,
            4.08,
            1.31,
            [
                {"text": title, "size": 11.5, "bold": True, "color": accent, "after": 3},
                {"text": body, "size": 9.1, "color": SLATE, "after": 0},
            ],
            fill=fill,
            border=accent,
            mt=0.07,
        )

    maturity = [
        ("M1 已有基线", "统一湖仓与流批", "统一入湖、Zeno/Saturn、跨域共部署"),
        ("M2 首版闭环", "治理质量与消费", "Owner、模型、数据质量、NISM、问数"),
        ("M3 在研增强", "分布式中枢", "联邦、服务分层、跨集群、Pipeline"),
        ("M4 时间盒", "近源与网元计算", "NDS近实时、6G流批、Saturn下沉"),
        ("M5 条件扩展", "动作证据与开放", "自有Agent先消费，再开放生态与API"),
    ]
    for i, (stage, title, body) in enumerate(maturity):
        x = 0.38 + i * 2.52
        accent = TEAL if i < 2 else (AMBER if i < 4 else NAVY)
        fill = SOFT_TEAL if i < 2 else (SOFT_AMBER if i < 4 else SOFT_NAVY)
        box(
            slide,
            x,
            2.61,
            2.42,
            1.03,
            [
                {"text": stage, "size": 8.5, "bold": True, "color": accent, "after": 1},
                {"text": title, "size": 10.5, "bold": True, "color": NAVY, "after": 2},
                {"text": body, "size": 8.5, "color": SLATE, "after": 0},
            ],
            fill=fill,
            border=accent,
            mt=0.05,
        )

    chain = [
        ("网络事实与数据源", "基站 · NDS · UME-NM · 外部数据"),
        ("智能数据中枢底座", "湖仓 · 流批 · 联邦 · 编排"),
        ("可信数据控制", "Owner · 语义 · 数据质量 · 版本"),
        ("统一消费服务", "NISM · NL2SQL/API · 工具"),
        ("产品与Agent", "RDE · NGI · IDA · iMind · 生态"),
    ]
    for i, (title, sub) in enumerate(chain):
        x = 0.38 + i * 2.52
        box(
            slide,
            x,
            3.77,
            2.42,
            0.87,
            [
                {"text": title, "size": 10.5, "bold": True, "color": NAVY, "after": 2, "align": PP_ALIGN.CENTER},
                {"text": sub, "size": 8.5, "color": MUTED, "after": 0, "align": PP_ALIGN.CENTER},
            ],
            fill=WHITE,
            border=TEAL if i in (1, 2, 3) else LINE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            mt=0.04,
        )

    box(
        slide,
        0.38,
        4.78,
        8.12,
        2.28,
        [
            {"text": "目标定位｜6G无线产品横向共用的智能数据中枢", "size": 13, "bold": True, "color": TEAL, "after": 5},
            {"text": "· 以RAN网管为内部Owner锚点，由无线大数据与数据引擎共同研发；首个版本目标为2026 Q4，作为内部研发里程碑，不作外部承诺。", "size": 9.6, "color": SLATE, "after": 4},
            {"text": "· 完全吸收现有数据引擎、湖仓、NISM、治理和语义资产，不另建孤立平台；向下连接网元与网络原生执行，向上支撑AIR Net自有Agent、产品应用、第三方Agent与开放API。", "size": 9.6, "color": SLATE, "after": 4},
            {"text": "· 以第二消费者、第二系统、两种后端和正式版本证明横向复用；组件存在、方案完成或单项目效果都不能替代产品成立条件。", "size": 9.6, "color": SLATE, "after": 0},
        ],
        fill=SOFT_TEAL,
        border=TEAL,
        mt=0.10,
        ml=0.14,
    )
    box(
        slide,
        8.64,
        4.78,
        4.29,
        2.28,
        [
            {"text": "边界与非目标", "size": 13, "bold": True, "color": CORAL, "after": 5},
            {"text": "· 不是覆盖公司全部数据、存储、模型和协议的大中台。", "size": 9.4, "color": SLATE, "after": 3},
            {"text": "· 不替代无线、核心网和OSS域控制器，不进入毫秒快环远端同步单点。", "size": 9.4, "color": SLATE, "after": 3},
            {"text": "· 不强制所有产品共享同一存储、运行时或Agent框架。", "size": 9.4, "color": SLATE, "after": 3},
            {"text": "· 通用湖仓、模型生态和渠道保持开放；中兴掌握电信语义、数据质量、近源约束和动作证据。", "size": 9.4, "color": SLATE, "after": 0},
        ],
        fill=SOFT_CORAL,
        border=CORAL,
        mt=0.10,
        ml=0.14,
    )
    footer(slide, 2, "5.3+5.4 现实基础与定位", "[S40][S45][S59]")


def slide_p3(prs) -> None:
    slide = new_slide(prs)
    header(slide, "6  ·  OPPORTUNITY × ARCHITECTURE", "一个目标架构整合机会点与选择：先闭合横向基线，再验证近源差异，最后开放")
    thesis(
        slide,
        "底层解决数据成本和分时运行，中层建立治理、数据质量、语义与服务控制，上层服务自有Agent和开放生态。优先级由架构依赖决定，不平均分配资源。",
    )

    layers = [
        (SOFT_NAVY, NAVY, "L5 产品与Agent消费", "RDE / NGI / IDA / iMind / 自有Agent / 第三方应用", "自有Agent先消费"),
        (SOFT_TEAL, TEAL, "L4 统一数据服务与语义", "NISM查询/目录/订阅 · NL2SQL/API · 本体 · 数据工具", "P0 首版"),
        (TEAL, WHITE, "L3 可信数据控制", "对象 · Owner · 数据产品 · 数据质量 · 血缘 · 版本/事件 · 权限", "P0 控制点"),
        (NAVY, WHITE, "L2 统一与分布式运行", "湖仓 · Zeno/Saturn/Spark · Pipeline · 联邦 · 跨集群 · 成本", "P0/P1 增强"),
        (SOFT_AMBER, AMBER, "L1 近源生产与分时处理", "基站/NDS · 秒/分钟/批 · 精准订阅 · 过滤/特征 · 断链/回证", "P1 时间盒"),
        (SOFT_CORAL, CORAL, "X 动作证据横切", "Agent身份 · 策略 · 数据/模型版本 · NDT · 灰度 · 回滚", "P1 场景验证"),
    ]
    for i, (fill, accent, title, body, status) in enumerate(layers):
        is_dark = fill in (TEAL, NAVY)
        box(
            slide,
            0.38,
            1.18 + i * 0.78,
            5.07,
            0.69,
            [
                {"text": title, "size": 10.5, "bold": True, "color": WHITE if is_dark else NAVY, "after": 1},
                {"text": body, "size": 8.4, "color": WHITE if is_dark else SLATE, "after": 0},
                {"text": status, "size": 8.3, "bold": True, "color": WHITE if is_dark else accent, "after": 0, "align": PP_ALIGN.RIGHT},
            ],
            fill=fill,
            border=accent if not is_dark else None,
            mt=0.04,
            mb=0.02,
        )

    add_table(
        slide,
        5.61,
        1.18,
        7.32,
        4.59,
        [
            ["机会选择", "落点", "最小闭环", "阶段门 / 停止线"],
            [
                "P0 可信数据产品基线",
                "L3+L4",
                "两类种子数据产品；对象、Owner、数据质量、语义、NISM；至少两个消费者。",
                "第二消费者仍重做模型/接口，或无Owner、无质量门、无版本承诺则不成立。",
            ],
            [
                "P0 湖仓与消费整合",
                "L2+L4+L5",
                "复用Zeno/Saturn/湖仓；统一查询/订阅；一个自有Agent贯通生产—消费—反馈。",
                "必须度量重复存储、接入周期、查询性能和资源成本；只搬组件不算升级。",
            ],
            [
                "P1 分布式与近源",
                "L1+L2",
                "一种高频无线数据；秒/分钟/批分时；断链、对账、回退与联邦/就近消费。",
                "收益不覆盖网元/NDS开销、必须依赖中心或侵入快环时退回域内能力。",
            ],
            [
                "P1/P2 动作与开放",
                "X+L5",
                "一个有界Agent动作场景；身份/策略/证据/回滚；成熟后映射第三方Agent和API。",
                "无可回放证据、无人工接管或归因责任不清，不扩大权限和结果承诺。",
            ],
        ],
        col_w=[1.52, 0.84, 2.72, 2.24],
        font=8.15,
    )

    box(
        slide,
        0.38,
        5.91,
        7.78,
        1.15,
        [
            {"text": "选择建议", "size": 11.5, "bold": True, "color": TEAL, "after": 3},
            {"text": "资源首先投向L3/L4可信数据产品基线，与L2已有湖仓/流批及L5自有Agent闭合2026 Q4首版；分布式/近源和动作证据各选一个真实场景时间盒验证。公共能力必须进入无线主航道版本。", "size": 9.4, "color": SLATE, "after": 0},
        ],
        fill=SOFT_TEAL,
        border=TEAL,
        mt=0.07,
    )
    box(
        slide,
        8.30,
        5.91,
        4.63,
        1.15,
        [
            {"text": "不选择", "size": 11.5, "bold": True, "color": CORAL, "after": 3},
            {"text": "不新建覆盖所有域的统一大平台；不同时铺开全部数据类型、Agent和接口；不以目录、接入量或API数量替代复用、数据质量、资源和版本指标。", "size": 9.4, "color": SLATE, "after": 0},
        ],
        fill=SOFT_CORAL,
        border=CORAL,
        mt=0.07,
    )
    footer(slide, 3, "第6章 机会与架构整合", "[S59]")


def slide_p4(prs) -> None:
    slide = new_slide(prs)
    header(slide, "7  ·  STRATEGY TABLE", "公司选择按依赖顺序推进：Q4首版闭合基线，2027验证差异，成熟后开放")
    thesis(
        slide,
        "本页严格使用“我司选择—具体策略（具体动作和时间）—不做的风险”三列表达。P0进入正式版本，P1必须有预算和停止线，P2以前序复用门通过为前提。",
    )
    add_table(
        slide,
        0.38,
        1.18,
        12.55,
        5.45,
        [
            ["我司选择", "具体策略（包含具体动作和时间）", "不做的风险"],
            [
                "P0｜以智能数据中枢统一无线数据能力主线",
                "现在—2026 Q4：由RAN网管牵引、无线大数据与数据引擎承担核心研发，把数据引擎、湖仓、NISM、治理和语义收敛到同一产品边界与版本基线；明确唯一产品Owner、架构Owner及两类种子数据产品，形成首个可用版本。每项能力映射代码、接口、消费者和验收指标。",
                "数据能力继续分散在NGI、NIA、NDS、网管和各Agent中；重复采集、重复存储和项目式治理持续，智能数据中枢退化为新名称而非共同底座。",
            ],
            [
                "P0｜先闭合数据质量、语义和统一消费",
                "2026 Q4首版：以CM/PM与MR/告警等两类高价值数据为种子，落实对象ID、Owner、统一模型、有效性/一致性/准确性/及时性/完整性/唯一性规则及质量诊断；NISM提供查询、订阅和目录，NL2SQL/API或本体语义至少服务一个自有Agent和一个产品应用。2027扩展跨域治理Agent、观测和分布式消费。",
                "Agent继续先取数、再清洗、后解释，可信度与交付周期不可控；统一湖仓只有物理汇聚，没有责任、数据质量和语义控制，无法支撑多Agent复用。",
            ],
            [
                "P1｜时间盒验证分布式、近源和动作证据",
                "2026 Q3–Q4完成方案穿刺，2027按门扩展：选择一种高频UE级数据和一个无线Agent场景，验证NDS分钟级、网元秒级或Saturn下沉，记录CPU/内存、P99、字节减量、断链、对账和回退；同步验证Agent身份、策略、数据/模型版本、NDT/灰度、人工接管与动作回证。收益不覆盖开销或侵入快环即停止。",
                "不做会丢失设备商近源数据与执行优势；无边界铺开则把远端平台和Agent引入快环，增加资源成本、故障耦合和动作责任风险。",
            ],
            [
                "P2｜成熟后开放第三方Agent、API与标准接口",
                "首版复用门通过后：将稳定的数据产品、语义、数据质量和工具接口开放给第三方Agent，并映射NEF/CAMARA Quality on Demand（QoD）等能力API；把ZSM、3GPP SA2/SA5、O-RAN主张绑定参考实现、信息模型和互操作测试。能力服务与结果服务分别核算履约、授权、归因和回滚成本。",
                "过早开放会放大不稳定契约和责任缺口、形成兼容债；长期不开放则形成封闭栈，失去标准话语权、生态入口和跨厂商部署机会。",
            ],
        ],
        col_w=[2.30, 6.56, 3.69],
        font=8.15,
    )
    box(
        slide,
        0.38,
        6.71,
        12.55,
        0.35,
        [{"text": "策略纪律：P0不能被P1/P2抢占资源；内部方案完成、PoC成功或API数量，均不能替代“第二消费者复用 + 数据质量闭环 + 正式版本承诺”。", "size": 9.5, "bold": True, "color": NAVY, "after": 0, "align": PP_ALIGN.CENTER}],
        fill=PALE,
        border=TEAL,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        mt=0.02,
        mb=0.01,
    )
    footer(slide, 4, "第7章 具体策略")


def slide_p5(prs) -> None:
    slide = new_slide(prs)
    header(slide, "8  ·  AP-P0", "最高优先级AP：90天闭合“可信数据产品基线 + 首个Agent消费”")
    thesis(
        slide,
        "目标不是90天完成全部智能数据中枢，而是用两类种子数据产品和两个真实消费者证明：同一对象、数据质量、语义与NISM接口可复用，并可进入2026 Q4正式版本。",
    )

    work = [
        (
            "W1｜第0–30天",
            "冻结边界与责任",
            "指定RAN网管产品Owner、无线大数据技术Owner和跨线架构Owner；盘点代码/接口/数据/消费者；冻结首版包含湖仓复用、数据质量、语义、NISM和两类数据产品。",
        ),
        (
            "W2｜第15–60天",
            "建立可信数据产品",
            "建议CM/PM与MR/告警等已有基础数据族；冻结对象ID、Owner、Schema、版本、用途、消费者、血缘、SLO和六维数据质量规则；建立监控、诊断、修复与回退。",
        ),
        (
            "W3｜第31–75天",
            "贯通消费与语义",
            "用NISM封装目录、查询和订阅，打通两个后端/中枢；完成指标、维度和API语义映射；为一个自有Agent提供NL2SQL/API或结构化工具并记录输入版本。",
        ),
        (
            "W4｜第61–90天",
            "完成复用与版本门",
            "一个Agent与一个非Agent产品共同消费；验证第二消费者无需重建核心模型；对比接入工期、重复存储、查询性能、资源成本及质量缺陷周期；通过后进入Q4版本。",
        ),
    ]
    for i, (phase, title, body) in enumerate(work):
        x = 0.38 + i * 3.15
        box(
            slide,
            x,
            1.18,
            3.03,
            1.52,
            [
                {"text": phase, "size": 8.8, "bold": True, "color": TEAL, "after": 1},
                {"text": title, "size": 11.5, "bold": True, "color": NAVY, "after": 3},
                {"text": body, "size": 8.7, "color": SLATE, "after": 0},
            ],
            fill=SOFT_TEAL if i < 2 else SOFT_NAVY,
            border=TEAL if i < 2 else NAVY,
            mt=0.06,
        )

    add_table(
        slide,
        0.38,
        2.83,
        12.55,
        2.89,
        [
            ["验收域", "90天必须交付", "通过门", "不通过如何处置"],
            [
                "产品",
                "首版章程、边界、版本清单、Owner/RACI、两类种子数据产品。",
                "能力均有代码、接口、消费者、指标、退出；产品Owner签字进入Q4版本。",
                "无Owner或无消费者的能力移出首版，不以“待定”占位。",
            ],
            [
                "数据质量",
                "六维规则、质量基线、监控看板、异常诊断与闭环工单。",
                "关键字段/链路有门槛；缺陷可发现、定位、修复、回归和回退。",
                "只展示分数、不能定位责任或触发处置时，退回治理基线重做。",
            ],
            [
                "复用",
                "一个自有Agent+一个产品应用；两个后端/中枢；统一NISM接口。",
                "第二消费者复用核心对象、质量规则和服务接口，主要工作是域适配。",
                "仍大量复制模型和管道时冻结扩面，先消除分叉。",
            ],
            [
                "价值",
                "接入周期、重复数据量、资源成本、查询P95/P99、缺陷周期基线。",
                "至少两项核心指标相对现状可量化改善，且不引入快环同步依赖。",
                "价值不可量化则不扩大资源，改善项和副作用带入阶段评审。",
            ],
            [
                "可信消费",
                "Agent身份、数据/模型版本、工具调用、结果和异常反馈证据。",
                "输入可追溯、越权可拦截、异常可回放；高风险动作仍由域控制器授权。",
                "无证据时保持只读/建议模式，不扩大自动动作权限。",
            ],
        ],
        col_w=[1.12, 3.43, 4.23, 3.77],
        font=8.05,
    )

    governance = [
        (
            TEAL,
            SOFT_TEAL,
            "治理节奏",
            "每周工作包、双周架构/Schema评审，第30/60/90天阶段门；产品Owner裁决范围，架构Owner裁决对象/接口，消费者Owner确认业务可用。",
        ),
        (
            NAVY,
            SOFT_NAVY,
            "关键依赖",
            "现有数据引擎与湖仓代码、NISM接口、两类数据Owner、Agent场景Owner、测试环境和现状成本基线必须在第30天前可用。",
        ),
        (
            CORAL,
            SOFT_CORAL,
            "停止线",
            "没有第二消费者、数据质量不能进入准入、收益不可度量、公共层成为快环单点或Q4版本无法承诺时，停止扩面并退回最小可信数据产品。",
        ),
    ]
    for i, (accent, fill, title, body) in enumerate(governance):
        x = 0.38 + i * 4.20
        box(
            slide,
            x,
            5.87,
            4.08,
            1.19,
            [
                {"text": title, "size": 11, "bold": True, "color": accent, "after": 3},
                {"text": body, "size": 9.1, "color": SLATE, "after": 0},
            ],
            fill=fill,
            border=accent,
            mt=0.07,
        )
    footer(slide, 5, "第8章 最高优先级AP", "[S59]")


def slide_blob(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return "\n".join(texts)


def count_cjk(prs: Presentation) -> list[int]:
    return [
        len(re.findall(r"[\u4e00-\u9fff]", slide_blob(slide)))
        for slide in prs.slides
    ]


def verify(prs: Presentation) -> None:
    if len(prs.slides) != TOTAL:
        raise SystemExit(f"expected {TOTAL} slides, got {len(prs.slides)}")

    textboxes = []
    overflow = []
    tables = []
    small_fonts = []
    for page, slide in enumerate(prs.slides, 1):
        table_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                textboxes.append((page, getattr(shape, "name", "?")))
            if shape.has_table:
                table_count += 1
            right = (shape.left + shape.width) / EMU
            bottom = (shape.top + shape.height) / EMU
            if (
                right > W_IN + 0.03
                or bottom > H_IN + 0.03
                or shape.left < Emu(-0.02 * EMU)
                or shape.top < Emu(-0.02 * EMU)
            ):
                overflow.append((page, getattr(shape, "name", "?"), round(right, 3), round(bottom, 3)))
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size.pt < 7.5:
                            small_fonts.append((page, run.text[:20], run.font.size.pt))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size.pt < 7.5:
                                    small_fonts.append((page, run.text[:20], run.font.size.pt))
        tables.append(table_count)

    if textboxes:
        raise SystemExit(f"standalone text boxes are not allowed: {textboxes}")
    if overflow:
        raise SystemExit(f"shapes overflow canvas: {overflow}")
    if small_fonts:
        raise SystemExit(f"font smaller than 7.5pt: {small_fonts}")
    if tables != [0, 0, 1, 1, 1]:
        raise SystemExit(f"expected tables [0,0,1,1,1], got {tables}")

    counts = count_cjk(prs)
    print("cjk chars:", counts, "total", sum(counts))
    for page, (got, need) in enumerate(zip(counts, CJK_MIN), 1):
        if got < need:
            raise SystemExit(f"P{page} CJK {got} < {need}")
    if sum(counts) < CJK_TOTAL_MIN:
        raise SystemExit(f"total CJK {sum(counts)} < {CJK_TOTAL_MIN}")

    blob = "\n".join(slide_blob(slide) for slide in prs.slides)
    required = [
        "智能数据中枢",
        "数据质量",
        "NISM",
        "2026 Q4",
        "RAN网管",
        "无线大数据",
        "近源",
        "动作证据",
        "Quality on Demand",
        "AP-P0",
    ]
    missing = [term for term in required if term not in blob]
    if missing:
        raise SystemExit(f"missing required terms: {missing}")

    forbidden = ["任务适用性", "Data Fitness for Use", "Quality of Data", "统一控制所有网元"]
    hits = [term for term in forbidden if term in blob]
    if hits:
        raise SystemExit(f"forbidden terms present: {hits}")
    ambiguous_qod = [
        line for line in blob.splitlines()
        if "QoD" in line and "Quality on Demand" not in line
    ]
    if ambiguous_qod:
        raise SystemExit(f"QoD must only mean Quality on Demand: {ambiguous_qod}")
    print("verify ok:", {"slides": TOTAL, "tables": tables, "textboxes": 0})


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_p1(prs)
    slide_p2(prs)
    slide_p3(prs)
    slide_p4(prs)
    slide_p5(prs)
    verify(prs)
    prs.save(OUT)
    print(f"wrote {OUT}")
    return OUT


if __name__ == "__main__":
    build()
