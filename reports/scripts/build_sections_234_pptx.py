#!/usr/bin/env python3
"""Generate a pyramid-style insight deck for report sections 2.3, 3 and 4."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parents[2]
ASSET = ROOT / "reports" / "data-fabric-assets"
OUT = ROOT / "reports" / "6g-data-fabric-sections-234-deck.pptx"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
LIME = RGBColor(0x2E, 0x8B, 0x57)
CORAL = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xB8, 0x86, 0x0B)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF4, 0xF7, 0xF8)
PALE = RGBColor(0xE8, 0xF0, 0xF2)
LINE = RGBColor(0xC5, 0xD0, 0xD4)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
FONT = "WenQuanYi Micro Hei"


def _set_run_east_asia(run) -> None:
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT)


def rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def lined(shape, fill: RGBColor, border: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)


def write_box(tf, lines, *, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT, spacing=1.05):
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
        _set_run_east_asia(run)


def add_text(slide, l, t, w, h, lines, **kw):
    box = slide.shapes.add_textbox(l, t, w, h)
    write_box(box.text_frame, lines, **kw)
    return box


def topbar(slide, chapter: str, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
    rgb(bar, TEAL)
    add_text(slide, Inches(0.38), Inches(0.14), Inches(12.5), Inches(0.22), [chapter], size=10, color=TEAL, bold=True)
    add_text(slide, Inches(0.38), Inches(0.32), Inches(12.5), Inches(0.36), [title], size=20, color=NAVY, bold=True)


def footer(slide, page: int, total: int, chapter: str) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.38), Inches(7.22), Inches(12.55), Inches(0.01))
    rgb(line, LINE)
    add_text(
        slide,
        Inches(0.38),
        Inches(7.24),
        Inches(9.2),
        Inches(0.22),
        [f"6G 数据架构 × 数据编织 · 行业深度洞察 · {chapter}"],
        size=9,
        color=MUTED,
    )
    add_text(slide, Inches(11.4), Inches(7.24), Inches(1.55), Inches(0.22), [f"{page} / {total}"], size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def thesis(slide, text: str, top=0.70, height=0.58) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(top), Inches(12.55), Inches(height))
    lined(box, PALE, TEAL)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    write_box(tf, [f"论点  {text}"], size=12, color=NAVY, bold=True, spacing=1.08)
    tf.paragraphs[0].space_before = Pt(4)


def card(slide, l, t, w, h, title, body, accent=TEAL, title_size=12, body_size=10):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    lined(shp, WHITE, LINE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.07), h)
    rgb(bar, accent)
    add_text(slide, l + Inches(0.16), t + Inches(0.06), w - Inches(0.24), Inches(0.28), [title], size=title_size, color=NAVY, bold=True)
    add_text(slide, l + Inches(0.16), t + Inches(0.34), w - Inches(0.24), h - Inches(0.40), body, size=body_size, color=SLATE, spacing=1.05)


def add_table(slide, l, t, w, h, rows, col_w=None, header=True, font=9):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), l, t, w, h)
    table = table_shape.table
    if col_w:
        for i, cw in enumerate(col_w):
            table.columns[i].width = cw
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.TOP
            fill = NAVY if header and r == 0 else (OFF if r % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font if r else font)
            run.font.bold = header and r == 0
            run.font.color.rgb = WHITE if header and r == 0 else INK
            run.font.name = FONT
            _set_run_east_asia(run)
            cell.text_frame.word_wrap = True
            cell.text_frame.margin_left = Inches(0.04)
            cell.text_frame.margin_right = Inches(0.04)
            cell.text_frame.margin_top = Inches(0.03)
            cell.text_frame.margin_bottom = Inches(0.03)
    return table


def ensure_png(svg_name: str, png_name: str) -> Path:
    png = ASSET / png_name
    svg = ASSET / svg_name
    if png.exists():
        return png
    if svg.exists():
        import cairosvg

        cairosvg.svg2png(url=str(svg), write_to=str(png), output_width=2200)
        return png
    raise FileNotFoundError(png_name)


def picture(slide, path: Path, l, t, w, h, caption: str | None = None):
    slide.shapes.add_picture(str(path), l, t, w, h)
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    frame.fill.background()
    frame.line.color.rgb = LINE
    frame.line.width = Pt(0.75)
    if caption:
        add_text(slide, l, t + h + Inches(0.02), w, Inches(0.22), [caption], size=8, color=MUTED)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    e_north = ensure_png("ericsson-ai-ready-data-management-2026.svg", "ericsson-ai-ready-data-management-2026.png")
    e_fig18 = ensure_png("ericsson-data-ingestion-architecture-2026.svg", "ericsson-data-ingestion-architecture-2026.png")
    huawei = ASSET / "huawei-6g-data-plane-architecture.jpg"
    nokia = ASSET / "nokia-autonomous-network-fabric.png"

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    total = 16

    # ------------------------------------------------------------------ 1 cover
    s = new_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rgb(bg, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), H)
    rgb(accent, TEAL)
    add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.35), ["6G 数据架构 × 数据编织  ·  行业深度洞察胶片"], size=14, color=GOLD, bold=True)
    add_text(s, Inches(0.7), Inches(1.65), Inches(12), Inches(1.4), ["从机制事实到价值链趋势：", "2.3 / 3 / 4 章金字塔推演"], size=32, color=WHITE, bold=True, spacing=1.15)
    add_text(
        s,
        Inches(0.7),
        Inches(3.35),
        Inches(11.5),
        Inches(1.3),
        [
            "本册不把 Data Fabric 写成既定终点。2.3 先用五类可证伪判断核对公开机制；第3章只从四条约束推出可验证趋势；",
            "第4章把趋势压回 3 年可验证事实与 5 年控制点，不再另造预测清单。",
            "评委不另做口头交接：每页先给论点，再给证据、边界和停止线。",
        ],
        size=14,
        color=RGBColor(0xD5, 0xDE, 0xE5),
        spacing=1.2,
    )
    for i, (title, body) in enumerate(
        [
            ("2.3  事实底座", "七个子节各一页：规则、标准、横向比较、三家架构竞情、相邻竞合。"),
            ("第3章  价值链", "先给整体推理，再按状态可知→数据适用→行动可信→价值可结算编排。"),
            ("第4章  时间窗", "一页收束：3年看透不变量，5年看清控制点。"),
        ]
    ):
        x = Inches(0.7 + i * 4.1)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.0), Inches(3.9), Inches(1.55))
        lined(box, RGBColor(0x12, 0x2B, 0x4A), TEAL)
        add_text(s, x + Inches(0.18), Inches(5.12), Inches(3.55), Inches(0.32), [title], size=14, color=GOLD, bold=True)
        add_text(s, x + Inches(0.18), Inches(5.48), Inches(3.55), Inches(0.9), [body], size=12, color=WHITE, spacing=1.1)
    add_text(s, Inches(0.7), Inches(6.75), Inches(11.5), Inches(0.3), ["依据报告正文 · 金字塔结构：先结论，后机制，再边界"], size=11, color=RGBColor(0x9A, 0xB0, 0xC0))

    # ------------------------------------------------------------------ 2 pyramid
    s = new_slide(prs)
    topbar(s, "本册结构  ·  先读这一页", "金字塔阅读顺序：事实底座支撑价值链，价值链支撑时间窗")
    thesis(s, "先核对公开机制，再推出可验证趋势，最后只投射时间窗。厂商案例是外部校验，不是趋势来源。")
    layers = [
        (0.38, 1.45, 12.55, 1.15, GOLD, "第4章  时间窗", "只投射第3章 T1–T8。3年看透不变量能否被合同、接口和证据固定；5年看清谁拥有语义/任务适用性、近源运行时、动作证据和跨网结算控制点。不再另造趋势。"),
        (0.90, 2.75, 11.50, 1.55, TEAL, "第3章  价值链推演", "不可变事实 → 矛盾 → 必须补齐的能力 → 可验证趋势。四条约束对应四个关口：物理分布→状态可知（T2/T3）；语义依赖→数据适用（T1/T6）；动作外部性→行动可信（T4/T5）；交易协同→价值可结算（T7）。T8 横切全部关口。"),
        (1.45, 4.45, 10.40, 2.45, NAVY, "2.3  机制事实底座", "零假设：SA2/SA5/O-RAN 与网络原生功能可能已足够。五类判断：网络原生解决 / 机制重合 / 需要补充 / Data Fabric 不适用 / 公开证据不足。2.3.1 定规则，2.3.2 定标准交界面，2.3.3 横向比较机制而非排名，2.3.4–2.3.6 把架构图放进三家竞情，2.3.7 划相邻竞合边界。"),
    ]
    for l, t, w, h, color, title, body in layers:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        rgb(box, color)
        add_text(s, Inches(l + 0.22), Inches(t + 0.10), Inches(w - 0.4), Inches(0.32), [title], size=16, color=WHITE, bold=True)
        add_text(s, Inches(l + 0.22), Inches(t + 0.44), Inches(w - 0.4), Inches(h - 0.55), [body], size=12, color=WHITE, spacing=1.12)
    footer(s, 2, total, "本册结构")

    # ------------------------------------------------------------------ 3  2.3.1
    s = new_slide(prs)
    topbar(s, "2.3.1  分析规则", "先看架构原本解决什么，再用 2.2 提问，只给五类可证伪判断")
    thesis(s, "2.3 不预设 6G 必须采用独立 Data Fabric。未描述的机制记为未知，不是缺失；名称相似不等于能力已实现。")
    card(
        s,
        Inches(0.38),
        Inches(1.42),
        Inches(4.05),
        Inches(2.55),
        "本节目的",
        [
            "2.2 已建立能力框架。2.3 只回答：公开架构、标准和产品实际提供了哪些机制，这些机制覆盖、部分覆盖还是替代了哪些能力。",
            "厂商叙述只作外部校验，不把其路线图写成行业必然。",
        ],
        NAVY,
        body_size=11,
    )
    card(
        s,
        Inches(4.58),
        Inches(1.42),
        Inches(4.05),
        Inches(2.55),
        "零假设",
        [
            "SA2、SA5、O-RAN 与网络原生功能可能已经足够。",
            "外部编织层不得进入毫秒级快环同步依赖。",
            "只有当公开机制无法覆盖那些能力时，才判断“需要补充”。",
        ],
        TEAL,
        body_size=11,
    )
    card(
        s,
        Inches(8.78),
        Inches(1.42),
        Inches(4.15),
        Inches(2.55),
        "禁止写法",
        [
            "不把 Mesh / Fabric / Data Plane 名称直接写成 2.2 能力已落地。",
            "不把研究架构写成商用能力。",
            "不把单域案例外推成统一跨域编织。",
        ],
        CORAL,
        body_size=11,
    )
    add_text(s, Inches(0.38), Inches(4.08), Inches(12.5), Inches(0.28), ["五类判断：每条后文结论必须落入其中一类，并可被后续公开材料证伪"], size=12, color=NAVY, bold=True)
    labels = [
        (LIME, "网络原生解决", "现有网元、接口或管理功能已覆盖该能力，不必再引入外部编织层。"),
        (TEAL, "机制重合", "公开架构已提供同类机制，但范围、接口或证据链仍不完整。"),
        (AMBER, "需要补充", "能力缺口已被公开材料承认，或现有机制明确无法覆盖。"),
        (CORAL, "Data Fabric 不适用", "该问题应由网络原生、标准接口或业务系统承担，不应交给编织层。"),
        (SLATE, "公开证据不足", "材料未描述该机制。记为未知，不得写成缺失或落后。"),
    ]
    for i, (color, title, body) in enumerate(labels):
        x = Inches(0.38 + i * 2.55)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.42), Inches(2.45), Inches(2.55))
        lined(box, WHITE, LINE)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4.42), Inches(2.45), Inches(0.10))
        rgb(head, color)
        add_text(s, x + Inches(0.10), Inches(4.60), Inches(2.25), Inches(0.55), [title], size=12, color=NAVY, bold=True)
        add_text(s, x + Inches(0.10), Inches(5.18), Inches(2.25), Inches(1.65), [body], size=11, color=SLATE, spacing=1.08)
    footer(s, 3, total, "2.3.1")

    # ------------------------------------------------------------------ 4  2.3.2
    s = new_slide(prs)
    topbar(s, "2.3.2  标准与接口落位", "标准分别解决局部问题，没有共同命名为 Data Fabric")
    thesis(s, "待验证的是共同外部行为，不是再加统一名称。更有把握冻结角色、稳定标识、最小元数据、交付外部行为，并与任务适用性/授权/证据挂钩。Kafka、湖仓、完整 KG、独立 Fabric 或专用 Data Plane 都仍是未收敛实现选择。")
    add_table(
        s,
        Inches(0.38),
        Inches(1.38),
        Inches(12.55),
        Inches(2.42),
        [
            ["标准域 / 状态", "实际定义或研究对象", "与2.2能力的关系", "不能据此推出"],
            ["ITU-R IMT-2030  Recommendation", "六类场景、总体能力与2030目标；不定义网元、协议或数据生命周期", "形成AI原生、通感、孪生和泛在智能的数据需求来源", "不能推出任务适用性、Catalog、血缘或“6G Data Fabric”接口已成为IMT要求  [S1]"],
            ["3GPP RAN  6G Study", "无线测量、RRC/MAC/PDCP/RRM、MDT；亚毫秒快环至秒级报告", "网络原生定义小区、波束、UE、频段、测量与配置等权威事实", "不能推出RAN还应承担跨域目录、企业湖仓或统一产品运营"],
            ["SA2 KI#21  Rel-20 Draft TR", "6GC系统级发现、注册、采集、处理、存储与暴露，及UE/RAN/OAM/AF关系", "可能原生覆盖角色、发现、交付控制、最小元数据和授权挂钩点", "TR 23.801-01尚未冻结独立数据面、必选NF或具体实现拓扑  [S5]"],
            ["SA5 TS积木 + DMFW Draft", "PM/CM/FM/Trace、MDA、AI/ML、NDT管理及管理侧数据框架", "原生提供管理对象、分析、模型/孪生生命周期、同步与部分质量语义", "管理能力≠毫秒数据总线、NDT保真度或动作安全已被证明  [S6][S26][S29]"],
            ["O-RAN 开放接口 / ZSM 029", "O1/O2/A1/E2/R1；ZSM研究数据注册/发现、资产、认证、采集与编排", "按时标开放数据、策略、模型和动作；ZSM与目录/资产机制明显重合", "五类接口不是统一语义/任务适用性/血缘总线；工作项≠最终规范或独立Fabric  [S7][S43]"],
            ["TMF / GSMA / CAMARA", "自治运营、服务包装、开发者消费、用量、认证与结算", "覆盖数据/能力产品的外部行为和商业运营接口", "不能推出内部采集、任务适用性生成或无线动作路径已经统一  [S15][S16]"],
        ],
        col_w=[Inches(2.35), Inches(3.45), Inches(3.45), Inches(3.30)],
        font=8,
    )
    add_table(
        s,
        Inches(0.38),
        Inches(3.88),
        Inches(12.55),
        Inches(1.88),
        [
            ["交界面", "现有分工", "可观察问题", "可能解法（不预设实现）"],
            ["SA2系统数据 ↔ SA5管理数据", "SA2研究系统级发现/交付；SA5管PM/CM/FM/Trace、MDA与NDT", "同一OAM数据可能重复注册，使用不同质量或授权表达", "共享稳定ID、最小元数据和质量语义；可由标准功能、适配层或外部目录实现"],
            ["RAN/O-RAN ↔ Core数据框架", "E2/O1提供无线与管理事实；SA2候选路径负责系统交付和暴露", "低时延直达与治理可见性冲突，跨接口对象难关联", "事件时间、配置版本、关联ID和异步证据回流；快环继续使用本地确定性状态"],
            ["O1/O2/A1/E2/R1之间", "接口对象、Schema、作用域和时标不同", "同一小区、策略和模型在不同接口难以追踪", "版本化映射、权威源和血缘引用；不要求一个接口或平台替代全部接口"],
            ["AI/ML/NDT ↔ 动作放行", "SA5管理生命周期，O-RAN或厂商控制器执行动作", "“可管理”容易被误读为“可证明安全”", "输入快照、模型/策略/动作版本、同步误差、灰度和回滚证据挂钩"],
            ["内部能力 ↔ CAMARA/Operate API", "内部网络能力经标准API向开发者暴露", "API一致不代表语义、任务适用性、责任、用途和计量一致", "内部契约到外部SLA的映射、跨网一致性测试和清结算字段"],
        ],
        col_w=[Inches(2.55), Inches(3.35), Inches(3.25), Inches(3.40)],
        font=8,
    )
    for i, (title, body, color) in enumerate(
        [
            ("形态A · 增强既有功能", "DCCF/ADRF/NWDAF/MDA扩展。若既有SBI与管理服务能满足发现、交付和生命周期，网络原生增强可能比新增平台成本更低。", TEAL),
            ("形态B · 新功能/新平面", "若隔离、吞吐、Fan-out或近源处理无法满足，可引入专用控制/处理/存储/分发；代价是会话、安全、移动性、容灾与状态一致性。", AMBER),
            ("形态C · 混合能力包", "能力标准化、拓扑可组合。事件/流/文件/模型可选不同后端。当前更有把握：先冻结外部行为与测试，而不是名称。", NAVY),
        ]
    ):
        card(s, Inches(0.38 + i * 4.2), Inches(5.84), Inches(4.05), Inches(1.15), title, [body], color, title_size=11, body_size=9)
    footer(s, 4, total, "2.3.2")

    # ------------------------------------------------------------------ 5  2.3.3
    s = new_slide(prs)
    topbar(s, "2.3.3  三家机制横向比较", "按问题比机制，不是按谁更接近 Data Fabric 排名")
    thesis(s, "三家是不同起点：Ericsson 联邦+产品组合，Huawei 网络原生拓扑，Nokia 数据管理与自治软件商品化。登记对象不同，不能按功能表做 Fabric 成熟度排名。")
    add_table(
        s,
        Inches(0.38),
        Inches(1.38),
        Inches(12.55),
        Inches(4.85),
        [
            ["6G数据面问题", "Ericsson公开机制", "Huawei公开机制", "Nokia公开机制", "客观判断"],
            [
                "发现与全局视图",
                "2021 GDC；2026参考架构提出统一目录体验与联邦数据岛",
                "DO维护DA能力注册并按服务目标构造数据拓扑",
                "Data Suite提供Catalog、多厂商适配和3GPP/TMF语义",
                "登记对象不同：数据资产、执行能力与产品目录不能直接等同；均需验证跨域稳定ID",
            ],
            [
                "摄取、近源处理与分发",
                "EDCA/DRG/DDC、Mediation、DataOps及智能数据移动；Figure 18仍保留EDCA/DRG",
                "DA/DPF/DSF近源算子，DCP以Pub/Sub支持多消费者",
                "域湖仓保留自治，Common Fabric/Lakehouse联邦访问",
                "这些首先属于数据平面或集成运行时；有无Data Fabric不决定吞吐、时延和容灾",
            ],
            [
                "语义、任务适用性与血缘",
                "2026参考架构描述质量、血缘、语义模型与知识图谱；统一产品部署证据不足",
                "研究输入含Data QoS，AUTINOps有语义底座；DO/DA/DCP统一血缘未公开",
                "Data Suite显式提供质量/validity、血缘、语义和产品能力",
                "机制重合程度不同；参考架构、标准输入和已发布产品必须分档，不能按功能表排名",
            ],
            [
                "策略与动作证据",
                "EIAP/rApps、NWDAF及AI agent架构消费数据并执行受限功能",
                "AUTINOps、NDT sandbox与AgenticCore分别位于管理、研究和Core愿景层",
                "AN Fabric Sense–Think–Act经AN Apps调用域控制器",
                "它们是数据消费者与动作系统，不等于数据管理层；共享快照和回滚证据仍需单证",
            ],
            [
                "数据产品与外部运营",
                "2026参考架构提出数据产品/Marketplace；CCES/Aduna覆盖能力渠道",
                "Open Gateway覆盖能力API；内部跨域数据产品契约公开较少",
                "Data Suite、Fabric Experience、Monetization与Marketplace形成产品叙事",
                "数据产品、网络能力API和渠道聚合是三个层次；收入证据不能反推内部统一Fabric",
            ],
            [
                "6G演进接口",
                "Common data plane与model orchestrator属于6G愿景，与2026参考架构无公开产品继承",
                "DO/DA/DCP是候选网络原生架构并有DCP原型",
                "Holistic data framework与现有Data Suite/AN Fabric做能力映射",
                "三条路线可能互补、替代或按时标共存；目前没有规范性结论支持单一路线",
            ],
        ],
        col_w=[Inches(1.85), Inches(2.7), Inches(2.65), Inches(2.65), Inches(2.70)],
        font=8,
    )
    add_text(
        s,
        Inches(0.38),
        Inches(6.30),
        Inches(12.55),
        Inches(0.70),
        [
            "比较结论：Ericsson更强调分散产品与数据岛的联邦，Huawei更强调网络原生拓扑和近源执行，Nokia更强调数据管理与自治软件的产品化。这是架构起点和机制选择，不是“谁更接近Data Fabric”的排名。",
            "后三页把各自官方架构图放进竞情：先看图里实际有哪些对象，再判断机制重合、网络原生替代、必要补充或证据空白。",
        ],
        size=11,
        color=SLATE,
        spacing=1.08,
    )
    footer(s, 5, total, "2.3.3")

    # ------------------------------------------------------------------ 6 Ericsson
    s = new_slide(prs)
    topbar(s, "2.3.4  Ericsson 竞情", "三条线必须分开：架构继承、产品演进、未来架构")
    thesis(s, "2026 Figure 18 仍保留 EDCA/DRG；官方产品血缘是 Mediation→Telco DataOps，不是 EDCA→DataOps；AI-ready 是 North Star，不是已交付统一 SKU。“no data copy”应读成减少不必要复制，而非绝对不移动数据。")
    picture(s, e_north, Inches(0.38), Inches(1.38), Inches(6.10), Inches(1.88), "图4a  Ericsson 2026 AI-ready North Star（白皮书p.16）：源/摄取/精炼治理/消费支撑 + 目录/质量/血缘/产品/KG")
    picture(s, e_fig18, Inches(0.38), Inches(3.52), Inches(6.10), Inches(1.58), "图4b  Figure 18 摄取架构：EDCA、DRG仍在；客户侧DMF + Ericsson侧EFDL/Data Fabric 构成新联邦边界")
    card(
        s,
        Inches(6.63),
        Inches(1.38),
        Inches(6.30),
        Inches(1.18),
        "架构继承线  ·  2021 A → 2026 A",
        ["2021 EDCA/DRG/DDC/GDC 解决“采集一次、跨域中继、联邦发现”。2026 Figure 18仍保留EDCA/DRG，并扩展DMF、EFDL、EFDL Data Fabric。DDC/GDC名称消失，公开材料未说明被替代或退役，只能判断功能被重新组织。"],
        TEAL,
        body_size=9,
    )
    card(
        s,
        Inches(6.63),
        Inches(2.62),
        Inches(6.30),
        Inches(1.12),
        "产品演进线  ·  Mediation → Telco DataOps",
        ["2025 solution brief给出的官方血缘是Mediation→Telco DataOps，增加管道、管理、治理和应用加速。EIAP（Swisscom合同）、NWDAF、CCES/Aduna分别有锚点，但不能画成同一版本链或统一元数据控制面。"],
        NAVY,
        body_size=9,
    )
    card(
        s,
        Inches(6.63),
        Inches(3.80),
        Inches(6.30),
        Inches(1.30),
        "未来架构线  ·  North Star ≠ 6G愿景",
        ["2026-01 AI-ready把范围扩到语义/KG、质量、血缘、数据产品和Agent消费，是future-state参考架构。6G common data plane + model orchestrator更靠近网络数据采集、传输、潜在存储、治理与模型放置。二者有机制交集，无公开产品继承。"],
        AMBER,
        body_size=9,
    )
    add_table(
        s,
        Inches(0.38),
        Inches(5.22),
        Inches(12.55),
        Inches(1.78),
        [
            ["已证明 / 架构事实", "尚未证明", "反证门 / 引用"],
            [
                "摄取加工、EIAP、NWDAF、API渠道分别有产品或部署锚点。Figure 18延续EDCA/DRG；North Star另行加入语义、治理、产品和Agent支撑。",
                "参考架构已成为统一SKU；DataOps/EIAP/NWDAF共享同一跨域语义、任务适用性和证据链；第三方Agent、异构湖仓和跨网API在相同契约下可双向迁移。",
                "若统一平台在多厂商网络以更低TCO替代这些组件，联邦组合的经济优势会下降。S9/S10/S11/S18/S21/S30/S35/S36/S57/S58",
            ],
        ],
        col_w=[Inches(4.2), Inches(4.35), Inches(4.00)],
        font=9,
    )
    footer(s, 6, total, "2.3.4")

    # ------------------------------------------------------------------ 7 Huawei
    s = new_slide(prs)
    topbar(s, "2.3.5  Huawei 竞情", "DO / DA / DCP 首先是 6G 数据面问题，不是 Data Fabric 宣言")
    thesis(s, "DO/DA/DCP 把多源—多处理—多消费者建模为网络内可编排拓扑，可作部分通用集成的网络原生替代；但仍是研究架构。必须同时保留“网络原生吸收”和“外部元数据补充”两种解释。")
    picture(s, huawei, Inches(0.38), Inches(1.38), Inches(5.95), Inches(2.72), "图5  Huawei 2024 6G Data Plane：DO编排拓扑，DA/DPF近源处理，DCP/Data Spine以Pub/Sub解耦。研究架构，不是3GPP既定方案")
    card(
        s,
        Inches(6.48),
        Inches(1.38),
        Inches(6.45),
        Inches(1.32),
        "先按原始设计读控制流 / 数据流",
        [
            "控制流：服务请求向DO声明采集/过滤/聚合/转换/分析目标；DO查询DA能力注册，分解操作链，选择UE/基站/边缘/Core上的DA/DPF/DSF，再下发拓扑、访问、隐私和资源策略。",
            "数据流：源DA发布topic，DCP按订阅或带内拓扑异步分发；上游模型输出可继续作为下游输入。这首先是6G数据面问题，不是Data Fabric问题。",
        ],
        NAVY,
        body_size=9,
    )
    card(
        s,
        Inches(6.48),
        Inches(2.76),
        Inches(6.45),
        Inches(1.34),
        "两种解释并行  ·  不把“补成Fabric”当唯一终点",
        [
            "原生解释：DO能力注册可吸收部分目录/工作流；DA/DPF保持近源执行；DCP针对Fan-out和低时延优化。",
            "补充解释：跨域资产语义、Owner、产品契约、持久化历史仍可能由外部控制层或湖仓承担。结果A入标 / B厂商实现 / C分层混合都可能；停止线是增强SBA/UP已够用则不新增平面。",
        ],
        TEAL,
        body_size=9,
    )
    add_table(
        s,
        Inches(0.38),
        Inches(4.22),
        Inches(12.55),
        Inches(2.75),
        [
            ["组件 / 证据层", "原生职责", "与Fabric参照系", "证据边界"],
            ["DO  研究原型P", "DA能力注册、目标分解、操作DAG、选点并下发拓扑/策略", "发现/策略/编排与部分Fabric机制重合；登记对象主要是执行能力，≠通用资产Catalog", "与SA2/SA5/管理职责边界、DO容灾和多厂商控制仍未冻结"],
            ["DA / DPF / DSF", "近源采集、过滤、缓存、推理；DPF承载算子，DSF承载持久存储", "网络原生近源运行时；Fabric不应替代其毫秒执行和断链状态", "资源预算、升级回退、可信执行、移动性和多厂商可移植性未证"],
            ["DCP 约2ms；10 broker/128KB约2394.6MB/s", "消息代理、队列和Pub/Sub；Stateful查表，Stateless把拓扑编码进包头", "可能替代通用MQ的部分交付功能；不负责证明语义、血缘或数据产品", "对照开启持久化而DCP移除，不公平；缺跨站容灾、租户、移动性和长期稳定性"],
            ["AUTINOps 管理域C2", "DTN、语义底座和Agent执行Non-RT跨域运维闭环", "管理域消费者/动作系统，可提供部分语义和证据", "部署案例不含DO/DA/DCP级网络数据面，不是DCP商用证明  [S12][S37][S38]"],
        ],
        col_w=[Inches(2.55), Inches(3.35), Inches(3.40), Inches(3.25)],
        font=8,
    )
    footer(s, 7, total, "2.3.5")

    # ------------------------------------------------------------------ 8 Nokia
    s = new_slide(prs)
    topbar(s, "2.3.6  Nokia 竞情", "Data Suite 与 L4/L5 机制重合最显式；AN Fabric 是智能消费层")
    thesis(s, "两套“Fabric”不是同一层。公开案例分别证明域级AI、处理规模或单项产品，不能合并外推为一个统一实例已跨RAN/Core/传输和多厂商生产部署。")
    picture(s, nokia, Inches(0.38), Inches(1.38), Inches(5.95), Inches(2.55), "图6  Nokia 2025 AN Fabric：Sense–Think–Act智能循环经AN Apps调用域控制器，再向Monetization/Marketplace暴露。智能消费层，不是数据管理Fabric")
    card(
        s,
        Inches(6.48),
        Inches(1.38),
        Inches(6.45),
        Inches(2.55),
        "先区分数据管理、智能消费与网络执行",
        [
            "数据流：各域本地采集形成domain lakehouse；common fabric/lakehouse经适配与联邦向Data Suite供给产品，再进入AN Fabric的Sense层。数据默认留域，公共层不要求物理集中。",
            "控制流：Sense统一观察和chain-of-custody，Think用telco-trained模型/KG/Glass-Box，Act经域AN Apps调用现有OSS/RIC/SON。",
            "Data Suite与L4/L5机制重合最显式；AN Fabric是智能消费层。2024蓝图承认CSP真实分布式部署仍有限。公共层—Data Suite—AN Fabric是否共享同一元数据/策略/版本，公开材料尚不足以回答。",
        ],
        TEAL,
        body_size=9,
    )
    add_table(
        s,
        Inches(0.38),
        Inches(4.08),
        Inches(12.55),
        Inches(1.48),
        [
            ["案例", "公开规模/结果", "能够证明", "不能证明"],
            ["KDDI AVA PDDR", "覆盖20万4G/5G RAN，AI检测silent cell并联动恢复", "RAN域异常检测和闭环恢复可在商用网运行", "AN Fabric/Data Suite跨RAN/Core/传输统一部署或多厂商中立"],
            ["stc Cognitive SON", "高流量期>1万自动动作；负载小区利用率约+30%、用户吞吐约+10%", "域级RAN自治优化与自动动作产生可量化结果", "这些动作由统一跨域Fabric驱动，或Core/传输共享同一产品契约"],
            ["Indosat / 匿名Tier-1", "节能扩至多区域Nokia RAN；25万小区、27万报告/秒", "域应用可SaaS化；Data Suite具备大规模摄取和产品化潜力", "未具名、无独立审计，不能外推广泛跨域商用或真正多厂商"],
        ],
        col_w=[Inches(2.15), Inches(3.55), Inches(3.45), Inches(3.40)],
        font=8,
    )
    qs = [
        ("同一控制体系？", "Data Suite/AN Fabric/AN Apps分层清晰，但AVA/MantaRay/节能等SKU可能只在品牌层组合。验证门：统一版本、对象模型和跨产品复用。"),
        ("真正多厂商？", "3GPP/TMF语义和多厂商适配进入产品说明；Nokia Schema/规则仍可能软锁定。验证门：第三方后端/Agent、公开契约和双向迁移。"),
        ("贯通动作？", "Sense–Think–Act连接分析与应用，但高风险动作仍依赖域控制器。验证门：输入快照、动作证据、灰度和第三方执行器。"),
        ("可替换后端？", "本地/GCP/混合提高选择空间，也可能带来云和元数据双重锁定。验证门：核心元数据可导出、多云TCO、主权部署和故障迁移。"),
    ]
    for i, (title, body) in enumerate(qs):
        x = Inches(0.38 + i * 3.2)
        card(s, x, Inches(5.64), Inches(3.05), Inches(1.35), title, [body], CORAL if i in (0, 2) else AMBER, title_size=11, body_size=8)
    footer(s, 8, total, "2.3.6")

    # ------------------------------------------------------------------ 9  2.3.7
    s = new_slide(prs)
    topbar(s, "2.3.7  相邻竞合者", "云、湖仓、聚合商可提供部分机制，但不能承接 RAN 快环与网元动作权限")
    thesis(s, "云和湖仓可以提供模型与数据运行时，设备商更接近网络语义、近源执行和动作权限，聚合商掌握开发者入口与结算。价值在于三层能否通过开放Schema、可导出证据、可替换后端和SLA解耦，而不是其中任何一层必须成为统一平台。")
    add_table(
        s,
        Inches(0.38),
        Inches(1.42),
        Inches(12.55),
        Inches(3.15),
        [
            ["玩家层", "可提供的机制", "相对设备商优势", "不应承担的网络职责", "客观竞合关系"],
            [
                "Microsoft / AWS / Google  云与Agent运行时",
                "IaaS、湖仓计算、模型训练/推理、Agent工具和多云运维",
                "算力弹性、模型生态、工程工具及企业数据入口",
                "不应进入RAN快环，或越过3GPP/O-RAN权限直接控制网元",
                "可作为数据/模型后端，也可能提供外部Fabric控制能力；取决于主权、成本和开放接口。Google Telecom Data Fabric仍为Private Preview  [S25]",
            ],
            [
                "Databricks / Snowflake  湖仓/Catalog/MLOps",
                "开放表格式、流批、Catalog、共享、特征/MLOps与跨企业协作",
                "通用数据工程成熟、客户存量大、后端规模效应强",
                "不应定义无线对象、Near-RT资源预算和网络动作安全",
                "适合作为可替换后端；设备商无需重复建设全栈湖仓，但需保留电信语义、任务适用性和策略接口的控制或可移植性",
            ],
            [
                "Aduna / CPaaS / Marketplace  API分发与计量",
                "跨运营商目录、认证、开发者触达、用量、计费和收入分成",
                "覆盖聚合与渠道网络效应，可把单网API变成规模产品",
                "不应反向定义内部数据采集、任务适用性生成和动作执行实现",
                "验证外部需求与结算，不证明内部Fabric；可倒逼上游形成稳定语义、SLA和责任映射",
            ],
        ],
        col_w=[Inches(2.35), Inches(2.55), Inches(2.35), Inches(2.55), Inches(2.75)],
        font=8,
    )
    card(
        s,
        Inches(0.38),
        Inches(4.70),
        Inches(6.15),
        Inches(2.25),
        "竞合边界",
        [
            "谁能把网络对象变成可运营产品，谁就与湖仓/渠道形成重叠。",
            "谁若把快环状态同步交给远端平台，谁就越过停止线。",
            "相邻竞合的关键不是“谁更像Fabric”，而是责任切分是否被合同和接口写清。",
            "Data Fabric视角的检查项：开放Schema、可导出证据、可替换后端、SLA解耦。",
        ],
        TEAL,
        body_size=11,
    )
    card(
        s,
        Inches(6.68),
        Inches(4.70),
        Inches(6.25),
        Inches(2.25),
        "对后文趋势的硬约束",
        [
            "T2/T3必须坚持全局控制、近源执行；快环不得外包给云目录。",
            "T4/T5的动作权限不能交给聚合商或行业平台。",
            "T7的结算可以发生在渠道和行业平台，但L1生产责任仍在网络侧。",
            "T8应冻结外部行为与测试，而不是把某一云SKU写成标准网元。",
        ],
        NAVY,
        body_size=11,
    )
    footer(s, 9, total, "2.3.7")

    # ------------------------------------------------------------------ 10  2.4
    s = new_slide(prs)
    topbar(s, "2.4  行业小结", "共同问题已经明确，Data Fabric 只是可能解法之一")
    thesis(s, "三家是不同起点，不是成熟度序列。目前没有端到端胜者；能被证明的是局部机制，不是统一编织已经成立。")
    points = [
        (TEAL, "共同问题已经出现", "异构发现、近源处理、跨时标交付、语义/任务适用性、动作证据和外部运营均有真实需求，区别在实现责任与时标，不在名称。"),
        (NAVY, "标准没有预设平台形态", "SA2/SA5仍在Rel-20研究；增强既有功能、新NF和混合能力包都可能。独立Fabric或专用Data Plane均非既定结论。"),
        (AMBER, "三家是不同起点，不是成熟度序列", "Ericsson强调联邦与产品组合，Huawei强调网络原生拓扑，Nokia强调数据管理和自治软件商品化；三者可替代也可互补。"),
        (CORAL, "证据必须按层拆开", "摄取、湖仓、流平台和域应用已有商用；统一语义、跨产品契约和高风险动作证据不能由单点案例外推。"),
        (SLATE, "Data Fabric有明确适用边界", "它适合跨域发现、治理和产品供给；近源毫秒执行、移动性和无线控制仍应由网络原生机制承担。"),
        (LIME, "尚无端到端胜者", "没有中立证据证明任何一家已在多运营商、多厂商环境同时实现公共语义、通信数据适用性、开放后端和可信动作闭环。"),
    ]
    for i, (color, title, body) in enumerate(points):
        x = Inches(0.38 + (i % 3) * 4.2)
        y = Inches(1.40 + (i // 3) * 1.48)
        card(s, x, y, Inches(4.05), Inches(1.38), title, [body], color, title_size=12, body_size=10)
    card(
        s,
        Inches(0.38),
        Inches(4.40),
        Inches(12.55),
        Inches(2.55),
        "转入第3章的接力规则  ·  最稳健断言：Data Fabric是有价值的观察语言，不是6G数据面的预设答案",
        [
            "2.3只提供事实约束，不提供趋势清单。第3章不得把“厂商正在做X”改写成“行业必然出现X”。",
            "推导顺序只能是：不可变事实 → 若不解决会产生的矛盾 → 必须补齐的能力 → 可验证趋势。标准与厂商仅作外部校验。",
            "若某趋势只能靠某家路线图成立，则该趋势降级为观察项，不得进入第4章时间窗。",
            "应让网络原生方案、外部编织方案和分时标混合方案，在同一组时延、可靠性、语义、任务适用性、证据、开放性与TCO指标下竞争。",
        ],
        NAVY,
        body_size=12,
    )
    footer(s, 10, total, "2.4")

    # ------------------------------------------------------------------ 11  ch3 logic
    s = new_slide(prs)
    topbar(s, "第3章  整体推理逻辑", "四条约束打开四个关口；趋势按价值链排列，而不是按口号排列")
    thesis(s, "先让状态可知，再让数据适用，再让行动可信，最后让价值可结算。T8 横切全部关口，不单列第五条主线。")
    add_table(
        s,
        Inches(0.38),
        Inches(1.40),
        Inches(12.55),
        Inches(2.35),
        [
            ["不可变约束", "若不解决的矛盾", "因此必须补齐", "必要关口 / 趋势"],
            ["物理分布：RAN/Core/Edge/ISAC持续产生高频数据；全量集中受带宽、时延、能耗、隐私和断链约束", "中心既看不到及时变化，又承受不可接受的搬运成本；本地执行与全局治理相互脱节", "在源附近处理，同时把结构、配置、质量、权限和血缘变化同步为可理解状态", "状态可知  ·  T2 变化事件  ·  T3 全局控制/近源执行"],
            ["语义依赖：无线数据含义依赖对象、时空、配置、采样和测量误差；能访问≠适合当前任务", "模型和消费者只能重复解释、清洗与对齐；字段完整仍可能产生错误比较、训练偏差和闭环误判", "把数据封装为带消费者、Owner、契约和生命周期的产品，并以面向用途的任务适用性证明适用边界", "数据适用  ·  T1 可运营产品  ·  T6 任务适用性信封"],
            ["动作外部性：AI/Agent输出会改变真实网络状态；局部错误可影响相邻域、业务SLA和其他闭环", "数据错误会直接升级为动作风险；仅有模型置信度无法回答谁授权、为何放行、失败如何接管", "建立供数—决策—执行—回证责任链，并用NDT反事实预演、策略、灰度和回滚约束高风险动作", "行动可信  ·  T4 Agent协同闭环  ·  T5 NDT动作关口"],
            ["交易协同：行业应用需要跨运营商覆盖、一致行为、用途约束、责任界定、计量与清结算", "API数量增加但体验不一致、无人持续付费；内部成本、外部SLA、责任与收入无法闭环", "按责任边界区分可组合能力与结果承诺，分别定义控制权、合同、计价、风险和验收", "价值可结算  ·  T7 责任边界双路径"],
        ],
        col_w=[Inches(3.3), Inches(3.7), Inches(2.1), Inches(3.45)],
        font=10,
    )
    add_text(s, Inches(0.38), Inches(3.90), Inches(12.5), Inches(0.28), ["价值链阅读顺序  ·  后一关口以前一关口成立为前提"], size=12, color=NAVY, bold=True)
    stages = [
        ("1", "状态可知", "先知道对象在哪、如何变化、谁能就近执行"),
        ("2", "数据适用", "再知道它能被谁、以何种质量、用于何种任务"),
        ("3", "行动可信", "然后才允许决策改网，并留下反证与回证"),
        ("4", "价值可结算", "最后才把内部可信生产接出标准服务与行业结果"),
    ]
    for i, (num, title, body) in enumerate(stages):
        x = Inches(0.38 + i * 3.2)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.25), Inches(3.05), Inches(1.55))
        lined(box, WHITE, LINE)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), Inches(4.38), Inches(0.32), Inches(0.32))
        rgb(badge, TEAL)
        add_text(s, x + Inches(0.12), Inches(4.40), Inches(0.32), Inches(0.28), [num], size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.52), Inches(4.38), Inches(2.4), Inches(0.32), [title], size=14, color=NAVY, bold=True)
        add_text(s, x + Inches(0.14), Inches(4.78), Inches(2.78), Inches(0.90), [body], size=12, color=SLATE)
        if i < 3:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.95), Inches(4.85), Inches(0.28), Inches(0.18))
            rgb(arr, TEAL)
    add_text(
        s,
        Inches(0.38),
        Inches(5.95),
        Inches(12.55),
        Inches(1.00),
        [
            "递进纪律：状态可知≠数据适用；数据适用≠动作安全；行动可信≠有人付费。T8横切：判断重点是能力、信息模型和接口能否互操作，而非产品名称。",
            "依赖阅读：T2/T3与T1/T6形成可信供给底座，T7-A能力服务可据此先行；只有承诺网络动作和业务结果的T7-B，才额外依赖T4/T5。厂商只作外部校验。",
        ],
        size=12,
        color=SLATE,
        spacing=1.12,
    )
    footer(s, 11, total, "第3章推理")

    # ------------------------------------------------------------------ 12 stage1
    s = new_slide(prs)
    topbar(s, "价值链 1  ·  状态可知", "T2 元数据事件化  +  T3 全局控制、近源执行")
    thesis(s, "重点是变更，不是全量搬家。完整机制=基线快照+变化事件+周期对账。快环不得依赖远端目录；全局控制语义/产品/策略，近源执行时限、减量、自治和断链生存。")
    card(
        s,
        Inches(0.38),
        Inches(1.40),
        Inches(6.25),
        Inches(3.55),
        "T2  元数据从静态目录升级为变化事件  ·  高确定性",
        [
            "推导：静态目录只描述扫描时刻；配置/拓扑/权限/质量持续变化。模型和Agent直接消费后，过期Schema不再只是报表错误，而可能转化为错误动作。消费者要的是“发生了什么、影响谁”，不是再次读取全量目录。",
            "机制：把Schema、对象关系、拓扑、配置、质量、权限、血缘、成本及模型关联的变化，表示为带事件时间、版本、来源、作用域和Correlation ID的可订阅事件。",
            "6G范围：小区、波束、频段、UE群、会话、云资源及其映射；高风险变化阻断/审批，低风险自动测试。",
            "边界：不是把所有业务数据改成事件，也不是让中央事件总线进入快环同步路径。",
            "领先信号：Schema/配置/任务适用性/权限变更自动生成影响清单并触发CI或策略检查。反证：目录仍靠人工扫描，事故早于目录发现变化。",
        ],
        TEAL,
        body_size=10,
    )
    card(
        s,
        Inches(6.78),
        Inches(1.40),
        Inches(6.15),
        Inches(3.55),
        "T3  部署收敛为全局控制、近源执行  ·  高确定性",
        [
            "推导：高频无线与感知数据全量远传受带宽、时延、能耗与隐私约束；快环/MEC/专网/NTN必须在中心断连时继续工作；跨域训练和用途治理又要求语义与策略全局一致。稳定解是控制与执行分层。",
            "全局：对象/语义、任务适用性模板、产品定义、策略、模型和版本发布。",
            "本地：过滤、聚合、特征、缓存、预编译策略、受限推理与动作。回流：异步上传摘要、执行证据、漂移与冲突，支持版本回退。",
            "领先信号：近源算子与全局目录/策略解耦，断链期间按已批准版本运行。反证：所有查询/审批必须经过中心，或各域完全无法共享语义和策略。",
            "外部校验：Ericsson动态放置、Nokia域自治加公共Fabric、Huawei近源算子共同指向两级控制；趋同的是部署约束，不是统一产品或网元。",
        ],
        NAVY,
        body_size=10,
    )
    add_table(
        s,
        Inches(0.38),
        Inches(5.05),
        Inches(12.55),
        Inches(1.90),
        [
            ["关口检验", "验证指标", "若不成立", "2.3外部校验"],
            [
                "状态可知是T1/T6的前提：后续产品化才有稳定对象可引用，任务适用性才有来源和时效。",
                "T2：变更发现P95、影响覆盖、误报/漏报、对账差异、失效传播时长。T3：P99、跨域字节减量、断链可用、策略同步、版本回退。",
                "T1/T6只能做出静态目录，T4的动作会打到过期对象上。",
                "三家都有摄取/近源叙述；跨域事件对账和统一控制均未证明。OpenLineage证明事件可工程化，电信关键是网络对象/配置/任务适用性作用域  [S27]",
            ],
        ],
        col_w=[Inches(2.7), Inches(3.3), Inches(3.15), Inches(3.40)],
        font=8,
    )
    footer(s, 12, total, "T2 / T3")

    # ------------------------------------------------------------------ 13 stage2
    s = new_slide(prs)
    topbar(s, "价值链 2  ·  数据适用", "T1 可运营数据产品  +  T6 面向特定用途的任务适用性")
    thesis(s, "价值单位从「表、Topic、接口和PB」转为「持续满足明确消费者与任务的数据供给」。任务适用性是面向用途的适用性信封，不是统一总分或3GPP既定指标；QoD仅指CAMARA Quality on Demand API。")
    card(
        s,
        Inches(0.38),
        Inches(1.40),
        Inches(6.25),
        Inches(3.55),
        "T1  网络数据封装为可运营产品  ·  高确定性",
        [
            "最小完备集：明确消费者与任务、稳定标识和版本、Owner、Schema/语义、任务适用性、SLO、用途/权限、血缘、成本/计量、保留期及退出。只有目录条目、接口包装或Marketplace页面，仍不是数据产品。",
            "推导：源×消费者增加会使点对点集成爆炸；可复用供给必须出现稳定契约。无线含义依赖小区/频段/波束/配置/采样窗口，产品必须携带上下文和质量边界。没有Owner就无法形成服务承诺。",
            "落地：优先在移动性、节能、故障、QoE和模型特征形成少量高复用产品；默认交付事件、聚合指标、特征或分析服务，而非全量原始流。",
            "领先信号：产品契约进入版本基线，生产者和消费者共同签署SLO/任务适用性。反证：每个场景仍重新采集清洗，产品长期无人消费，Owner与退出机制缺失。",
        ],
        TEAL,
        body_size=10,
    )
    card(
        s,
        Inches(6.78),
        Inches(1.40),
        Inches(6.15),
        Inches(3.55),
        "T6 任务适用性成为模型与闭环的准入标准  ·  中高确定性",
        [
            "定义：用一组可测试的适用性信封说明，数据在何种对象、时空、配置与采样条件下，可以被哪个模型或闭环安全使用。同一KPI在不同小区/频段/波束/配置/采样窗口下不可直接比较。",
            "信封至少包含——通用：准确、完整、一致、新鲜、唯一、可追溯。电信：采样代表性、时空覆盖、同步误差、配置一致、测量置信、suspect/incomplete。AI：标签来源、训练—服务偏差、特征漂移、适用对象与禁用边界。契约：按用途设阈值与降级，而非简单合格/不合格。",
            "领先信号：任务适用性进入数据契约、模型卡、跨域接口和动作放行条件。反证：仍只统计记录量、接口成功率、完整性，或用一个总分掩盖差异。",
            "外部校验：Huawei Data QoS输入、Nokia数据有效性、Ericsson质量/生命周期主张说明竞争正在上移；短期更可能先形成「数据族×场景」契约。",
        ],
        NAVY,
        body_size=10,
    )
    add_table(
        s,
        Inches(0.38),
        Inches(5.05),
        Inches(12.55),
        Inches(1.90),
        [
            ["关口检验", "验证指标", "若不成立", "2.3外部校验"],
            [
                "数据适用是T4/T5的前提：高风险动作才有资格引用该产品；跨域复用不会silent错用。",
                "T1：交付周期、复用消费者、契约违规、重复采集下降、单位任务成本。T6：覆盖/代表性、同步误差、配置一致、漂移、误动作、降级触发。",
                "T4/T5会把不合格数据送进闭环，使T7-B结果承诺失去责任基础；T7-A仍可在稳定产品和能力SLA下先行。",
                "Nokia Data Suite、Ericsson AI-ready和Open Gateway说明产品化已进入供需两侧；跨产品共享与用途重评均未证明  [S10][S24]",
            ],
        ],
        col_w=[Inches(2.7), Inches(3.3), Inches(3.15), Inches(3.40)],
        font=8,
    )
    footer(s, 13, total, "T1 / T6")

    # ------------------------------------------------------------------ 14 stage3
    s = new_slide(prs)
    topbar(s, "价值链 3  ·  行动可信", "T4 供数—决策—执行—回证  +  T5 NDT 作为动作前反事实验证")
    thesis(s, "不是含糊的“耦合或融合”：融合的是接口与证据图，分离的是治理责任。NDT是动作前反事实验证的证据来源，不是绝对安全证明；Fabric供快照，NDT产证据。")
    card(
        s,
        Inches(0.38),
        Inches(1.40),
        Inches(6.25),
        Inches(3.55),
        "T4  供数—决策—执行—回证  ·  中高确定性",
        [
            "四段：Fabric提供版本化产品、上下文、任务适用性与策略 → 模型/Agent依据意图选择数据、工具和候选动作 → 执行面按身份、作用域、租约和风险预算动作 → 动作、结果、回滚与漂移写回共同证据图。",
            "推导：Agent必须持续获得当前对象、配置、质量与权限，而非一次性把文档放进向量库。一旦能调用工具，数据错误就会转化为动作风险。责任不能合并：数据管用途/保留，模型管性能/漂移，Agent管身份/权限/动作。",
            "落地顺序：接口融合（目录检索、订阅、血缘、策略校验以受控工具供Agent调用）→ 证据融合 → 治理分离（IPOE；Reasoner与Actor分离）。先做检索、调查、工单编排和建议，再扩大到白名单低风险动作；快环保留确定性执行。",
            "领先信号：工具调用携带身份、产品/模型版本和策略决策。反证：只有聊天或RAG，无法说明用了哪份数据、谁授权、执行了什么。",
        ],
        TEAL,
        body_size=10,
    )
    card(
        s,
        Inches(6.78),
        Inches(1.40),
        Inches(6.15),
        Inches(3.55),
        "T5  高风险动作以NDT预演为放行关口  ·  中高确定性",
        [
            "关系：Data Fabric向NDT提供带时间、配置、任务适用性与血缘的版本化状态快照；NDT对候选动作做反事实预演并输出收益、冲突和副作用；策略引擎据此放行、灰度或拒绝；生产结果再回流校准。NDT是证据链中的消费者和生产者，不是Fabric同义词。",
            "路径：可信快照 → 孪生预演 → 策略闸门（误差门槛·风险预算·审批·灰度）→ 生产反馈（执行·熔断·回滚·回放·校准）。",
            "成为关口的必要条件：按动作类别定义保真度、最大同步滞后、场景覆盖和可接受预测误差；保存输入快照、孪生/模型版本、候选动作、预演结果、放行理由和生产结果；误差超门槛自动退回建议模式。",
            "领先信号：NDT进入变更审批和Agent动作强制路径。反证：只有拓扑可视化，没有误差门槛、生产反馈或超限退回。",
        ],
        NAVY,
        body_size=10,
    )
    add_table(
        s,
        Inches(0.38),
        Inches(5.05),
        Inches(12.55),
        Inches(1.90),
        [
            ["关口检验", "验证指标", "若不成立", "2.3外部校验"],
            [
                "行动可信是T7-B结果服务的前提：才有资格承诺闭环结果；T7-A能力服务不必等待高风险自动化成熟。",
                "T4：证据完整、越权拦截、错误上下文阻断、人工接管、回滚成功。T5：同步滞后、预测误差、覆盖、回放一致、灰度失败、回滚成功。",
                "外部API与行业结算会放大不可追溯动作，形成责任真空。全面无人治理与跨域自由行动仍不可信。",
                "3GPP和厂商分别推进数据框架、AI生命周期、Agent工具链和自治平台；公开证据不足以支持一个统一Agent Fabric接管全网。SA5已有NDT管理积木≠已能证明任意动作安全  [S4][S23][S29]",
            ],
        ],
        col_w=[Inches(2.7), Inches(3.3), Inches(3.15), Inches(3.40)],
        font=8,
    )
    footer(s, 14, total, "T4 / T5")

    # ------------------------------------------------------------------ 15 stage4
    s = new_slide(prs)
    topbar(s, "价值链 4  ·  价值可结算  ·  T8 横切", "T7沿责任边界分化：能力服务卖输入，结果服务承担输出责任")
    thesis(s, "分叉点不是数据形态，而是决策权、动作权和结果责任由谁承担。稳定产品只决定能否履约；客户集成能力、责任转移意愿与可归因性决定走向哪条路径。T8继续横切承载方式。")
    card(
        s,
        Inches(0.38),
        Inches(1.40),
        Inches(6.25),
        Inches(1.72),
        "路径A  能力型服务  ·  中高确定性",
        [
            "买方保留决策和动作权；供应方保证接口、能力行为与SLA。供给Quality on Demand（QoD）、位置、身份、状态、分析或网络能力。",
            "经聚合层解决跨网覆盖、认证、用量和结算；按调用、会话、覆盖、用量或订阅收费。接口越标准，规模越大但差异化下降。",
        ],
        TEAL,
        title_size=12,
        body_size=10,
    )
    card(
        s,
        Inches(6.78),
        Inches(1.40),
        Inches(6.15),
        Inches(1.72),
        "路径B  结果型托管/自治服务  ·  低—中确定性",
        [
            "供应方参与决策和动作，并承诺体验、成功率、节能或风险改善；客户购买效果和责任转移，而不是单项接口。",
            "需要T4/T5、结果基线、归因、授权、风险上限和回滚；按托管费、效果增益、风险降低或结果分成，只能在边界清晰场景扩张。",
        ],
        AMBER,
        title_size=12,
        body_size=10,
    )
    card(
        s,
        Inches(0.38),
        Inches(3.22),
        Inches(6.25),
        Inches(2.15),
        "T7分化机制  ·  两条路径长期并存",
        [
            "有集成能力的买方希望自行组合输入并保留控制权，形成A；缺少持续运营能力、希望转移责任的客户购买结果，形成B。",
            "A标准化后易复制但议价下降，推动供应方向B延伸；B受归因、授权和责任约束，难完全标准化。两者不是简单成熟度先后。",
            "A验证持续付费、跨网结算和续费；B验证结果基线、归因置信、误动作、风险、争议和回滚成本。两条路径必须分开核算。",
        ],
        TEAL,
        body_size=10,
    )
    card(
        s,
        Inches(6.78),
        Inches(3.22),
        Inches(6.15),
        Inches(2.15),
        "T8横切  ·  低—中确定性",
        [
            "发现、采集、处理、存储、暴露和治理横跨不同域与时标，不同用例不可能共享完全相同的在线流程。标准应优先冻结跨厂商必须一致的外部行为。",
            "最稳定的标准对象是能力语义、信息模型、生命周期、接口行为与一致性测试，而非厂商品牌式拓扑。既可增强DCCF/ADRF/NWDAF/MDA，也可部署SMO服务、边缘Agent或新NF。",
            "领先信号：Rel-21收敛功能、信息模型与接口，同时允许多种部署拓扑。反证：3GPP明确冻结独立统一数据平面及必选NF和流程。",
        ],
        NAVY,
        body_size=10,
    )
    add_text(
        s,
        Inches(0.38),
        Inches(5.48),
        Inches(12.55),
        Inches(1.48),
        [
            "产业断言：6G Data Fabric不会以一个“大一统平台”胜出，而会沿「状态可知 → 数据适用 → 行动可信 → 价值可结算」进入网络。架构上表现为统一控制语义与分布式执行，标准上优先收敛可互操作的能力、信息模型、接口和证据。",
            "趋势边界：事件化≠全量流事件化；产品化≠上架目录；协同闭环≠Agent接管快环；NDT≠安全证明；任务适用性≠单一总分；商业链≠API数量；能力集合≠标准已放弃新NF。越过这些边界，趋势判断就会退化为营销术语。",
            "第4章不再新增趋势，只问：这些关口在3年内能否变成可验证事实，在5年内控制点落在谁手里。",
        ],
        size=11,
        color=SLATE,
        spacing=1.08,
    )
    footer(s, 15, total, "T7 / T8")

    # ------------------------------------------------------------------ 16 ch4
    s = new_slide(prs)
    topbar(s, "第4章  3年看透 / 5年看清", "只投射 T1–T8：3年固定不变量，5年看清控制点，不猜网元名称")
    thesis(s, "「看透」是把趋势变成可验证、可证伪、可停止的工程命题；「看清」是观察谁掌握语义/任务适用性、近源运行时、Agent/NDT证据、跨网结算和互操作规则。不猜网元名称，不再另造趋势。")
    add_table(
        s,
        Inches(0.38),
        Inches(1.36),
        Inches(12.55),
        Inches(4.42),
        [
            ["趋势", "1–3年必须出现的可验证事实（约2026–2028）", "3–5年可能看清的结构性结果（约2029–2031+）", "停止线 / 仍需保留的不确定性"],
            ["T1 数据产品  高", "故障/节能/移动性/QoE或模型特征中，至少少量产品具有消费者、Owner、Schema、任务适用性、SLO、用途与退出；同一产品跨第二场景/后端复用", "跨域供给从项目接口转为版本化产品组合；高价值产品进入主航道版本和运营体系，部分按任务或结果计量", "若仍按项目重复取数清洗，“产品化”只是目录改名；原始数据出售也未必成为主模式"],
            ["T2 元数据事件  高", "Schema、配置、拓扑、任务适用性、权限和血缘变化能触发影响分析、契约测试、缓存/模型失效或审批；同时保留基线快照和周期对账", "元数据事件成为质量治理、模型重验和策略执行的共同状态输入；跨域对象ID、事件模型、血缘与影响图形成控制点", "若事故仍早于目录发现变化，控制面仍是被动盘点；事件不应成为快环在线依赖"],
            ["T3 两级部署  高", "出现“SMO/Non-RT全局定义 + RAN/边缘本地处理”的时间盒PoC；证明断链运行、异步回证、版本同步与回退，不进入<10ms同步路径", "全局控制与分布执行成为默认形态；近源运行时、策略/模型分发、断链与回退构成设备商和云平台的分工边界", "若收益覆盖不了站点开销或必须依赖中心，应停止近源子Fabric产品化"],
            ["T4 Agent协同  中高", "目录/血缘/订阅/策略成为受控工具；Agent调用带身份、作用域、租约、数据/模型版本和结果证据。只读调查与建议规模化，低风险动作小范围白名单化", "低风险调查、编排和白名单动作常态化；数据/模型/工具/动作共享证据图，但治理责任继续分离", "若只有聊天或RAG而没有工具责任链，就不能称为可操作协同；全面无人治理仍不可信"],
            ["T5 NDT关口  中高", "至少一个高风险变更完成历史回放与影子运行；为保真度、同步滞后、误差和覆盖设门槛，超限自动退回建议/人工模式", "高风险动作形成“可信快照—孪生预演—策略—灰度—回滚—生产校准”强制路径；场景库、误差门槛和放行规则成为控制点", "若孪生仍只有可视化，就不能承担放行证据；NDT始终是证据源之一，不是绝对安全证明"],
            ["T6 任务适用性准入  中高", "按数据族形成适用性信封并进入数据契约、模型卡或接口；至少覆盖采样代表性、时空覆盖、同步/配置一致、测量置信及训练—服务偏差", "不同数据族的适用性信封收敛公共属性并进入互操作测试；任务适用性成为模型、闭环和跨域产品的共同验收语言", "若仍只报完整性和接口成功率，适用性尚未被管理；短期难出现覆盖所有数据和任务的统一总分"],
            ["T7 责任双路径  A中高/B低—中", "A：至少一个标准能力出现持续付费、跨网一致和可核算履约成本。B：至少一个场景在合同中明确结果基线、动作权限、归因、责任和回滚", "有集成能力的客户购买可组合输入，希望转移运营责任的客户购买可归因结果；A标准化扩张，B在责任清晰场景渐进发展", "A若只有API数量无持续收入则不成立；B若不能定义基线、归因、定责或回退则不能规模化；两条路径不得混算"],
            ["T8 能力集合  低—中", "Rel-21讨论开始收敛发现、处理、存储、暴露、数据质量/适用性属性和生命周期等能力及接口；产品原型证明可插拔后端与多种部署映射", "标准优先冻结能力语义、信息模型、接口行为和一致性测试；增强既有功能、边缘Agent、SMO服务或新NF可按部署组合", "若3GPP明确冻结独立统一数据平面及必选NF，则需调整；在此之前不押注名称，也不预设标准一定拒绝新功能"],
        ],
        col_w=[Inches(1.85), Inches(3.65), Inches(3.55), Inches(3.50)],
        font=7,
    )
    add_text(
        s,
        Inches(0.38),
        Inches(5.86),
        Inches(12.55),
        Inches(1.15),
        [
            "1–3年决策：先做可证伪底座。优先建立T2/T3事件和两级运行，再用T1/T6证明可复用、可验收；T4/T5从只读、建议、影子运行开始；T7-A看真实付费、跨网一致和履约成本，T7-B看基线、归因、责任和回滚；T8以进入工作文本、参考实现和互操作测试为标准验证。",
            "3–5年决策：看谁拥有控制点。设备商的稳定优势在网络语义、任务适用性和近源执行；云/数据平台在模型、Catalog与开发工具；聚合商在分发与结算。运营商将要求对象、策略、证据可导出，运行时和后端可替换。",
            "收束：3年让每条趋势都有可验证、可停止的证据；5年看清谁拥有语义/任务适用性、近源运行、动作证据、标准接口与跨网结算。投射基于截至2026-08的公开材料，不等于3GPP已采纳，也不构成厂商产品承诺。",
        ],
        size=10,
        color=NAVY,
        spacing=1.06,
    )
    footer(s, 16, total, "第4章")

    prs.save(OUT)
    print(f"Wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
